"""Safe, report-only comparisons for a candidate lecture-document parser."""

from __future__ import annotations

import hashlib
import json
import os
from collections import Counter
from collections.abc import Iterable
from dataclasses import dataclass
from pathlib import Path
from time import perf_counter
from typing import Any
from uuid import uuid4

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from oms_hub.document_processing.domain import (
    DocumentLocator,
    DocumentProcessor,
    ParsedAsset,
    ParsedDocument,
    ParsedSegment,
    SegmentKind,
    SourceSnapshot,
)
from oms_hub.document_processing.router import ParserMode


@dataclass(frozen=True, slots=True)
class ShadowComparison:
    legacy_document: ParsedDocument | None
    candidate_document: ParsedDocument | None
    report: dict[str, object]


@dataclass(frozen=True, slots=True)
class ShadowParseResult:
    document: ParsedDocument
    report: dict[str, object]
    degraded: bool


class LegacyPptxProcessor:
    """Local, non-Anki PowerPoint text extraction used as the shadow baseline."""

    name = "legacy-pptx"
    version = "1"

    def supports(self, snapshot: SourceSnapshot) -> bool:
        return snapshot.path.suffix.casefold() == ".pptx"

    def parse(self, snapshot: SourceSnapshot, asset_root: Path) -> ParsedDocument:
        del asset_root
        if not self.supports(snapshot):
            raise ValueError(f"legacy PowerPoint parser does not support {snapshot.path.name!r}")
        presentation = Presentation(str(snapshot.path))
        segments: list[ParsedSegment] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                text = ""
                kind = SegmentKind.PARAGRAPH
                if getattr(shape, "has_table", False):
                    kind = SegmentKind.TABLE
                    text = "\n".join(
                        " | ".join(cell.text.strip() for cell in row.cells)
                        for row in shape.table.rows
                    ).strip()
                elif getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                if text:
                    index = len(segments) + 1
                    segments.append(
                        ParsedSegment(
                            key=f"slide-{slide_number}-block-{index}",
                            kind=kind,
                            text=text,
                            locator=DocumentLocator(
                                f"slide {slide_number}",
                                slide_number=slide_number,
                                block_index=index,
                            ),
                        )
                    )
            notes = getattr(getattr(slide, "notes_slide", None), "notes_text_frame", None)
            note_text = (notes.text or "").strip() if notes is not None else ""
            if note_text:
                index = len(segments) + 1
                segments.append(
                    ParsedSegment(
                        key=f"slide-{slide_number}-note-{index}",
                        kind=SegmentKind.NOTE,
                        text=note_text,
                        locator=DocumentLocator(
                            f"slide {slide_number} notes",
                            slide_number=slide_number,
                            block_index=index,
                        ),
                    )
                )
        return ParsedDocument(
            source_id=snapshot.id,
            source_sha256=snapshot.sha256,
            source_format="pptx",
            parser_name=self.name,
            parser_version=self.version,
            segments=tuple(segments),
            assets=_pptx_image_assets(presentation),
            warnings=(),
        )


def _pptx_image_assets(presentation: Any) -> tuple[ParsedAsset, ...]:
    """Describe unique embedded images without copying their source bytes."""
    discovered: dict[str, tuple[str, list[DocumentLocator]]] = {}
    for slide_number, slide in enumerate(presentation.slides, start=1):
        image_number = 0
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type is not MSO_SHAPE_TYPE.PICTURE:
                continue
            image_number += 1
            image = shape.image
            digest = hashlib.sha256(image.blob).hexdigest()
            media_type = str(image.content_type)
            occurrence = DocumentLocator(
                label=f"slide {slide_number} image {image_number}",
                slide_number=slide_number,
            )
            if digest not in discovered:
                discovered[digest] = (media_type, [])
            discovered[digest][1].append(occurrence)
    assets: list[ParsedAsset] = []
    for digest, (media_type, occurrences) in discovered.items():
        slide_numbers = {occurrence.slide_number for occurrence in occurrences}
        locator = (
            occurrences[0]
            if len(slide_numbers) == 1
            else DocumentLocator(label="embedded image")
        )
        assets.append(
            ParsedAsset(
                key=f"image-{digest}",
                path=None,
                media_type=media_type,
                sha256=digest,
                locator=locator,
            )
        )
    return tuple(assets)


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type is MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


class DocumentShadowEvaluator:
    """Run a candidate parser without ever exposing document text in reports."""

    def __init__(self, candidate: DocumentProcessor, legacy: DocumentProcessor) -> None:
        self.candidate = candidate
        self.legacy = legacy

    def compare(self, snapshot: SourceSnapshot, asset_root: Path) -> ShadowComparison:
        started = perf_counter()
        legacy_document, legacy_error, legacy_duration = self._parse(
            self.legacy, snapshot, asset_root / "legacy", "legacy"
        )
        candidate_document, candidate_error, candidate_duration = self._parse(
            self.candidate, snapshot, asset_root / "candidate", "candidate"
        )
        blockers = list(_comparison_blockers(legacy_document, candidate_document))
        if legacy_error:
            blockers.append("legacy parser failed")
        if candidate_error:
            blockers.append("candidate parser failed")
        report: dict[str, object] = {
            "source_sha256": snapshot.sha256,
            "mode": ParserMode.SHADOW.value,
            "duration_ms": round((perf_counter() - started) * 1000, 3),
            "legacy": _document_report(
                legacy_document, self.legacy, legacy_duration, legacy_error
            ),
            "candidate": _document_report(
                candidate_document, self.candidate, candidate_duration, candidate_error
            ),
            "candidate_error": candidate_error,
            "fallback_used": False,
            "degraded": False,
            "promotion_blockers": tuple(sorted(set(blockers))),
        }
        return ShadowComparison(legacy_document, candidate_document, report)

    def parse_primary(self, snapshot: SourceSnapshot, asset_root: Path) -> ShadowParseResult:
        comparison = self.compare(snapshot, asset_root)
        report = dict(comparison.report)
        report["mode"] = ParserMode.ANYDOC.value
        if (
            comparison.candidate_document is not None
            and not report["promotion_blockers"]
        ):
            return ShadowParseResult(comparison.candidate_document, report, False)
        if comparison.legacy_document is not None:
            report["fallback_used"] = True
            report["degraded"] = True
            return ShadowParseResult(comparison.legacy_document, report, True)
        raise ValueError("both candidate and legacy document parsers failed")

    def parse(
        self,
        snapshot: SourceSnapshot,
        asset_root: Path,
        mode: ParserMode,
    ) -> ShadowParseResult:
        if mode is ParserMode.ANYDOC:
            return self.parse_primary(snapshot, asset_root)
        if mode is ParserMode.LEGACY:
            legacy_document, legacy_error, legacy_duration = self._parse(
                self.legacy, snapshot, asset_root / "legacy", "legacy"
            )
            if legacy_document is None:
                raise ValueError("legacy document parser failed")
            report = self._report(
                snapshot,
                ParserMode.LEGACY,
                legacy_document,
                legacy_error,
                legacy_duration,
                None,
                None,
                0.0,
                (),
            )
            return ShadowParseResult(legacy_document, report, False)
        comparison = self.compare(snapshot, asset_root)
        if comparison.legacy_document is None:
            raise ValueError("legacy document parser failed")
        report = dict(comparison.report)
        report["mode"] = mode.value
        return ShadowParseResult(comparison.legacy_document, report, False)

    def _report(
        self,
        snapshot: SourceSnapshot,
        mode: ParserMode,
        legacy_document: ParsedDocument | None,
        legacy_error: str | None,
        legacy_duration: float,
        candidate_document: ParsedDocument | None,
        candidate_error: str | None,
        candidate_duration: float,
        blockers: tuple[str, ...],
    ) -> dict[str, object]:
        return {
            "source_sha256": snapshot.sha256,
            "mode": mode.value,
            "duration_ms": round(legacy_duration + candidate_duration, 3),
            "legacy": _document_report(
                legacy_document, self.legacy, legacy_duration, legacy_error
            ),
            "candidate": _document_report(
                candidate_document, self.candidate, candidate_duration, candidate_error
            ),
            "candidate_error": candidate_error,
            "fallback_used": False,
            "degraded": False,
            "promotion_blockers": tuple(sorted(set(blockers))),
        }

    def exceptional_report(
        self,
        source_sha256: str,
        mode: ParserMode,
        diagnostic_code: str,
    ) -> dict[str, object]:
        return {
            "source_sha256": source_sha256,
            "mode": mode.value,
            "duration_ms": 0.0,
            "legacy": _document_report(None, self.legacy, 0.0, "legacy_not_available"),
            "candidate": _document_report(
                None, self.candidate, 0.0, diagnostic_code
            ),
            "candidate_error": diagnostic_code,
            "fallback_used": False,
            "degraded": mode is ParserMode.ANYDOC,
            "promotion_blockers": ("document evaluation failed",),
        }

    @staticmethod
    def write_report(report: dict[str, object], destination: Path) -> None:
        destination.parent.mkdir(parents=True, exist_ok=True)
        temporary = destination.with_name(f".{destination.name}.partial-{uuid4().hex}")
        try:
            with temporary.open("x", encoding="utf-8") as stream:
                json.dump(report, stream, sort_keys=True, separators=(",", ":"))
                stream.write("\n")
                stream.flush()
                os.fsync(stream.fileno())
            os.replace(temporary, destination)
        finally:
            temporary.unlink(missing_ok=True)

    @staticmethod
    def _parse(
        processor: DocumentProcessor,
        snapshot: SourceSnapshot,
        asset_root: Path,
        role: str,
    ) -> tuple[ParsedDocument | None, str | None, float]:
        started = perf_counter()
        if not processor.supports(snapshot):
            return None, f"{role}_unsupported", 0.0
        try:
            return processor.parse(snapshot, asset_root), None, round(
                (perf_counter() - started) * 1000, 3
            )
        except Exception:  # noqa: BLE001 - parser failures become safe diagnostics
            return None, f"{role}_parse_failed", round((perf_counter() - started) * 1000, 3)


def _document_report(
    document: ParsedDocument | None,
    processor: DocumentProcessor,
    duration_ms: float,
    error: str | None,
) -> dict[str, object]:
    if document is None:
        return {
            "parser_name": processor.name,
            "parser_version": processor.version,
            "duration_ms": duration_ms,
            "segment_counts": {},
            "page_coverage": (),
            "slide_coverage": (),
            "notes": 0,
            "tables": 0,
            "assets": 0,
            "warnings": (),
            "normalized_text_sha256": None,
            "error": error,
        }
    counts = Counter(segment.kind.value for segment in document.segments)
    normalized = "\n".join(_normalize_text(segment.text) for segment in document.segments)
    return {
        "parser_name": document.parser_name,
        "parser_version": document.parser_version,
        "duration_ms": duration_ms,
        "segment_counts": dict(sorted(counts.items())),
        "page_coverage": tuple(
            sorted(
                {
                    segment.locator.page_number
                    for segment in document.segments
                    if segment.locator.page_number
                }
            )
        ),
        "slide_coverage": tuple(
            sorted(
                {
                    segment.locator.slide_number
                    for segment in document.segments
                    if segment.locator.slide_number
                }
            )
        ),
        "notes": counts[SegmentKind.NOTE.value],
        "tables": counts[SegmentKind.TABLE.value],
        "assets": len(document.assets),
        "warnings": tuple(
            hashlib.sha256(warning.encode("utf-8")).hexdigest()
            for warning in document.warnings
        ),
        "normalized_text_sha256": hashlib.sha256(normalized.encode("utf-8")).hexdigest(),
        "error": None,
    }


def _comparison_blockers(
    legacy: ParsedDocument | None,
    candidate: ParsedDocument | None,
) -> tuple[str, ...]:
    if candidate is None:
        return ("candidate parser failed",)
    blockers: list[str] = []
    if not candidate.segments:
        blockers.append("candidate has no segments")
    if candidate.warnings:
        blockers.append("candidate emitted warnings")
    if legacy is None:
        blockers.append("legacy parser failed")
        return tuple(blockers)
    if _normalized_document_hash(legacy) != _normalized_document_hash(candidate):
        blockers.append("normalized text differs")
    if not _coverage(legacy, "page_number").issubset(_coverage(candidate, "page_number")):
        blockers.append("candidate reduced legacy page coverage")
    if not _coverage(legacy, "slide_number").issubset(_coverage(candidate, "slide_number")):
        blockers.append("candidate reduced legacy slide coverage")
    for kind, label in (
        (SegmentKind.NOTE, "notes"),
        (SegmentKind.TABLE, "tables"),
    ):
        if _count_kind(candidate, kind) < _count_kind(legacy, kind):
            blockers.append(f"candidate has fewer {label}")
    if len(candidate.assets) < len(legacy.assets):
        blockers.append("candidate has fewer assets")
    return tuple(blockers)


def _normalize_text(value: str) -> str:
    return " ".join(value.casefold().split())


def _normalized_document_hash(document: ParsedDocument) -> str:
    normalized = "\n".join(_normalize_text(segment.text) for segment in document.segments)
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


def _coverage(document: ParsedDocument, field: str) -> set[int]:
    return {
        value
        for segment in document.segments
        if (value := getattr(segment.locator, field)) is not None
    }


def _count_kind(document: ParsedDocument, kind: SegmentKind) -> int:
    return sum(segment.kind is kind for segment in document.segments)
