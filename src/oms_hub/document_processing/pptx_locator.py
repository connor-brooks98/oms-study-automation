"""PowerPoint-specific provenance restored from package relationships."""

from __future__ import annotations

from collections.abc import Iterable
from dataclasses import replace
from pathlib import PurePosixPath
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from oms_hub.document_processing.domain import (
    DocumentLocator,
    ParsedAsset,
    ParsedDocument,
    ParsedSegment,
    SegmentKind,
    SourceSnapshot,
)
from oms_hub.document_processing.ocr import LocalOcr

_NATIVE_WORDS_WITH_OPTIONAL_IMAGE_OCR = 8


class PptxLocatorEnricher:
    """Restore unambiguous slide, notes, and image locations for an Anydoc result."""

    def __init__(self, ocr: LocalOcr | None = None) -> None:
        self.ocr = ocr or LocalOcr()

    def enrich(self, snapshot: SourceSnapshot, parsed: ParsedDocument) -> ParsedDocument:
        presentation = Presentation(str(snapshot.path))
        image_locations = _image_locations(presentation)
        assets, asset_keys_by_location, warnings = _locate_assets(parsed.assets, image_locations)
        source_segments = tuple(_slide_segments(presentation, asset_keys_by_location))
        segments = _enrich_segments(
            parsed.segments,
            source_segments,
            assets,
        )
        ocr_segments, ocr_warnings = _ocr_slide_images(
            presentation,
            segments,
            source_segments,
            assets,
            asset_keys_by_location,
            self.ocr,
        )
        style_segments = _unmatched_style_segments(segments, source_segments)
        return replace(
            parsed,
            segments=tuple(
                sorted(
                    (*segments, *ocr_segments, *style_segments),
                    key=lambda segment: (
                        segment.locator.slide_number or 10**9,
                        1 if segment.key.endswith("-ocr") else 0,
                        segment.key,
                    ),
                )
            ),
            assets=assets,
            warnings=(*parsed.warnings, *warnings, *ocr_warnings),
        )


def _enrich_segments(
    parsed_segments: tuple[ParsedSegment, ...],
    source_segments: tuple[ParsedSegment, ...],
    assets: tuple[ParsedAsset, ...],
) -> tuple[ParsedSegment, ...]:
    """Restore only unambiguous provenance; preserve Anydoc semantics verbatim."""
    asset_slides = {
        asset.key: asset.locator.slide_number
        for asset in assets
        if asset.locator.slide_number is not None
    }
    locations: dict[str, DocumentLocator] = {}
    for kind in SegmentKind:
        candidates = tuple(segment for segment in parsed_segments if segment.kind is kind)
        references = tuple(segment for segment in source_segments if segment.kind is kind)
        for candidate in candidates:
            if not candidate.text.strip():
                continue
            matching_locations = {
                reference.locator
                for reference in source_segments
                if _normalize_text(reference.text) == _normalize_text(candidate.text)
            }
            if len(matching_locations) == 1:
                locations[candidate.key] = matching_locations.pop()
        if any(candidate.key in locations for candidate in candidates):
            continue
        if len(candidates) == len(references):
            locations.update(
                {
                    candidate.key: reference.locator
                    for candidate, reference in zip(candidates, references, strict=True)
                }
            )
            continue
        reference_slides = tuple(
            dict.fromkeys(
                reference.locator.slide_number
                for reference in references
                if reference.locator.slide_number is not None
            )
        )
        if len(candidates) == len(reference_slides):
            locations.update(
                {
                    candidate.key: DocumentLocator(
                        label=f"slide {slide_number}", slide_number=slide_number
                    )
                    for candidate, slide_number in zip(
                        candidates, reference_slides, strict=True
                    )
                }
            )
    enriched: list[ParsedSegment] = []
    styles_by_text: dict[str, tuple[str, ...]] = {}
    for source in source_segments:
        normalized = _normalize_text(source.text)
        if normalized and source.style_metadata:
            styles_by_text.setdefault(normalized, source.style_metadata)
    for segment in parsed_segments:
        if segment.key in locations:
            enriched.append(
                replace(
                    segment,
                    locator=locations[segment.key],
                    style_metadata=styles_by_text.get(
                        _normalize_text(segment.text), segment.style_metadata
                    ),
                )
            )
            continue
        slides = {asset_slides[key] for key in segment.asset_keys if key in asset_slides}
        if len(slides) == 1:
            slide_number = next(iter(slides))
            enriched.append(
                replace(
                    segment,
                    locator=DocumentLocator(
                        label=f"slide {slide_number}", slide_number=slide_number
                    ),
                )
            )
            continue
        enriched.append(segment)
    # Some Anydoc PPTX outputs omit speaker notes.  Retain the candidate sequence
    # verbatim and append only the missing, source-proven note records.
    if not any(segment.kind is SegmentKind.NOTE for segment in parsed_segments):
        enriched.extend(
            segment
            for segment in source_segments
            if segment.kind is SegmentKind.NOTE
        )
    return _restore_slide_styles(tuple(enriched), source_segments)


def _restore_slide_styles(
    segments: tuple[ParsedSegment, ...],
    source_segments: tuple[ParsedSegment, ...],
) -> tuple[ParsedSegment, ...]:
    """Attach unmatched run cues once to a text segment on their source slide."""
    styles_by_slide: dict[int, list[str]] = {}
    for source in source_segments:
        slide_number = source.locator.slide_number
        if slide_number is not None and source.style_metadata:
            styles_by_slide.setdefault(slide_number, []).extend(source.style_metadata)
    restored_cues = {
        (segment.locator.slide_number, cue)
        for segment in segments
        if segment.locator.slide_number is not None
        for cue in segment.style_metadata
    }
    restored = list(segments)
    for slide_number, styles in styles_by_slide.items():
        remaining = tuple(dict.fromkeys(
            cue for cue in styles if (slide_number, cue) not in restored_cues
        ))
        if not remaining:
            continue
        index = next(
            (
                index
                for index, segment in enumerate(restored)
                if segment.locator.slide_number == slide_number and segment.text.strip()
            ),
            None,
        )
        if index is not None:
            restored[index] = replace(
                restored[index],
                style_metadata=tuple(dict.fromkeys((*restored[index].style_metadata, *remaining))),
            )
    return tuple(restored)


def _unmatched_style_segments(
    segments: tuple[ParsedSegment, ...],
    source_segments: tuple[ParsedSegment, ...],
) -> tuple[ParsedSegment, ...]:
    """Keep slide-level formatting cues even when candidate text has no slide locator."""
    restored_slides = {
        segment.locator.slide_number
        for segment in segments
        if segment.locator.slide_number is not None and segment.style_metadata
    }
    styles_by_slide: dict[int, list[str]] = {}
    for source in source_segments:
        slide_number = source.locator.slide_number
        if (
            slide_number is not None
            and slide_number not in restored_slides
            and source.style_metadata
        ):
            styles_by_slide.setdefault(slide_number, []).extend(source.style_metadata)
    return tuple(
        ParsedSegment(
            key=f"slide-{slide_number}-style-metadata",
            kind=SegmentKind.NOTE,
            text="",
            locator=DocumentLocator(
                f"slide {slide_number} formatting", slide_number=slide_number
            ),
            style_metadata=tuple(dict.fromkeys(styles)),
        )
        for slide_number, styles in sorted(styles_by_slide.items())
    )


def _locate_assets(
    parsed_assets: tuple[ParsedAsset, ...], image_locations: dict[str, tuple[_ImageLocation, ...]]
) -> tuple[tuple[ParsedAsset, ...], dict[tuple[int, int], str], tuple[str, ...]]:
    assets: list[ParsedAsset] = []
    keys_by_location: dict[tuple[int, int], str] = {}
    warnings: list[str] = []
    for asset in parsed_assets:
        matches = image_locations.get(_normalize_part_name(asset.origin), ())
        slide_numbers = {match.slide_number for match in matches}
        if asset.path is not None:
            for occurrence in matches:
                keys_by_location[(occurrence.slide_number, occurrence.image_number)] = asset.key
        if len(slide_numbers) == 1:
            match = matches[0]
            located = replace(
                asset,
                locator=DocumentLocator(
                    label=f"slide {match.slide_number} image {match.image_number}",
                    slide_number=match.slide_number,
                ),
            )
            assets.append(located)
            continue
        assets.append(asset)
        if slide_numbers:
            warnings.append(
                f"asset {asset.key!r} occurs on multiple slides and was not automatically bound"
            )
        else:
            warnings.append(f"asset {asset.key!r} could not be matched to a PowerPoint slide")
    return tuple(assets), keys_by_location, tuple(warnings)


class _ImageLocation:
    def __init__(self, slide_number: int, image_number: int) -> None:
        self.slide_number = slide_number
        self.image_number = image_number


def _image_locations(presentation: Any) -> dict[str, tuple[_ImageLocation, ...]]:
    locations: dict[str, list[_ImageLocation]] = {}
    for slide_number, slide in enumerate(presentation.slides, start=1):
        image_number = 0
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type is not MSO_SHAPE_TYPE.PICTURE:
                continue
            image_number += 1
            part_name = _relationship_image_part_name(slide, shape)
            if part_name is None:
                continue
            locations.setdefault(part_name, []).append(_ImageLocation(slide_number, image_number))
    return {part_name: tuple(values) for part_name, values in locations.items()}


def _relationship_image_part_name(slide: Any, shape: Any) -> str | None:
    relationship_id = shape._element.blip_rId  # pyright: ignore[reportPrivateUsage]
    if relationship_id is None:
        return None
    relationship = slide.part.rels.get(relationship_id)
    if relationship is None or relationship.target_part is None:
        return None
    return _normalize_part_name(str(relationship.target_part.partname))


def _slide_segments(
    presentation: Any, asset_keys_by_location: dict[tuple[int, int], str]
) -> Iterable[ParsedSegment]:
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_number = 0
        image_number = 0
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type is MSO_SHAPE_TYPE.PICTURE:
                image_number += 1
                key = asset_keys_by_location.get((slide_number, image_number))
                if key is not None:
                    yield ParsedSegment(
                        key=f"slide-{slide_number}-image-{image_number}",
                        kind=SegmentKind.IMAGE,
                        text="",
                        locator=DocumentLocator(
                            label=f"slide {slide_number} image {image_number}",
                            slide_number=slide_number,
                        ),
                        asset_keys=(key,),
                    )
                continue
            text, kind = _shape_content(shape)
            if not text:
                continue
            text_number += 1
            yield ParsedSegment(
                key=f"slide-{slide_number}-content-{text_number}",
                kind=kind,
                text=text,
                locator=DocumentLocator(
                    label=f"slide {slide_number} content {text_number}",
                    slide_number=slide_number,
                ),
                style_metadata=_shape_style_metadata(shape),
            )
        notes = _notes_text(slide)
        if notes:
            yield ParsedSegment(
                key=f"slide-{slide_number}-notes",
                kind=SegmentKind.NOTE,
                text=notes,
                locator=DocumentLocator(
                    label=f"slide {slide_number} notes", slide_number=slide_number
                ),
            )


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type is MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _shape_content(shape: Any) -> tuple[str, SegmentKind]:
    if shape.has_table:
        table = shape.table
        return (
            "\n".join(
                " | ".join(cell.text.strip() for cell in row.cells) for row in table.rows
            ).strip(),
            SegmentKind.TABLE,
        )
    if shape.has_text_frame:
        return shape.text.strip(), SegmentKind.PARAGRAPH
    return "", SegmentKind.PARAGRAPH


def _shape_style_metadata(shape: Any) -> tuple[str, ...]:
    """Keep answer-format clues as provenance rather than modifying displayed text."""
    if not getattr(shape, "has_text_frame", False):
        return ()
    cues: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text.strip()
            if not text:
                continue
            if run.font.bold:
                cues.append(f"bold: {text}")
            if run.font.italic:
                cues.append(f"italic: {text}")
            if run.font.underline:
                cues.append(f"underline: {text}")
            if run._r.xpath("./a:rPr/a:highlight"):  # pyright: ignore[reportPrivateUsage]
                cues.append(f"highlighted: {text}")
            color = getattr(getattr(run.font, "color", None), "rgb", None)
            if color is not None:
                cues.append(f"color #{color}: {text}")
    return tuple(dict.fromkeys(cues))


def _ocr_slide_images(
    presentation: Any,
    segments: tuple[ParsedSegment, ...],
    source_segments: tuple[ParsedSegment, ...],
    assets: tuple[ParsedAsset, ...],
    asset_keys_by_location: dict[tuple[int, int], str],
    ocr: LocalOcr,
) -> tuple[tuple[ParsedSegment, ...], tuple[str, ...]]:
    """OCR image-only slides and visual question screenshots, never small logos.

    # ponytail: 30% slide-area threshold; tune only if real decks show false negatives.
    """
    text_slides = {
        segment.locator.slide_number
        for segment in segments
        if segment.locator.slide_number is not None and segment.text.strip()
    }
    native_words_by_slide: dict[int, int] = {}
    for segment in source_segments:
        slide_number = segment.locator.slide_number
        if slide_number is not None and segment.kind is not SegmentKind.NOTE:
            native_words_by_slide[slide_number] = (
                native_words_by_slide.get(slide_number, 0) + len(segment.text.split())
            )
    assets_by_slide: dict[int, list[ParsedAsset]] = {}
    large_asset_locations: set[tuple[int, str]] = set()
    slide_area = presentation.slide_width * presentation.slide_height
    for slide_number, slide in enumerate(presentation.slides, start=1):
        image_number = 0
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type is not MSO_SHAPE_TYPE.PICTURE:
                continue
            image_number += 1
            if shape.width * shape.height >= slide_area * 0.30:
                key = asset_keys_by_location.get((slide_number, image_number))
                if key is not None:
                    large_asset_locations.add((slide_number, key))
    assets_by_key = {asset.key: asset for asset in assets if asset.path is not None}
    for (slide_number, _image_number), key in asset_keys_by_location.items():
        asset = assets_by_key.get(key)
        if asset is not None and asset not in assets_by_slide.setdefault(slide_number, []):
            assets_by_slide[slide_number].append(asset)
    added: list[ParsedSegment] = []
    blockers: list[str] = []
    for slide_number, slide_assets in sorted(assets_by_slide.items()):
        candidates = (
            slide_assets
            if slide_number not in text_slides
            else [
                asset
                for asset in slide_assets
                if (slide_number, asset.key) in large_asset_locations
            ]
        )
        if not candidates:
            continue
        recognized = "\n".join(
            text
            for asset in candidates
            if asset.path is not None and (text := ocr.text(asset.path))
        ).strip()
        if recognized:
            added.append(
                ParsedSegment(
                    key=f"slide-{slide_number}-ocr",
                    kind=SegmentKind.PARAGRAPH,
                    text=recognized,
                    locator=DocumentLocator(f"slide {slide_number} OCR", slide_number=slide_number),
                    asset_keys=tuple(asset.key for asset in candidates),
                )
            )
        elif native_words_by_slide.get(slide_number, 0) < _NATIVE_WORDS_WITH_OPTIONAL_IMAGE_OCR:
            blockers.append(
                f"BLOCKER: OCR is required but unavailable or empty for slide {slide_number}"
            )
    return tuple(added), tuple(blockers)


def _notes_text(slide: Any) -> str:
    try:
        return str(slide.notes_slide.notes_text_frame.text).strip()
    except (AttributeError, KeyError):
        return ""


def _normalize_part_name(origin: str | None) -> str:
    if not origin:
        return ""
    normalized = origin.replace("\\", "/").lstrip("/")
    return str(PurePosixPath(normalized))


def _normalize_text(value: str) -> str:
    return " ".join(value.casefold().split())
