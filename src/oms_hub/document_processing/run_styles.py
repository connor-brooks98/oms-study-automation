"""Immutable, parallel PowerPoint run-style extraction."""

from __future__ import annotations

import colorsys
import hashlib
from collections import deque
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pydantic import BaseModel, ConfigDict, Field, field_validator, model_validator

from oms_hub.anki.contracts import canonical_payload_sha256
from oms_hub.document_processing.pptx_locator import walk_shapes_with_paths

RUN_STYLE_PARSER_VERSION = "pptx-run-styles-v1"


class StyledTextRun(BaseModel):
    """One physical PowerPoint run, including blank runs for stable locators."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    source_id: str = Field(min_length=1)
    source_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    slide_number: int = Field(gt=0)
    locator: str = Field(min_length=1)
    text: str
    normalized_text_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    explicit_rgb: str | None = Field(default=None, pattern=r"^[0-9A-F]{6}$")
    theme_color: str | None = Field(default=None, min_length=1)
    inherited_color_resolution: str | None = None
    resolved_color: str | None = Field(default=None, pattern=r"^[0-9A-F]{6}$")
    brightness: float | None = Field(default=None, ge=-1, le=1)
    bold: bool | None = None
    italic: bool | None = None
    underline: bool | None = None
    highlight: str | None = None
    color_attempted: bool = False

    @field_validator("explicit_rgb", "resolved_color", mode="before")
    @classmethod
    def _normalize_rgb(cls, value: object) -> str | None:
        return None if value is None else str(value).strip().upper()

    @field_validator("theme_color", "inherited_color_resolution", mode="before")
    @classmethod
    def _normalize_optional_text(cls, value: object) -> str | None:
        return None if value is None else str(value).strip().casefold() or None

    @model_validator(mode="after")
    def _validate_normalized_text_hash(self) -> StyledTextRun:
        if self.normalized_text_sha256 != normalized_text_sha256(self.text):
            raise ValueError("style run normalized text hash does not match its text")
        return self


class StyledTextRunSidecar(BaseModel):
    """Hash-bound style records that deliberately do not alter ParsedDocument."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    source_id: str = Field(min_length=1)
    source_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    parser_version: str = Field(min_length=1)
    runs: tuple[StyledTextRun, ...]
    sidecar_sha256: str = ""

    def canonical_payload(self) -> dict[str, object]:
        return self.model_dump(mode="json", exclude={"sidecar_sha256"})

    @model_validator(mode="after")
    def _validate_hash(self) -> StyledTextRunSidecar:
        if any(
            run.source_id != self.source_id or run.source_sha256 != self.source_sha256
            for run in self.runs
        ):
            raise ValueError("style run source identity does not match sidecar")
        locators = tuple(run.locator for run in self.runs)
        if len(locators) != len(set(locators)):
            raise ValueError("style run locators must be unique")
        expected = canonical_payload_sha256(self.canonical_payload())
        if self.sidecar_sha256 not in {"", expected}:
            raise ValueError("sidecar hash does not match its canonical payload")
        if not self.sidecar_sha256:
            object.__setattr__(self, "sidecar_sha256", expected)
        return self


class StyledTextRunExtractor:
    """Offline PPTX style reader. Malformed style XML fails closed per run."""

    def extract(
        self,
        path: Path,
        *,
        source_id: str,
        source_sha256: str | None = None,
    ) -> StyledTextRunSidecar:
        payload_sha256 = hashlib.sha256(path.read_bytes()).hexdigest()
        if source_sha256 is not None and payload_sha256 != source_sha256:
            raise ValueError("PowerPoint SHA-256 does not match its immutable source pin")
        source_sha256 = source_sha256 or payload_sha256
        presentation = Presentation(str(path))
        records: list[StyledTextRun] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            scheme = _theme_scheme(slide.part)
            clr_map = _clr_map(slide.part)
            for shape, path_parts in walk_shapes_with_paths(slide.shapes):
                shape_path = ".".join(map(str, path_parts))
                if getattr(shape, "has_table", False):
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        for column_index, cell in enumerate(row.cells, start=1):
                            records.extend(
                                _paragraph_records(
                                    cell.text_frame.paragraphs,
                                    f"slide:{slide_number}:shape:{shape_path}:cell:{row_index}:{column_index}",
                                    source_id, source_sha256, slide_number, scheme, clr_map,
                                )
                            )
                elif getattr(shape, "has_text_frame", False):
                    records.extend(
                        _paragraph_records(
                            shape.text_frame.paragraphs,
                            f"slide:{slide_number}:shape:{shape_path}",
                            source_id, source_sha256, slide_number, scheme, clr_map,
                        )
                    )
            try:
                records.extend(
                    _paragraph_records(
                        slide.notes_slide.notes_text_frame.paragraphs,
                        f"slide:{slide_number}:notes",
                        source_id, source_sha256, slide_number,
                        _theme_scheme(slide.notes_slide.part), _clr_map(slide.notes_slide.part),
                    )
                )
            except (AttributeError, KeyError):
                pass
        return StyledTextRunSidecar(
            source_id=source_id,
            source_sha256=source_sha256,
            parser_version=RUN_STYLE_PARSER_VERSION,
            runs=tuple(records),
        )


def extract_styled_text_run_sidecar(
    path: Path, *, source_id: str, source_sha256: str | None = None
) -> StyledTextRunSidecar:
    return StyledTextRunExtractor().extract(
        path, source_id=source_id, source_sha256=source_sha256
    )


def normalized_text_sha256(text: str) -> str:
    return hashlib.sha256(" ".join(text.split()).encode("utf-8")).hexdigest()


def matches_policy_color(run: StyledTextRun, colors: Any) -> bool:
    """Return true only for a supported resolved color / effective theme match."""
    if not run.resolved_color:
        return False
    for color in colors:
        rgb = getattr(color, "rgb", None)
        theme_ref = getattr(color, "theme_ref", None)
        rgb_match = rgb is not None and _rgb_delta(run.resolved_color, rgb) <= 8
        theme_match = theme_ref is not None and run.theme_color == str(theme_ref).casefold()
        if rgb_match or theme_match:
            return True
    return False


def _paragraph_records(
    paragraphs: Any, locator_prefix: str, source_id: str, source_sha256: str,
    slide_number: int, scheme: dict[str, str], clr_map: dict[str, str],
) -> list[StyledTextRun]:
    records: list[StyledTextRun] = []
    for paragraph_index, paragraph in enumerate(paragraphs, start=1):
        for run_index, run in enumerate(paragraph.runs, start=1):
            explicit, theme, resolved, brightness, attempted, inherited = _run_color(
                run, scheme, clr_map
            )
            records.append(
                StyledTextRun(
                    source_id=source_id, source_sha256=source_sha256,
                    slide_number=slide_number,
                    locator=f"{locator_prefix}:p:{paragraph_index}:r:{run_index}",
                    text=run.text or "",
                    normalized_text_sha256=normalized_text_sha256(run.text or ""),
                    explicit_rgb=explicit, theme_color=theme, resolved_color=resolved,
                    brightness=brightness, color_attempted=attempted,
                    inherited_color_resolution=inherited,
                    bold=run.font.bold, italic=run.font.italic,
                    underline=_optional_bool(run.font.underline),
                    highlight=_highlight(run),
                )
            )
    return records


def _run_color(
    run: Any, scheme: dict[str, str], clr_map: dict[str, str]
) -> tuple[str | None, str | None, str | None, float | None, bool, str | None]:
    try:
        r_pr = run._r.rPr
        if r_pr is None:
            return None, None, None, None, False, None
        solid = next((child for child in r_pr if _local(child) == "solidFill"), None)
        if solid is None or len(solid) != 1:
            return None, None, None, None, False, None
        color = solid[0]
        kind = _local(color)
        transforms = tuple(_local(child) for child in color)
        brightness, supported = _brightness(color, transforms)
        if not supported:
            return None, None, None, None, True, "unresolved"
        if kind == "srgbClr":
            value = _hex(color.get("val"))
            if value is None:
                return None, None, None, None, True, "unresolved"
            return value, None, _apply_brightness(value, brightness), brightness, True, None
        if kind == "sysClr":
            value = _hex(color.get("lastClr"))
            if value is None:
                return None, None, None, None, True, "unresolved"
            return value, None, _apply_brightness(value, brightness), brightness, True, None
        if kind == "schemeClr":
            reference = str(color.get("val") or "").casefold()
            effective = clr_map.get(reference, reference)
            value = scheme.get(effective)
            if not reference or value is None:
                return None, effective or None, None, brightness, True, "unresolved"
            return None, effective, _apply_brightness(value, brightness), brightness, True, None
        return None, None, None, None, True, "unresolved"
    except (AttributeError, KeyError, TypeError, ValueError):
        return None, None, None, None, True, "unresolved"


def _brightness(color: Any, transforms: tuple[str, ...]) -> tuple[float | None, bool]:
    if not transforms:
        return None, True
    if any(transform not in {"lumMod", "lumOff"} for transform in transforms):
        return None, False
    if transforms.count("lumMod") > 1 or transforms.count("lumOff") > 1:
        return None, False
    values = { _local(child): child.get("val") for child in color }
    try:
        mod = int(values.get("lumMod", "100000"))
        off = int(values.get("lumOff", "0"))
    except (TypeError, ValueError):
        return None, False
    if not 0 <= mod <= 100000 or not 0 <= off <= 100000:
        return None, False
    if "lumOff" in values:
        if mod + off != 100000:
            return None, False
        return off / 100000, True
    return mod / 100000 - 1, True


def _apply_brightness(value: str | None, brightness: float | None) -> str | None:
    if value is None:
        return None
    if brightness is None:
        return value
    red, green, blue = (int(value[index:index + 2], 16) / 255 for index in range(0, 6, 2))
    hue, lightness, saturation = colorsys.rgb_to_hls(red, green, blue)
    if brightness >= 0:
        lightness = lightness + brightness * (1 - lightness)
    else:
        lightness = lightness * (1 + brightness)
    rgb = colorsys.hls_to_rgb(hue, lightness, saturation)
    return "".join(f"{round(channel * 255):02X}" for channel in rgb)


def _theme_scheme(part: Any) -> dict[str, str]:
    theme = _related_theme(part)
    if theme is None:
        return {}
    try:
        root = parse_xml(theme.blob)
        scheme = next(node for node in root.iter() if _local(node) == "clrScheme")
        colors: dict[str, str] = {}
        for entry in scheme:
            if len(entry):
                child = entry[0]
                kind = _local(child)
                if kind == "srgbClr":
                    value = _hex(child.get("val"))
                elif kind == "sysClr":
                    value = _hex(child.get("lastClr"))
                else:
                    continue
                if len(child):
                    continue
                if value:
                    colors[_local(entry)] = value
        return colors
    except (AttributeError, StopIteration, TypeError, ValueError):
        return {}


def _related_theme(start: Any) -> Any | None:
    queue: deque[Any] = deque([start])
    seen: set[str] = set()
    allowed = {RT.SLIDE_LAYOUT, RT.SLIDE_MASTER, RT.NOTES_MASTER}
    while queue:
        part = queue.popleft()
        name = str(getattr(part, "partname", id(part)))
        if name in seen:
            continue
        seen.add(name)
        for relation in part.rels.values():
            if relation.reltype == RT.THEME:
                return relation.target_part
            if relation.reltype in allowed and relation.target_part is not None:
                queue.append(relation.target_part)
    return None


def _clr_map(part: Any) -> dict[str, str]:
    """Take the nearest override, otherwise the master mapping, then identity."""
    queue: deque[Any] = deque([part])
    seen: set[str] = set()
    allowed = {RT.SLIDE_LAYOUT, RT.SLIDE_MASTER, RT.NOTES_MASTER}
    while queue:
        current = queue.popleft()
        name = str(getattr(current, "partname", id(current)))
        if name in seen:
            continue
        seen.add(name)
        mapping = _part_clr_map(current)
        if mapping is not None:
            return mapping
        for relation in current.rels.values():
            if relation.reltype in allowed and relation.target_part is not None:
                queue.append(relation.target_part)
    return {}


def _part_clr_map(part: Any) -> dict[str, str] | None:
    try:
        for node in part._element.iter():
            if _local(node) in {"overrideClrMapping", "clrMap"}:
                return {
                    str(key).casefold(): str(value).casefold()
                    for key, value in node.attrib.items()
                }
    except AttributeError:
        return None
    return None


def _hex(value: object) -> str | None:
    text = str(value or "").strip().upper()
    if len(text) != 6 or not all(character in "0123456789ABCDEF" for character in text):
        return None
    return text


def _local(node: Any) -> str:
    return str(node.tag).rsplit("}", 1)[-1]


def _optional_bool(value: object) -> bool | None:
    return value if isinstance(value, bool) else None


def _highlight(run: Any) -> str | None:
    try:
        r_pr = run._r.rPr
        if r_pr is None:
            return None
        highlight = next((child for child in r_pr if _local(child) == "highlight"), None)
        if highlight is None or not len(highlight):
            return None
        color = highlight[0]
        return _hex(color.get("lastClr") if _local(color) == "sysClr" else color.get("val"))
    except AttributeError:
        return None


def _rgb_delta(left: str, right: str) -> int:
    return max(
        abs(int(left[index : index + 2], 16) - int(str(right)[index : index + 2], 16))
        for index in range(0, 6, 2)
    )
