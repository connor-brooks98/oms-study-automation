import hashlib
import re
from collections.abc import Sequence
from dataclasses import dataclass, replace
from io import BytesIO
from pathlib import Path
from typing import Literal, Protocol

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel, ConfigDict, Field, model_validator
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from oms_hub.anki.contracts import canonical_payload_sha256
from oms_hub.anki.course_policy import CourseCurationPolicy
from oms_hub.anki.domain import SourceKind
from oms_hub.document_processing.run_styles import (
    StyledTextRunSidecar,
    matches_policy_color,
    normalized_text_sha256,
)
from oms_hub.ingestion.domain import StudyRevision, UploadKind
from oms_hub.study_generation.domain import OutlineRecord

EXTRACTION_VERSION = "source-passages-v1"
ExtractionStatus = Literal["text", "vision", "vision_unavailable"]

_SPACE = re.compile(r"\s+")
_SENTENCE_BOUNDARY = re.compile(r"(?<=[.!?])\s+")
_TIME = r"(?:\d{1,2}:)?\d{1,2}:\d{2}(?:[.,]\d+)?"
_TIMESTAMP_RANGE = re.compile(
    rf"^\s*(?P<start>{_TIME})\s*-->\s*(?P<end>{_TIME})\s*$"
)
_TIMESTAMP_PREFIX = re.compile(
    rf"^\s*\[?(?P<start>{_TIME})\]?\s+(?P<text>.+?)\s*$"
)
_SUMMARY_CITATION = re.compile(r"\[([0-9,\s]+)\]")
SummarySection = Literal["core", "depth", "emphasis"]
_SUMMARY_HEADINGS: dict[str, SummarySection] = {
    "CORE CONCEPTS": "core",
    "DEPTH MAP": "depth",
    "PROFESSOR EMPHASIS FLAGS": "emphasis",
}


@dataclass(frozen=True, slots=True)
class SourcePassage:
    passage_id: str
    source_id: str
    revision_id: int
    lecture_id: int
    artifact_id: str
    source_kind: SourceKind
    locator: str
    text: str
    content_hash: str
    extraction_status: ExtractionStatus
    slide_number: int | None = None
    start_seconds: float | None = None
    end_seconds: float | None = None
    summary_backrefs: tuple[str, ...] = ()
    summary_section: SummarySection | None = None

    @classmethod
    def create(
        cls,
        *,
        revision_id: int,
        lecture_id: int,
        artifact_id: str,
        source_kind: SourceKind,
        locator: str,
        text: str,
        extraction_status: ExtractionStatus = "text",
        slide_number: int | None = None,
        start_seconds: float | None = None,
        end_seconds: float | None = None,
        source_id: str | None = None,
        summary_backrefs: tuple[str, ...] = (),
        summary_section: SummarySection | None = None,
    ) -> "SourcePassage":
        normalized = _normalize_text(text)
        if revision_id <= 0 or lecture_id <= 0:
            raise ValueError("source identifiers must be positive")
        if not artifact_id.strip() or not locator.strip():
            raise ValueError("source artifact and locator cannot be blank")
        if not normalized and extraction_status != "vision_unavailable":
            raise ValueError("source passage text cannot be blank")
        if slide_number is not None and slide_number <= 0:
            raise ValueError("slide number must be positive")
        if start_seconds is not None and start_seconds < 0:
            raise ValueError("transcript start cannot be negative")
        if end_seconds is not None and (
            end_seconds < 0
            or (
                start_seconds is not None
                and end_seconds < start_seconds
            )
        ):
            raise ValueError("transcript end is invalid")
        content_hash = hashlib.sha256(normalized.encode("utf-8")).hexdigest()
        identity = "\0".join(
            (
                EXTRACTION_VERSION,
                str(revision_id),
                source_kind.value,
                locator.strip(),
                extraction_status,
                content_hash,
            )
        )
        return cls(
            passage_id=hashlib.sha256(identity.encode("utf-8")).hexdigest(),
            source_id=source_id or _default_source_id(
                lecture_id=lecture_id,
                source_kind=source_kind,
                locator=locator,
                slide_number=slide_number,
            ),
            revision_id=revision_id,
            lecture_id=lecture_id,
            artifact_id=artifact_id.strip(),
            source_kind=source_kind,
            locator=locator.strip(),
            text=normalized,
            content_hash=content_hash,
            extraction_status=extraction_status,
            slide_number=slide_number,
            start_seconds=start_seconds,
            end_seconds=end_seconds,
            summary_backrefs=summary_backrefs,
            summary_section=summary_section,
        )

    @property
    def citation(self) -> str:
        if self.slide_number is not None:
            suffix = (
                " speaker notes"
                if self.source_kind is SourceKind.SPEAKER_NOTES
                else ""
            )
            return (
                f"Lecture {self.lecture_id}, slide "
                f"{self.slide_number}{suffix}"
            )
        if self.start_seconds is not None:
            interval = _format_time(self.start_seconds)
            if self.end_seconds is not None:
                interval = (
                    f"{interval}\N{EN DASH}{_format_time(self.end_seconds)}"
                )
            return f"Lecture {self.lecture_id}, transcript {interval}"
        if self.source_kind is SourceKind.SUMMARY:
            return f"Lecture {self.lecture_id}, NotebookLM summary"
        return f"Lecture {self.lecture_id}, transcript"


class SourceEmphasisEvidence(BaseModel):
    """Policy-specific, additive evidence projected from an immutable style sidecar."""

    model_config = ConfigDict(extra="forbid", frozen=True)

    evidence_id: str = ""
    source_id: str = Field(min_length=1)
    source_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    sidecar_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    locator: str = Field(min_length=1)
    text: str = Field(min_length=1)
    normalized_text_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    emphasis_kind: Literal["colored_text"] = "colored_text"
    normalized_color: str = Field(pattern=r"^[0-9A-F]{6}$")
    policy_match: bool = True
    policy_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    provenance_hash: str = ""

    def canonical_payload(self) -> dict[str, object]:
        return self.model_dump(mode="json", exclude={"evidence_id", "provenance_hash"})

    @model_validator(mode="after")
    def _validate_identity(self) -> "SourceEmphasisEvidence":
        if not self.text.strip():
            raise ValueError("source emphasis evidence text cannot be blank")
        if not self.policy_match:
            raise ValueError("source emphasis evidence must be a policy match")
        if self.normalized_text_sha256 != normalized_text_sha256(self.text):
            raise ValueError("emphasis normalized text hash does not match its text")
        provenance = canonical_payload_sha256(self.canonical_payload())
        evidence_id = hashlib.sha256(
            ("source-emphasis-evidence-v1\0" + provenance).encode("utf-8")
        ).hexdigest()
        if self.provenance_hash not in {"", provenance}:
            raise ValueError("emphasis provenance hash does not match its canonical payload")
        if self.evidence_id not in {"", evidence_id}:
            raise ValueError("emphasis evidence ID does not match its canonical payload")
        if not self.provenance_hash:
            object.__setattr__(self, "provenance_hash", provenance)
        if not self.evidence_id:
            object.__setattr__(self, "evidence_id", evidence_id)
        return self


def project_source_emphasis_evidence(
    sidecar: StyledTextRunSidecar, policy: CourseCurationPolicy
) -> tuple[SourceEmphasisEvidence, ...]:
    """Select only nonblank colored runs matching this frozen policy."""
    if policy.emphasis_mode not in {"colored_text", "combined"}:
        return ()
    evidence: list[SourceEmphasisEvidence] = []
    for run in sidecar.runs:
        if not run.text.strip() or not matches_policy_color(run, policy.emphasis_colors):
            continue
        resolved_color = run.resolved_color
        if resolved_color is None:
            continue
        evidence.append(
            SourceEmphasisEvidence(
                source_id=run.source_id,
                source_sha256=run.source_sha256,
                sidecar_sha256=sidecar.sidecar_sha256,
                locator=run.locator,
                text=" ".join(run.text.split()),
                normalized_text_sha256=run.normalized_text_sha256,
                normalized_color=resolved_color,
                policy_sha256=policy.policy_sha256,
            )
        )
    return tuple(evidence)


class SummaryMalformedError(ValueError):
    """Raised when a NotebookLM outline is missing required V2 sections."""


class NotebookSummaryParser:
    def parse(self, record: OutlineRecord) -> list[SourcePassage]:
        if not record.path.is_file():
            raise FileNotFoundError(record.path)
        payload = record.path.read_bytes()
        digest = hashlib.sha256(payload).hexdigest()
        if digest != record.sha256:
            raise ValueError("NotebookLM summary SHA-256 does not match its pin")
        try:
            text = "\n".join(
                page.extract_text() or ""
                for page in PdfReader(BytesIO(payload)).pages
            )
        except PdfReadError as exc:
            raise SummaryMalformedError(
                "NotebookLM summary is not a readable PDF"
            ) from exc
        return self.parse_text(
            lecture_id=record.lecture_id,
            outline_id=record.id,
            text=text,
        )

    def parse_text(
        self,
        *,
        lecture_id: int,
        outline_id: int,
        text: str,
    ) -> list[SourcePassage]:
        sections: dict[SummarySection, list[str]] = {
            "core": [],
            "depth": [],
            "emphasis": [],
        }
        current: SummarySection | None = None
        seen: set[SummarySection] = set()
        for raw_line in text.splitlines():
            line = raw_line.strip().lstrip("#").strip()
            heading = _SUMMARY_HEADINGS.get(line.upper())
            if heading is not None:
                current = heading
                seen.add(heading)
                continue
            if current is None or not line:
                continue
            cleaned = re.sub(
                r"^(?:[-*•\x7f]|\d+[.)])\s*",
                "",
                line,
            ).strip()
            if cleaned:
                sections[current].append(cleaned)

        missing = [
            heading
            for heading, section in _SUMMARY_HEADINGS.items()
            if section in {"depth", "emphasis"} and section not in seen
        ]
        if missing:
            raise SummaryMalformedError(
                "NotebookLM summary is missing required section(s): "
                + ", ".join(missing)
            )

        passages: list[SourcePassage] = []
        prefixes = {"core": "CORE", "depth": "DEPTH", "emphasis": "EMPH"}
        counters = {"core": 0, "depth": 0, "emphasis": 0}
        for section in ("core", "depth", "emphasis"):
            for item in sections[section]:
                counters[section] += 1
                number = counters[section]
                suffix = (
                    f"{number:02d}"
                    if section == "core"
                    else f"{'D' if section == 'depth' else 'E'}{number}"
                )
                backrefs = tuple(
                    value
                    for match in _SUMMARY_CITATION.findall(item)
                    for value in re.findall(r"\d+", match)
                )
                passages.append(
                    SourcePassage.create(
                        revision_id=outline_id,
                        lecture_id=lecture_id,
                        artifact_id=f"outline:{outline_id}",
                        source_kind=SourceKind.SUMMARY,
                        locator=f"summary:{section}:{number}",
                        text=item,
                        source_id=(
                            f"SUM:{lecture_id}:{prefixes[section]}:{suffix}"
                        ),
                        summary_backrefs=backrefs,
                        summary_section=section,
                    )
                )
        return passages


class RevisionRepository(Protocol):
    def get_study_revision(
        self,
        revision_id: int,
    ) -> StudyRevision: ...

    def imported_derived_audit_matches(self, revision: StudyRevision) -> bool: ...

    def has_imported_derived_audit(self, revision_id: int) -> bool: ...


class OutlineRepository(Protocol):
    def outline(self, outline_id: int) -> OutlineRecord | None: ...


class SlideVisionExtractor(Protocol):
    def describe_slide(
        self,
        presentation_path: Path,
        slide_number: int,
    ) -> str | None: ...


@dataclass(frozen=True, slots=True)
class _TranscriptUnit:
    text: str
    start_seconds: float | None
    end_seconds: float | None


class LectureSourceExtractor:
    def __init__(
        self,
        revisions: RevisionRepository,
        *,
        outlines: OutlineRepository | None = None,
        summary_parser: NotebookSummaryParser | None = None,
        vision: SlideVisionExtractor | None = None,
        transcript_max_chars: int = 800,
        transcript_overlap_sentences: int = 1,
    ) -> None:
        if transcript_max_chars < 50:
            raise ValueError("transcript passage size is too small")
        if transcript_overlap_sentences < 0:
            raise ValueError("transcript overlap cannot be negative")
        self.revisions = revisions
        self.outlines = outlines
        self.summary_parser = summary_parser or NotebookSummaryParser()
        self.vision = vision
        self.transcript_max_chars = transcript_max_chars
        self.transcript_overlap_sentences = (
            transcript_overlap_sentences
        )

    def extract(
        self,
        revision_ids: Sequence[int],
        *,
        summary_outline_id: int | None = None,
    ) -> list[SourcePassage]:
        if len(set(revision_ids)) != len(revision_ids):
            raise ValueError("source revision IDs must be unique")
        passages: list[SourcePassage] = []
        for revision_id in revision_ids:
            revision = self.revisions.get_study_revision(revision_id)
            if (
                revision.provenance_kind == "imported_derived"
                or self.revisions.has_imported_derived_audit(revision.id)
            ) and not self.revisions.imported_derived_audit_matches(revision):
                raise ValueError("imported-derived slide provenance is no longer ready")
            if revision.kind is UploadKind.SLIDES:
                passages.extend(self._extract_slides(revision))
            elif revision.kind is UploadKind.TRANSCRIPTS:
                passages.extend(self._extract_transcript(revision))
            else:
                raise ValueError(
                    f"unsupported source revision kind: {revision.kind}"
                )
        if summary_outline_id is not None:
            if self.outlines is None:
                raise ValueError("summary outline repository is not configured")
            outline = self.outlines.outline(summary_outline_id)
            if outline is None:
                raise KeyError(summary_outline_id)
            lecture_ids = {passage.lecture_id for passage in passages}
            if lecture_ids and lecture_ids != {outline.lecture_id}:
                raise ValueError("summary outline belongs to another lecture")
            passages.extend(self.summary_parser.parse(outline))
        return passages

    def _extract_slides(
        self,
        revision: StudyRevision,
    ) -> list[SourcePassage]:
        path = revision.immutable_source_path
        if not path.is_file():
            raise FileNotFoundError(path)
        presentation = Presentation(str(path))
        passages: list[SourcePassage] = []
        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            text = _normalize_text(
                " ".join(
                    shape.text
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                    and getattr(shape, "text", "").strip()
                )
            )
            if text:
                passages.append(
                    SourcePassage.create(
                        revision_id=revision.id,
                        lecture_id=revision.lecture_id,
                        artifact_id=revision.upload_item_id,
                        source_kind=SourceKind.SLIDE,
                        locator=f"slide:{slide_number}",
                        text=text,
                        slide_number=slide_number,
                    )
                )
            notes = ""
            if slide.has_notes_slide:
                notes = _normalize_text(
                    slide.notes_slide.notes_text_frame.text
                )
            if notes:
                passages.append(
                    SourcePassage.create(
                        revision_id=revision.id,
                        lecture_id=revision.lecture_id,
                        artifact_id=revision.upload_item_id,
                        source_kind=SourceKind.SPEAKER_NOTES,
                        locator=f"slide:{slide_number}:notes",
                        text=notes,
                        slide_number=slide_number,
                    )
                )
            has_image = any(
                shape.shape_type is MSO_SHAPE_TYPE.PICTURE
                for shape in slide.shapes
            )
            if not text and has_image:
                description = (
                    self.vision.describe_slide(path, slide_number)
                    if self.vision is not None
                    else None
                )
                passages.append(
                    SourcePassage.create(
                        revision_id=revision.id,
                        lecture_id=revision.lecture_id,
                        artifact_id=revision.upload_item_id,
                        source_kind=SourceKind.VISION,
                        locator=f"slide:{slide_number}:image",
                        text=description or "",
                        extraction_status=(
                            "vision"
                            if description and description.strip()
                            else "vision_unavailable"
                        ),
                        slide_number=slide_number,
                    )
                )
        return passages

    def _extract_transcript(
        self,
        revision: StudyRevision,
    ) -> list[SourcePassage]:
        path = (
            revision.immutable_derived_path
            if revision.immutable_derived_path is not None
            and revision.immutable_derived_path.is_file()
            else revision.immutable_source_path
        )
        raw_text = path.read_text(encoding="utf-8")
        units = _transcript_units(raw_text)
        windows = _overlapping_windows(
            units,
            max_chars=self.transcript_max_chars,
            overlap=self.transcript_overlap_sentences,
        )
        passages: list[SourcePassage] = []
        for position, window in enumerate(windows, start=1):
            start = next(
                (
                    unit.start_seconds
                    for unit in window
                    if unit.start_seconds is not None
                ),
                None,
            )
            end = next(
                (
                    unit.end_seconds
                    for unit in reversed(window)
                    if unit.end_seconds is not None
                ),
                None,
            )
            time_locator = (
                f"{_compact_time(start)}-{_compact_time(end)}"
                if start is not None
                else "untimed"
            )
            passages.append(
                SourcePassage.create(
                    revision_id=revision.id,
                    lecture_id=revision.lecture_id,
                    artifact_id=revision.upload_item_id,
                    source_kind=SourceKind.TRANSCRIPT,
                    locator=f"transcript:{position}:{time_locator}",
                    text=" ".join(unit.text for unit in window),
                    start_seconds=start,
                    end_seconds=end,
                )
            )
        return passages


def _transcript_units(value: str) -> list[_TranscriptUnit]:
    units: list[_TranscriptUnit] = []
    pending_range: tuple[float, float] | None = None
    for raw_line in value.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        range_match = _TIMESTAMP_RANGE.fullmatch(line)
        if range_match is not None:
            pending_range = (
                _parse_time(range_match.group("start")),
                _parse_time(range_match.group("end")),
            )
            continue
        if line.isdigit() and pending_range is None:
            continue
        prefix_match = _TIMESTAMP_PREFIX.fullmatch(line)
        if prefix_match is not None:
            start = _parse_time(prefix_match.group("start"))
            end = None
            text = prefix_match.group("text")
        elif pending_range is not None:
            start, end = pending_range
            pending_range = None
            text = line
        else:
            start = None
            end = None
            text = line
        sentences = [
            _normalize_text(sentence)
            for sentence in _SENTENCE_BOUNDARY.split(text)
            if _normalize_text(sentence)
        ]
        units.extend(
            _TranscriptUnit(
                text=sentence,
                start_seconds=start,
                end_seconds=end,
            )
            for sentence in sentences
        )
    for index, unit in enumerate(units[:-1]):
        next_start = units[index + 1].start_seconds
        if unit.end_seconds is None and next_start is not None:
            units[index] = replace(unit, end_seconds=next_start)
    return units


def _overlapping_windows(
    units: Sequence[_TranscriptUnit],
    *,
    max_chars: int,
    overlap: int,
) -> list[tuple[_TranscriptUnit, ...]]:
    windows: list[tuple[_TranscriptUnit, ...]] = []
    start = 0
    while start < len(units):
        end = start
        characters = 0
        while end < len(units):
            added = len(units[end].text) + (1 if end > start else 0)
            if end > start and characters + added > max_chars:
                break
            characters += added
            end += 1
        windows.append(tuple(units[start:end]))
        if end >= len(units):
            break
        start = max(start + 1, end - overlap)
    return windows


def _parse_time(value: str) -> float:
    parts = value.replace(",", ".").split(":")
    seconds = float(parts[-1])
    minutes = int(parts[-2])
    hours = int(parts[-3]) if len(parts) == 3 else 0
    return hours * 3_600 + minutes * 60 + seconds


def _format_time(value: float) -> str:
    total = int(value)
    hours, remainder = divmod(total, 3_600)
    minutes, seconds = divmod(remainder, 60)
    if hours:
        return f"{hours:02d}:{minutes:02d}:{seconds:02d}"
    return f"{minutes:02d}:{seconds:02d}"


def _compact_time(value: float | None) -> str:
    return "open" if value is None else str(int(value))


def _default_source_id(
    *,
    lecture_id: int,
    source_kind: SourceKind,
    locator: str,
    slide_number: int | None,
) -> str:
    if source_kind in {
        SourceKind.SLIDE,
        SourceKind.SPEAKER_NOTES,
        SourceKind.VISION,
    }:
        number = slide_number or 0
        suffix = {
            SourceKind.SLIDE: "",
            SourceKind.SPEAKER_NOTES: ":NOTES",
            SourceKind.VISION: ":IMAGE",
        }[source_kind]
        return f"SLD:{lecture_id}:{number:04d}{suffix}"
    if source_kind is SourceKind.TRANSCRIPT:
        match = re.match(r"transcript:(\d+):", locator)
        position = int(match.group(1)) if match is not None else 0
        return f"TRX:{lecture_id}:{position:04d}"
    return f"SUM:{lecture_id}:UNKNOWN:00"


def _normalize_text(value: str) -> str:
    return _SPACE.sub(" ", value).strip()
