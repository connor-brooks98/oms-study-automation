import asyncio
import hashlib
import json
from datetime import UTC, datetime
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import numpy as np
import pytest
from pptx import Presentation

from oms_hub.anki.cost_estimator import FrozenRateTable, ModelRate
from oms_hub.anki.course_policy import CourseCurationPolicy
from oms_hub.anki.domain import (
    CreateCurationJob,
    CurationStage,
    CurationState,
    PipelineContractVersion,
    ResolvedModelConfiguration,
    ResolvedStageModel,
    ReviewChangeSet,
)
from oms_hub.anki.index import AnkiIndex
from oms_hub.anki.normalize import NormalizedNote, semantic_text
from oms_hub.anki.pipeline import (
    CurationPipeline,
    PinnedInputChanged,
    StageArtifactStore,
    StageContext,
)
from oms_hub.anki.provider_attempts import (
    current_provider_attempt_identity,
    provider_replay_identity_document,
)
from oms_hub.anki.rehearsal.structured import (
    ReplayStructuredTextGenerator,
    structured_request_key,
)
from oms_hub.anki.rehearsal.vectors import ReplayEmbeddingClient
from oms_hub.anki.repository import AnkiCurationRepository
from oms_hub.anki.semantic.domain import DocumentRecord
from oms_hub.anki.semantic.service import SemanticIndexService, content_hash
from oms_hub.anki.semantic.store import SemanticSnapshotStore
from oms_hub.anki.sources import LectureSourceExtractor
from oms_hub.anki.stages import (
    CurationServicesRunner,
    PinnedCurationInputValidator,
    revision_fingerprint,
)
from oms_hub.anki.worker import AnkiCurationWorker
from oms_hub.db import Database
from oms_hub.ingestion.domain import StudyRevision, UploadKind
from oms_hub.llm.domain import (
    DEFAULT_GENERATION_OPTIONS,
    GeneratedText,
    GenerationOptions,
    ProviderName,
)
from oms_hub.llm.structured import StructuredTextService
from oms_hub.models import LectureModel


class _OfflineStructuredGenerator:
    """Offline-only normal StructuredTextService delegate for the R3/R7 contract."""

    offline_replay_only = True

    def __init__(self) -> None:
        self.calls: list[dict[str, object]] = []
        self.manifest: dict[str, dict[str, object]] = {}

    def generate_text(self, instruction: str, input_text: str, **kwargs: Any) -> GeneratedText:
        self.calls.append({"instruction": instruction, "input_text": input_text, **kwargs})
        request = json.loads(input_text)
        if "source_bundle" in request:
            evidence_id = request["source_bundle"]["evidence"][0]["evidence_id"]
            response: dict[str, object] = {
                "concepts": [
                    {
                        "canonical_statement": "Iron deficiency anemia",
                        "primary_entity": "Iron deficiency anemia",
                        "aliases": [],
                        "exact_terms": ["Ferritin"],
                        "depth_tier": 1,
                        "priority": 100,
                        "facts": [
                            {
                                "statement": "Ferritin is low in iron deficiency anemia.",
                                "evidence_ids": [evidence_id],
                                "generation_allowed": True,
                                "forbidden_cloze_targets": [],
                            }
                        ],
                        "source_evidence_ids": [evidence_id],
                        "retrieval_queries": ["Ferritin low"],
                    }
                ]
            }
        else:
            bundles = request["bundles"]
            response = {
                "rows": [
                    {
                        "bundle_id": bundle["bundle_id"],
                        "candidate_id": bundle["candidate"]["candidate_id"],
                        "disposition": "keep",
                        "confidence_bps": 10000,
                        "supporting_passage_ids": [bundle["selected_passages"][0]["passage_id"]],
                        "conflicting_passage_ids": [],
                        "redundant_with_candidate_id": None,
                        "reason": "The existing card states the supported fact.",
                    }
                    for bundle in bundles
                ]
            }
        generated = GeneratedText(
            text=json.dumps(response),
            provider=ProviderName.OPENAI,
            model=str(kwargs["model"]),
            request_id=f"offline-{len(self.calls)}",
            input_tokens=11,
            output_tokens=7,
            cost_microusd=0,
        )
        identity = current_provider_attempt_identity()
        assert identity is not None
        options = kwargs.get("options", DEFAULT_GENERATION_OPTIONS)
        assert isinstance(options, GenerationOptions)
        key = structured_request_key(
            instruction,
            input_text,
            output_schema=kwargs["output_schema"],  # type: ignore[arg-type]
            provider=kwargs["provider"],  # type: ignore[arg-type]
            model=str(kwargs["model"]),
            options=options,
            attempt_identity=identity,
        )
        self.manifest[key] = {
            "text": generated.text,
            "text_sha256": hashlib.sha256(generated.text.encode()).hexdigest(),
            "provider": generated.provider.value,
            "model": generated.model,
            "request_id": generated.request_id,
            "input_tokens": generated.input_tokens,
            "output_tokens": generated.output_tokens,
            "cost_microusd": generated.cost_microusd,
            "attempt_identity": provider_replay_identity_document(identity),
        }
        return generated


class _RevisionRepository:
    def __init__(self, revisions: dict[int, StudyRevision]) -> None:
        self.revisions = revisions

    def get_study_revision(self, revision_id: int) -> StudyRevision:
        return self.revisions[revision_id]

    def has_imported_derived_audit(self, revision_id: int) -> bool:
        del revision_id
        return False

    def imported_derived_audit_matches(self, revision: StudyRevision) -> bool:
        raise AssertionError(f"unexpected imported-derived audit check: {revision.id}")


def _revision(
    revision_id: int,
    lecture_id: int,
    kind: UploadKind,
    source: Path,
    *,
    derived: Path | None = None,
) -> StudyRevision:
    return StudyRevision(
        id=revision_id,
        upload_item_id=f"upload-{revision_id}",
        lecture_id=lecture_id,
        kind=kind,
        source_sha256=hashlib.sha256(source.read_bytes()).hexdigest(),
        immutable_source_path=source,
        derived_sha256=None
        if derived is None
        else hashlib.sha256(derived.read_bytes()).hexdigest(),
        immutable_derived_path=derived,
        canonical_source_path=None,
        canonical_derived_path=None,
        icloud_path=None,
        prompt_sha256=None,
        state="accepted",
        current=True,
    )


def _presentation(path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Iron deficiency anemia"
    slide.placeholders[1].text = "Ferritin is low."
    deck.save(path)


def _note() -> NormalizedNote:
    text = "Ferritin is low in iron deficiency anemia."
    return NormalizedNote(
        note_id=1,
        model_name="Basic",
        text=text,
        extra="",
        raw_fields={"Text": text},
        tags=("OMS::Heme",),
        card_ids=(101,),
        media=(),
        token_signature="ferritin iron deficiency anemia low",
        content_sha256=hashlib.sha256(text.encode()).hexdigest(),
        deck_names=("Deck",),
    )


def test_v3_r0_rejects_live_looking_clients_before_any_provider_call() -> None:
    class LiveLookingGenerator:
        calls = 0

        def generate_text(self, *_args: object, **_kwargs: object) -> GeneratedText:
            self.calls += 1
            raise AssertionError("live structured generator must not be called")

    generator = LiveLookingGenerator()
    runner = object.__new__(CurationServicesRunner)
    runner.structured = StructuredTextService(generator)
    runner.embedder = SimpleNamespace(offline_replay_only=False)
    runner.semantic = SimpleNamespace(embedder=runner.embedder)
    runner.repository = SimpleNamespace(
        get_policy_by_sha256=lambda _sha: (_ for _ in ()).throw(
            AssertionError("R0 must stop first")
        )
    )
    context = StageContext(
        job=SimpleNamespace(offline_replay_only=True, policy_sha256="a" * 64),
        stage=CurationStage.V3_R0_PREFLIGHT,
        input_sha256="a" * 64,
        prior_artifacts=(),
        prior_payloads={},
    )

    with pytest.raises(PinnedInputChanged, match="offline-only"):
        asyncio.run(runner.run(context))
    assert generator.calls == 0


def test_v3_resumed_stage_rejects_replaced_live_client_before_artifact_lookup() -> None:
    class LiveLookingGenerator:
        calls = 0

        def generate_text(self, *_args: object, **_kwargs: object) -> GeneratedText:
            self.calls += 1
            raise AssertionError("live structured generator must not be called")

    class UnreadableArtifacts(dict[CurationStage, dict[str, object]]):
        def __getitem__(self, key: CurationStage) -> dict[str, object]:
            raise AssertionError(f"R7 must not read {key.value} before the offline guard")

    generator = LiveLookingGenerator()
    runner = object.__new__(CurationServicesRunner)
    runner.structured = StructuredTextService(generator)
    runner.embedder = SimpleNamespace(offline_replay_only=False)
    runner.semantic = SimpleNamespace(embedder=SimpleNamespace(offline_replay_only=True))
    context = StageContext(
        job=SimpleNamespace(
            pipeline_contract_version=PipelineContractVersion.CARD_CENTRIC_V3,
            offline_replay_only=True,
        ),
        stage=CurationStage.V3_R7_CLASSIFICATION,
        input_sha256="a" * 64,
        prior_artifacts=(),
        prior_payloads=UnreadableArtifacts(),
    )

    with pytest.raises(PinnedInputChanged, match="offline-only"):
        asyncio.run(runner.run(context))
    assert generator.calls == 0


def test_v3_live_execution_requires_the_capture_repository_and_capture_only_clients() -> None:
    runner = object.__new__(CurationServicesRunner)
    runner.structured = SimpleNamespace(generator=SimpleNamespace(capture_only=True))
    runner.embedder = SimpleNamespace(capture_only=True)
    runner.semantic = SimpleNamespace(embedder=SimpleNamespace(capture_only=True))
    runner.repository = SimpleNamespace(allows_v3_live_capture=lambda: True)
    context = StageContext(
        job=SimpleNamespace(offline_replay_only=False),
        stage=CurationStage.V3_R0_PREFLIGHT,
        input_sha256="a" * 64,
        prior_artifacts=(),
        prior_payloads={},
    )

    runner._require_v3_offline_execution(context)  # noqa: SLF001
    runner.repository = SimpleNamespace(allows_v3_live_capture=lambda: False)
    with pytest.raises(PinnedInputChanged, match="capture-only"):
        runner._require_v3_offline_execution(context)  # noqa: SLF001


def test_v3_offline_golden_runner_reaches_review_with_durable_provider_evidence(
    tmp_path: Path,
) -> None:
    database = Database(f"sqlite:///{tmp_path / 'hub.db'}")
    database.migrate()
    with database.session() as session:
        lecture = LectureModel(
            subject="Heme", exam_number=1, lecture_number=1, topic="Anemia", lecturer="P"
        )
        session.add(lecture)
        session.flush()
        lecture_id = lecture.id
    slides = tmp_path / "lecture.pptx"
    transcript = tmp_path / "transcript.txt"
    _presentation(slides)
    transcript.write_text(
        "[00:00] Professor emphasis: ferritin is low in iron deficiency anemia.", encoding="utf-8"
    )
    revisions = _RevisionRepository(
        {
            1: _revision(1, lecture_id, UploadKind.SLIDES, slides),
            2: _revision(2, lecture_id, UploadKind.TRANSCRIPTS, transcript, derived=transcript),
        }
    )
    companion = AnkiIndex(tmp_path / "companion")
    note = _note()
    companion.rebuild_companion([note], snapshot_id="companion-1", fingerprint="a" * 64)
    replay = ReplayEmbeddingClient(tmp_path / "replay", model="fixture-embedding", dimensions=2)
    replay.seed(
        (
            "Ferritin",
            "Iron deficiency anemia",
            "Ferritin is low in iron deficiency anemia.",
            "transcript emphasis",
            "Ferritin low",
        ),
        input_type="query",
        vectors=np.asarray(((1.0, 0.0),) * 5, dtype=np.float32),
    )
    semantic_store = SemanticSnapshotStore(tmp_path / "semantic")
    semantic_text_value = semantic_text(note)
    manifest = semantic_store.replace(
        [DocumentRecord(note.note_id, semantic_text_value, content_hash(semantic_text_value))],
        np.asarray(((1.0, 0.0),), dtype=np.float32),
        model="fixture-embedding",
    )
    semantic = SemanticIndexService(
        semantic_store,
        replay,
        model="fixture-embedding",
        dimensions=2,
        min_coverage=1.0,
        query_cache_size=16,
    )
    repository = AnkiCurationRepository(database)
    policy = CourseCurationPolicy(
        policy_id="fixture",
        revision=1,
        course_id="heme",
        professor_label="P",
        scope_instruction="Use transcript emphasis.",
        emphasis_mode="transcript_emphasis",
        missing_emphasis_fallback="block",
        tag_scope_mode="disabled",
        classification_strictness="strict",
        generation_style_profile="cloze",
        ordinary_cost_limit_microusd=1_000_000,
        hard_stop_cost_limit_microusd=1_000_000,
    )
    repository.create_policy_revision(policy)
    route = ResolvedStageModel("openai", "fixture", thinking_mode="disabled")
    config = ResolvedModelConfiguration(
        "fixture",
        route,
        route,
        route,
        route,
        scope_r3=route,
        cheap_classify_r7=route,
        thorough_classify_r7=route,
        generation_r9=route,
    )
    table = FrozenRateTable(
        (ModelRate("fixture", 1, 1, 1, 1, 1), ModelRate("fixture-embedding", 1, 1, 1, 1, 1)),
        datetime(2026, 8, 17, tzinfo=UTC),
        "fixture",
    )
    request = CreateCurationJob(
        lecture_id=lecture_id,
        block_id=None,
        source_revision_ids=(1, 2),
        deck_allowlist=("Deck",),
        tag_allowlist=("OMS::Heme",),
        instruction_text="",
        target_deck="Deck",
        target_tag="Tag",
        index_snapshot_id="companion-1",
        lcl_prompt_version="lcl",
        judgment_rubric_version="judge",
        gap_prompt_version="gap",
        provider="openai",
        model="fixture",
        pipeline_contract_version=PipelineContractVersion.CARD_CENTRIC_V3,
        resolved_model_config=config,
        policy_sha256=policy.policy_sha256,
        rate_table_document=table.document(),
        offline_replay_only=True,
        semantic_generation=str(manifest.generation),
        companion_generation="companion-1",
        source_revision_hashes={
            revision_id: revision_fingerprint(revision)
            for revision_id, revision in revisions.revisions.items()
        },
    )
    job = repository.create_job(request)

    def source_indexes(_lecture_id: int) -> object:
        raise AssertionError("R0-R11 must not build a source index")

    input_validator = PinnedCurationInputValidator(
        repository,
        revisions,
        companion,
        semantic_store,
        source_indexes,
        semantic_model="fixture-embedding",
        semantic_dimensions=2,
    )
    generator = _OfflineStructuredGenerator()
    runner = CurationServicesRunner(
        runtime=SimpleNamespace(),
        repository=repository,
        source_extractor=LectureSourceExtractor(revisions),
        source_indexes=source_indexes,
        companion=companion,
        semantic=semantic,
        structured=StructuredTextService(generator),
        embedder=replay,
        focused_retrieval_limit=4,
        global_retrieval_limit=4,
        llm_settings=SimpleNamespace(),
    )
    artifact_store = StageArtifactStore(tmp_path / "artifacts")
    pipeline = CurationPipeline(repository, artifact_store, runner, input_validator=input_validator)
    worker = AnkiCurationWorker(repository, pipeline, worker_id="golden-worker")
    while repository.require_job(job.id).state is not CurationState.READY_FOR_REVIEW:
        assert asyncio.run(worker.run_once())

    artifacts = repository.list_stage_artifacts(job.id)
    events = repository.list_provider_attempt_events(job.id)
    assert repository.require_job(job.id).state is CurationState.READY_FOR_REVIEW
    assert {artifact.stage for artifact in artifacts} == {
        CurationStage.V3_R0_PREFLIGHT,
        CurationStage.V3_R1_SOURCE_INDEX,
        CurationStage.V3_R2_FIDELITY,
        CurationStage.V3_R3_SCOPE,
        CurationStage.V3_R4_INDEX_VERIFICATION,
        CurationStage.V3_R5_RETRIEVAL,
        CurationStage.V3_R6_CALIBRATION,
        CurationStage.V3_R7_CLASSIFICATION,
        CurationStage.V3_R8_GAP_CONFIRMATION,
        CurationStage.V3_R9_GENERATION,
        CurationStage.V3_R10_DEDUPE,
        CurationStage.V3_R11_REVIEW,
    }
    assert len(generator.calls) == 2
    assert not asyncio.run(worker.run_once())
    candidates = repository.list_candidates(job.id)
    assert [(item.note_id, item.content_hash, item.selected) for item in candidates] == [
        (1, note.content_sha256, True)
    ]
    assert repository.list_gap_cards(job.id) == []
    r11_artifact = next(item for item in artifacts if item.stage is CurationStage.V3_R11_REVIEW)
    r11 = artifact_store.read(r11_artifact, job=repository.require_job(job.id))
    with pytest.raises(ValueError, match="invisible"):
        repository.save_review(
            job.id,
            ReviewChangeSet(expected_revision=0, candidate_selections={999: True}),
            card_centric_snapshot=r11["snapshot"],
            v3_review_artifact_sha256=r11["artifact_sha256"],
            v3_cost_ledger_sha256=r11["cost_ledger_sha256"],
        )
    assert repository.list_candidates(job.id)[0].selected is True
    saved = repository.save_review(
        job.id,
        ReviewChangeSet(expected_revision=0, candidate_selections={1: False}),
        card_centric_snapshot=r11["snapshot"],
        v3_review_artifact_sha256=r11["artifact_sha256"],
        v3_cost_ledger_sha256=r11["cost_ledger_sha256"],
    )
    assert saved.revision == 1
    assert repository.list_candidates(job.id)[0].selected is False
    assert replay.evidence.document_replay_hits == replay.evidence.live_document_calls == 0
    assert replay.evidence.live_query_calls == replay.evidence.replay_misses == 0
    assert replay.evidence.query_replay_hits > 0
    assert all(row["event"] != "transport_failed" for row in events)
    for call_events in _events_by_call(events).values():
        assert [row["event"] for row in call_events] == [
            "begun",
            "dispatched",
            "response_received",
            "accepted",
        ]
        assert all(row["cost_reservation"] is not None for row in call_events)

    structured_manifest = tmp_path / "structured.json"
    structured_manifest.write_text(json.dumps(generator.manifest), encoding="utf-8")
    assert {
        record["attempt_identity"]["stage"]  # type: ignore[index]
        for record in generator.manifest.values()
    } == {CurationStage.V3_R3_SCOPE.value, CurationStage.V3_R7_CLASSIFICATION.value}
    assert all(
        len(record["attempt_identity"]["stage_input_sha256"]) == 64  # type: ignore[index]
        for record in generator.manifest.values()
    )
    replay_generator = ReplayStructuredTextGenerator(
        structured_manifest,
        require_attempt_identity=True,
    )
    first_r5 = artifact_store.read(
        next(item for item in artifacts if item.stage is CurationStage.V3_R5_RETRIEVAL),
        job=repository.require_job(job.id),
    )
    first_query_hits = replay.evidence.query_replay_hits
    replay_runner = CurationServicesRunner(
        runtime=SimpleNamespace(),
        repository=repository,
        source_extractor=LectureSourceExtractor(revisions),
        source_indexes=source_indexes,
        companion=companion,
        semantic=semantic,
        structured=StructuredTextService(replay_generator),
        embedder=replay,
        focused_retrieval_limit=4,
        global_retrieval_limit=4,
        llm_settings=SimpleNamespace(),
    )
    replay_job = repository.create_job(request)
    replay_artifact_store = StageArtifactStore(tmp_path / "replay-artifacts")
    replay_pipeline = CurationPipeline(
        repository,
        replay_artifact_store,
        replay_runner,
        input_validator=input_validator,
    )
    replay_worker = AnkiCurationWorker(
        repository,
        replay_pipeline,
        worker_id="replay-golden-worker",
    )
    while repository.require_job(replay_job.id).state is not CurationState.READY_FOR_REVIEW:
        assert asyncio.run(replay_worker.run_once())

    assert replay_job.id != job.id
    assert replay_generator.evidence.hits == 2
    assert replay_generator.evidence.misses == replay_generator.evidence.live_calls == 0
    assert repository.require_job(replay_job.id).state is CurationState.READY_FOR_REVIEW
    assert not asyncio.run(replay_worker.run_once())
    second_r5 = replay_artifact_store.read(
        next(
            item
            for item in repository.list_stage_artifacts(replay_job.id)
            if item.stage is CurationStage.V3_R5_RETRIEVAL
        ),
        job=repository.require_job(replay_job.id),
    )
    assert first_r5["cost_ledger"] != second_r5["cost_ledger"]
    assert replay.evidence.query_replay_hits == first_query_hits
    assert {artifact.stage for artifact in repository.list_stage_artifacts(replay_job.id)} == {
        artifact.stage for artifact in artifacts
    }
    database.close()


def _events_by_call(
    events: list[dict[str, object]],
) -> dict[tuple[object, ...], list[dict[str, object]]]:
    grouped: dict[tuple[object, ...], list[dict[str, object]]] = {}
    for row in events:
        grouped.setdefault(
            (row["stage"], row["stage_attempt"], row["call_index"], row["subcall_ordinal"]), []
        ).append(row)
    return grouped
