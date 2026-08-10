import hashlib
from contextlib import contextmanager
from dataclasses import dataclass, replace
from pathlib import Path
from types import SimpleNamespace
from uuid import UUID

import pytest
from pptx import Presentation

from oms_hub import cli
from oms_hub.anki.card_centric import build_source_index
from oms_hub.anki.domain import CreateCurationJob, CurationStage, PipelineContractVersion
from oms_hub.anki.pipeline import PinnedInputChanged
from oms_hub.anki.repository import AnkiCurationRepository
from oms_hub.anki.sources import LectureSourceExtractor
from oms_hub.anki.stages import PinnedCurationInputValidator, revision_fingerprint
from oms_hub.artifact_writes import ArtifactWriteClaimLost
from oms_hub.config import Settings
from oms_hub.db import Database
from oms_hub.domain import LectureKey
from oms_hub.existing_artifact_import import (
    ExistingArtifactImporter,
    ExistingArtifactImportRequest,
)
from oms_hub.ingestion.domain import StudyRevision, UploadKind
from oms_hub.ingestion.repository import IngestionRepository
from oms_hub.models import (
    ExistingArtifactImportModel,
    GenerationJobModel,
    OutlineOutputModel,
    OutlineReplacementReviewModel,
    StudyRevisionModel,
    UploadBatchModel,
    UploadItemModel,
)
from oms_hub.repositories import CatalogRepository, LectureInput
from oms_hub.study_generation.domain import GenerationKind, NotebookAnswer, SourceKind
from oms_hub.study_generation.outline import OutlinePdfRenderer, OutlineService
from oms_hub.study_generation.repository import (
    GenerationRepository,
    ImportedOutlineReplacementRequired,
)
from oms_hub.web.anki_routes import _outline_ready_for_curation


@dataclass(frozen=True)
class ImportedBundle:
    database: Database
    settings: Settings
    lecture_id: int
    slide_revision_id: int
    imported_transcript_id: int
    outline_id: int
    import_id: str
    transcript_path: Path
    outline_path: Path
    immutable_transcript_path: Path
    immutable_outline_path: Path


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _write_slides(path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Synaptic transmission"
    slide.placeholders[1].text = "Calcium influx triggers vesicle fusion."
    deck.save(path)


@pytest.fixture
def imported_bundle(tmp_path: Path) -> ImportedBundle:
    settings = Settings(
        _env_file=None,
        data_dir=tmp_path / "data",
        database_url=f"sqlite:///{tmp_path / 'hub.db'}",
        study_root=tmp_path / "study",
    )
    database = Database(settings.database_url)
    database.migrate()
    lecture_id = CatalogRepository(database).upsert_lecture(
        LectureInput("Neuro", 1, 24, "Synapse", "Professor", None)
    )
    slides = settings.data_dir / "artifacts" / "v2" / "slides" / "slides.pptx"
    slides.parent.mkdir(parents=True, exist_ok=True)
    _write_slides(slides)
    canonical_slides = settings.study_root / "Neuro" / "slides.pptx"
    canonical_slides.parent.mkdir(parents=True, exist_ok=True)
    canonical_slides.write_bytes(slides.read_bytes())
    slide_sha = _sha256(slides)
    slide_pdf = settings.data_dir / "artifacts" / "v2" / "slides" / "slides.pdf"
    slide_pdf.write_bytes(OutlinePdfRenderer().render("Slides", "# CORE CONCEPTS\n- source"))
    canonical_slide_pdf = settings.study_root / "Neuro" / "slides.pdf"
    canonical_slide_pdf.write_bytes(slide_pdf.read_bytes())
    slide_pdf_sha = _sha256(slide_pdf)
    with database.session() as session:
        session.add(UploadBatchModel(id="slides-batch", kind="slides", state="complete"))
        session.flush()
        session.add(
            UploadItemModel(
                id="slides-item",
                batch_id="slides-batch",
                kind="slides",
                original_filename=slides.name,
                staged_path=str(slides),
                sha256=slide_sha,
                size_bytes=slides.stat().st_size,
                state="complete",
                lecture_id=lecture_id,
                confidence=1,
                manual_assignment=True,
            )
        )
        session.flush()
        slide = StudyRevisionModel(
            upload_item_id="slides-item",
            lecture_id=lecture_id,
            kind="slides",
            source_sha256=slide_sha,
            immutable_source_path=str(slides),
            derived_sha256=slide_pdf_sha,
            immutable_derived_path=str(slide_pdf),
            canonical_source_path=str(canonical_slides),
            canonical_derived_path=str(canonical_slide_pdf),
            state="current",
            current=True,
        )
        session.add(slide)
        session.flush()
        slide_revision_id = slide.id
    transcript = tmp_path / "cleaned.txt"
    transcript.write_text("Calcium influx triggers synaptic vesicle fusion.", encoding="utf-8")
    outline = tmp_path / "outline.pdf"
    outline.write_bytes(
        OutlinePdfRenderer().render(
            "Synapse outline",
            "# CORE CONCEPTS\n- Calcium triggers vesicle fusion [1]\n\n"
            "# DEPTH MAP\n- DEEP: SNARE proteins [2]\n\n"
            "# PROFESSOR EMPHASIS FLAGS\n- Repeated: calcium entry [3]",
        )
    )
    imported = ExistingArtifactImporter(database, settings).import_artifacts(
        ExistingArtifactImportRequest(
            lecture_id,
            slide_revision_id,
            slide_sha,
            slide_pdf_sha,
            transcript,
            _sha256(transcript),
            outline,
            _sha256(outline),
        )
    )
    yield ImportedBundle(
        database,
        settings,
        lecture_id,
        slide_revision_id,
        imported.transcript_revision_id,
        imported.outline_id,
        imported.import_id,
        imported.transcript_path,
        imported.outline_path,
        imported.immutable_transcript_path,
        imported.immutable_outline_path,
    )
    database.close()


class _Companion:
    def snapshot_id(self) -> str:
        return "companion-imported"


class _Semantic:
    generation = UUID("4438eabc-3da1-4d6d-a6af-2302de092f8e")

    def load(self, **_: object) -> object:
        return SimpleNamespace(manifest=SimpleNamespace(generation=self.generation))


def _queued_job(bundle: ImportedBundle):
    revisions = IngestionRepository(bundle.database)
    slide = revisions.get_study_revision(bundle.slide_revision_id)
    transcript = revisions.get_study_revision(bundle.imported_transcript_id)
    outline = GenerationRepository(bundle.database).outline(bundle.outline_id)
    assert outline is not None
    job = AnkiCurationRepository(bundle.database).create_job(
        CreateCurationJob(
            lecture_id=bundle.lecture_id,
            block_id="synapse-block",
            source_revision_ids=(slide.id, transcript.id),
            source_revision_hashes={
                slide.id: revision_fingerprint(slide),
                transcript.id: revision_fingerprint(transcript),
            },
            summary_outline_id=outline.id,
            summary_outline_sha256=outline.sha256,
            deck_allowlist=("AnKing Step Deck",),
            tag_allowlist=("#AK_Step2_v12::Neurology",),
            instruction_text="",
            target_deck="OMS::Neuro::Synapse",
            target_tag="AnkiHub_Optional::LMU_OMS_II::Neuro::Synapse",
            index_snapshot_id="companion-imported",
            lcl_prompt_version="lcl-v1",
            judgment_rubric_version="judgment-v1",
            gap_prompt_version="gap-v1",
            provider="anthropic",
            model="claude-sonnet-5",
            pipeline_contract_version=PipelineContractVersion.CARD_CENTRIC_V2,
            semantic_generation=str(_Semantic.generation),
            companion_generation="companion-imported",
        )
    )
    return job, revisions, GenerationRepository(bundle.database)


def _validator(bundle: ImportedBundle) -> tuple[PinnedCurationInputValidator, object]:
    job, revisions, outlines = _queued_job(bundle)
    validator = PinnedCurationInputValidator(
        AnkiCurationRepository(bundle.database),
        revisions,
        _Companion(),  # type: ignore[arg-type]
        _Semantic(),  # type: ignore[arg-type]
        lambda _job_id: (_ for _ in ()).throw(AssertionError("source index is not pinned")),
        outlines=outlines,
        semantic_model="voyage-4-large",
        semantic_dimensions=1024,
    )
    return validator, job


def _failed_imported_outline_job(
    bundle: ImportedBundle, repository: GenerationRepository, answer: str
):
    job = repository.queue(bundle.lecture_id, GenerationKind.OUTLINE)
    revisions = IngestionRepository(bundle.database)
    slide = revisions.get_study_revision(bundle.slide_revision_id)
    transcript = revisions.get_study_revision(bundle.imported_transcript_id)
    mapping = repository.save_notebook_mapping(
        "Neuro", "neuro", 1, "notebook-imported-outline", "Neuro Exam 1"
    )
    repository.bind_source(
        mapping.id,
        bundle.lecture_id,
        slide.id,
        SourceKind.LECTURE_PDF,
        slide.derived_sha256 or "",
        "slides-source",
        "slides",
    )
    repository.bind_source(
        mapping.id,
        bundle.lecture_id,
        transcript.id,
        SourceKind.CLEANED_TRANSCRIPT,
        transcript.derived_sha256 or "",
        "transcript-source",
        "transcript",
    )
    with bundle.database.session() as session:
        stored = session.get(GenerationJobModel, job.id)
        assert stored is not None
        stored.state = "failed"
        stored.stage = "pdf"
        stored.error = "current imported outline requires a durable replacement review decision"
        stored.notebook_answer = answer
        stored.notebook_id = mapping.remote_notebook_id
        stored.pdf_source_id = "slides-source"
        stored.transcript_source_id = "transcript-source"
        stored.pdf_revision_id = bundle.slide_revision_id
        stored.transcript_revision_id = bundle.imported_transcript_id
    return repository.get(job.id)


def _mark_imported_outline_job_running(bundle: ImportedBundle, job_id: str) -> None:
    with bundle.database.session() as session:
        stored = session.get(GenerationJobModel, job_id)
        assert stored is not None
        stored.state = "running"


def test_imported_bundle_uses_ordinary_anki_source_snapshot_and_replay(
    imported_bundle: ImportedBundle,
) -> None:
    job, revisions, outlines = _queued_job(imported_bundle)
    transcript = revisions.get_study_revision(imported_bundle.imported_transcript_id)
    outline = outlines.outline(imported_bundle.outline_id)
    assert outline is not None
    assert transcript.provenance_kind == "imported_cleaned"
    assert transcript.import_id == outline.import_id == imported_bundle.import_id
    assert outline.provenance_kind == "imported_notebooklm"
    assert outline.slide_revision_id == imported_bundle.slide_revision_id
    assert outline.transcript_revision_id == transcript.id
    extractor = LectureSourceExtractor(revisions, outlines=outlines)
    passages = extractor.extract(job.source_revision_ids, summary_outline_id=job.summary_outline_id)
    snapshot = build_source_index(
        passages,
        snapshot_id=job.index_snapshot_id,
        source_revision_hashes=job.source_revision_hashes,
        summary_outline_sha256=job.summary_outline_sha256,
    )
    replay = AnkiCurationRepository(imported_bundle.database).prepare_stage_replay_inputs(
        job.id,
        stage=CurationStage.PREFLIGHT,
    )
    replay_again = AnkiCurationRepository(imported_bundle.database).prepare_stage_replay_inputs(
        job.id,
        stage=CurationStage.PREFLIGHT,
    )
    assert job.pipeline_contract_version is PipelineContractVersion.CARD_CENTRIC_V2
    assert job.source_revision_hashes[transcript.id] == revision_fingerprint(transcript)
    assert snapshot.source_revision_hashes == job.source_revision_hashes
    assert snapshot.summary_outline_sha256 == outline.sha256
    assert (
        snapshot.source_sha256
        == build_source_index(
            passages,
            snapshot_id=job.index_snapshot_id,
            source_revision_hashes=job.source_revision_hashes,
            summary_outline_sha256=job.summary_outline_sha256,
        ).source_sha256
    )
    assert {passage.revision_id for passage in passages} >= {transcript.id, outline.id}
    assert replay.canonical_json == replay_again.canonical_json
    assert replay.sha256 == replay_again.sha256
    validator, pinned_job = _validator(imported_bundle)
    validator.validate(pinned_job.id)


@pytest.mark.parametrize(
    "tamper",
    [
        "immutable-transcript-file",
        "canonical-transcript-file",
        "immutable-outline-file",
        "canonical-outline-file",
        "transcript-provenance",
        "outline-provenance",
        "outline-import-id",
        "outline-link",
        "outline-link-hash",
        "outline-slide-hash",
        "transcript-current",
        "outline-current",
    ],
)
def test_imported_anki_pins_fail_closed_on_identity_tampering(imported_bundle, tamper):
    validator, job = _validator(imported_bundle)
    if tamper.endswith("file"):
        {
            "immutable-transcript-file": imported_bundle.immutable_transcript_path,
            "canonical-transcript-file": imported_bundle.transcript_path,
            "immutable-outline-file": imported_bundle.immutable_outline_path,
            "canonical-outline-file": imported_bundle.outline_path,
        }[tamper].write_bytes(b"tampered")
    else:
        with imported_bundle.database.session() as session:
            transcript = session.get(StudyRevisionModel, imported_bundle.imported_transcript_id)
            outline = session.get(OutlineOutputModel, imported_bundle.outline_id)
            assert transcript is not None and outline is not None
            if tamper == "transcript-provenance":
                transcript.provenance_kind = "llm_cleaned"
            elif tamper == "outline-provenance":
                outline.provenance_kind = "notebooklm_generated"
            elif tamper == "outline-import-id":
                # Keep the relational graph valid so the pinned-input
                # validator, rather than SQLite, proves the logical import
                # identity mismatch is fail-closed.
                slide = session.get(StudyRevisionModel, imported_bundle.slide_revision_id)
                assert slide is not None and slide.derived_sha256 is not None
                session.add(
                    ExistingArtifactImportModel(
                        id="other-import",
                        bundle_sha256="f" * 64,
                        lecture_id=imported_bundle.lecture_id,
                        slide_revision_id=imported_bundle.slide_revision_id,
                        slide_source_sha256=slide.source_sha256,
                        slide_pdf_sha256=slide.derived_sha256,
                        transcript_sha256=transcript.source_sha256,
                        outline_sha256=outline.sha256,
                        subject="Neuro",
                        exam_number=1,
                        lecture_number=24,
                        topic="Synapse",
                        status="complete",
                        transcript_revision_id=imported_bundle.imported_transcript_id,
                        outline_id=imported_bundle.outline_id,
                    )
                )
                outline.import_id = "other-import"
            elif tamper == "outline-link":
                outline.transcript_revision_id = None
            elif tamper == "outline-link-hash":
                outline.transcript_sha256 = "0" * 64
            elif tamper == "outline-slide-hash":
                outline.slide_sha256 = "0" * 64
            elif tamper == "transcript-current":
                transcript.current = False
            else:
                outline.current = False
    with pytest.raises(PinnedInputChanged):
        validator.validate(job.id)


def test_imported_outline_readiness_requires_complete_links(
    imported_bundle: ImportedBundle,
) -> None:
    revisions = IngestionRepository(imported_bundle.database)
    outlines = GenerationRepository(imported_bundle.database)
    outline = outlines.outline(imported_bundle.outline_id)
    assert outline is not None
    assert _outline_ready_for_curation(outline, revisions)
    with imported_bundle.database.session() as session:
        stored = session.get(OutlineOutputModel, imported_bundle.outline_id)
        assert stored is not None
        stored.transcript_revision_id = None
    incomplete = outlines.outline(imported_bundle.outline_id)
    assert incomplete is not None
    assert not _outline_ready_for_curation(incomplete, revisions)


def test_ordinary_outline_record_keeps_generation_provenance_and_job_id(tmp_path: Path) -> None:
    database = Database(f"sqlite:///{tmp_path / 'hub.db'}")
    database.migrate()
    lecture_id = CatalogRepository(database).upsert_lecture(
        LectureInput("Neuro", 1, 1, "Normal", "", None)
    )
    path = tmp_path / "normal.pdf"
    path.write_bytes(OutlinePdfRenderer().render("Normal", "# CORE CONCEPTS\n- One"))
    repository = GenerationRepository(database)
    job = repository.queue(lecture_id, GenerationKind.OUTLINE)
    record = repository.record_outline(lecture_id, job.id, path, _sha256(path))
    assert record.job_id == job.id
    assert record.provenance_kind == "notebooklm_generated"
    assert record.import_id is None
    database.close()


def test_imported_outline_requires_explicit_replacement_and_rolls_back_copy(
    imported_bundle: ImportedBundle, monkeypatch: pytest.MonkeyPatch
) -> None:
    repository = GenerationRepository(imported_bundle.database)
    lecture = LectureKey("Neuro", 1, 24, "Synapse")
    answer = NotebookAnswer("# CORE CONCEPTS\n- Replacement")
    job = _failed_imported_outline_job(imported_bundle, repository, answer.text)

    class _StableRenderer:
        payload = OutlinePdfRenderer().render("Stable", answer.text)

        def render(self, _title: str, _content: str) -> bytes:
            return self.payload

    service = OutlineService(imported_bundle.settings, repository, _StableRenderer())
    before = imported_bundle.outline_path.read_bytes()

    with pytest.raises(ImportedOutlineReplacementRequired):
        service.file(job, lecture, answer)
    assert imported_bundle.outline_path.read_bytes() == before
    current = repository.current_outline(imported_bundle.lecture_id)
    assert current is not None and current.id == imported_bundle.outline_id

    def reject_record(*_args: object, **_kwargs: object) -> object:
        raise RuntimeError("database write failed")

    monkeypatch.setattr(cli, "Settings", lambda: imported_bundle.settings)
    monkeypatch.setattr(cli, "Database", lambda _url: imported_bundle.database)
    args = cli.build_parser().parse_args(
        [
            "approve-imported-outline-replacement",
            "--lecture-id",
            str(imported_bundle.lecture_id),
            "--generation-job-id",
            job.id,
            "--operator",
            "operator@example.test",
            "--reason",
            "Reviewed replacement",
            "--confirm",
        ]
    )
    assert args.handler(args) == 0
    review = repository.imported_outline_replacement_review(imported_bundle.lecture_id, job.id)
    assert review is not None
    _mark_imported_outline_job_running(imported_bundle, job.id)
    job = repository.get(job.id)
    rollback = imported_bundle.outline_path.with_name(
        f".{imported_bundle.outline_path.name}.rollback-{job.id}"
    )
    monkeypatch.setattr(repository, "record_outline", reject_record)
    with pytest.raises(RuntimeError, match="database write failed"):
        service.file(job, lecture, answer, replacement_review=review)
    assert imported_bundle.outline_path.read_bytes() == before
    assert not rollback.exists()
    current = repository.current_outline(imported_bundle.lecture_id)
    assert current is not None and current.id == imported_bundle.outline_id

    monkeypatch.undo()
    record = service.file(job, lecture, answer, replacement_review=review)
    assert record.job_id == job.id
    assert record.provenance_kind == "notebooklm_generated"
    assert repository.current_outline(imported_bundle.lecture_id) == record
    assert not rollback.exists()
    repository.complete(job.id)
    # An approved replacement makes the imported outline historical without
    # invalidating its immutable provenance graph at the next startup.
    imported_bundle.database.migrate()
    later_job = repository.queue(imported_bundle.lecture_id, GenerationKind.OUTLINE)
    later_path = imported_bundle.outline_path.with_name("later-generated-outline.pdf")
    later_path.write_bytes(OutlinePdfRenderer().render("Later", answer.text))
    later = repository.record_outline(
        imported_bundle.lecture_id, later_job.id, later_path, _sha256(later_path)
    )
    repository.complete(later_job.id)
    imported_bundle.database.migrate()
    assert repository.outline(record.id).current is False  # type: ignore[union-attr]
    assert repository.current_outline(imported_bundle.lecture_id) == later


@pytest.mark.parametrize(
    ("field", "value"),
    (
        ("state", "queued"),
        ("state", "running"),
        ("state", "paused"),
        ("state", "complete"),
        ("stage", "notebook"),
        ("notebook_answer", None),
        ("pdf_revision_id", None),
        ("transcript_revision_id", None),
        ("pdf_source_id", None),
        ("transcript_source_id", None),
        ("notebook_id", "wrong-notebook"),
        ("pdf_source_id", "wrong-pdf-source"),
        ("transcript_source_id", "wrong-transcript-source"),
    ),
)
def test_replacement_approval_cli_rejects_ineligible_job_without_review_write(
    imported_bundle: ImportedBundle,
    monkeypatch: pytest.MonkeyPatch,
    field: str,
    value: object,
) -> None:
    repository = GenerationRepository(imported_bundle.database)
    job = _failed_imported_outline_job(imported_bundle, repository, "# CORE CONCEPTS\n- One")
    with imported_bundle.database.session() as session:
        stored = session.get(GenerationJobModel, job.id)
        assert stored is not None
        setattr(stored, field, value)
    monkeypatch.setattr(cli, "Database", lambda _url: imported_bundle.database)
    args = cli.build_parser().parse_args(
        [
            "approve-imported-outline-replacement",
            "--lecture-id",
            str(imported_bundle.lecture_id),
            "--generation-job-id",
            job.id,
            "--operator",
            "operator@example.test",
            "--reason",
            "Reviewed replacement",
            "--confirm",
        ]
    )

    assert args.handler(args) == 2
    assert (
        repository.imported_outline_replacement_review(imported_bundle.lecture_id, job.id)
        is None
    )
    rejected = repository.get(job.id)
    assert getattr(rejected, field) == value


def test_outline_claim_loss_preserves_successor_and_rollback_evidence(
    imported_bundle: ImportedBundle, monkeypatch: pytest.MonkeyPatch
) -> None:
    repository = GenerationRepository(imported_bundle.database)
    lecture = LectureKey("Neuro", 1, 24, "Synapse")
    answer = NotebookAnswer("# CORE CONCEPTS\n- Replacement")
    job = _failed_imported_outline_job(imported_bundle, repository, answer.text)
    review = repository.approve_imported_outline_replacement(
        imported_bundle.lecture_id, job.id, "operator@example.test", "Reviewed replacement"
    )
    _mark_imported_outline_job_running(imported_bundle, job.id)
    job = repository.get(job.id)
    before = imported_bundle.outline_path.read_bytes()
    successor = b"successor canonical bytes"

    class _StableRenderer:
        payload = OutlinePdfRenderer().render("Stable", answer.text)

        def render(self, _title: str, _content: str) -> bytes:
            return self.payload

    class _LostClaim:
        calls = 0

        def assert_owned(self) -> None:
            self.calls += 1
            if self.calls >= 3:
                imported_bundle.outline_path.write_bytes(successor)
                raise ArtifactWriteClaimLost("successor owns the writer claim")

    class _LostCoordinator:
        def __init__(self, *_args: object) -> None:
            self.claim_state = _LostClaim()

        @contextmanager
        def claim(self, *_args: object):
            yield self.claim_state

    import oms_hub.study_generation.outline as outline_module

    monkeypatch.setattr(outline_module, "ArtifactWriteCoordinator", _LostCoordinator)
    service = OutlineService(imported_bundle.settings, repository, _StableRenderer())
    with pytest.raises(ArtifactWriteClaimLost, match="successor owns"):
        service.file(job, lecture, answer, replacement_review=review)

    rollback = imported_bundle.outline_path.with_name(
        f".{imported_bundle.outline_path.name}.rollback-{job.id}"
    )
    assert imported_bundle.outline_path.read_bytes() == successor
    assert rollback.read_bytes() == before
    current = repository.current_outline(imported_bundle.lecture_id)
    assert current is not None and current.id == imported_bundle.outline_id


def test_approved_outline_replacement_requires_running_pdf_job_at_promotion(
    imported_bundle: ImportedBundle,
) -> None:
    repository = GenerationRepository(imported_bundle.database)
    answer = NotebookAnswer("# CORE CONCEPTS\n- Replacement")
    job = _failed_imported_outline_job(imported_bundle, repository, answer.text)
    review = repository.approve_imported_outline_replacement(
        imported_bundle.lecture_id, job.id, "operator@example.test", "Reviewed replacement"
    )
    lecture = LectureKey("Neuro", 1, 24, "Synapse")
    service = OutlineService(imported_bundle.settings, repository, OutlinePdfRenderer())

    with pytest.raises(ImportedOutlineReplacementRequired, match="no longer eligible"):
        service.file(job, lecture, answer, replacement_review=review)
    assert repository.current_outline(imported_bundle.lecture_id).id == imported_bundle.outline_id  # type: ignore[union-attr]


def test_promotion_revalidates_audit_outline_edge_after_approval(
    imported_bundle: ImportedBundle,
) -> None:
    repository = GenerationRepository(imported_bundle.database)
    answer = NotebookAnswer("# CORE CONCEPTS\n- Replacement")
    job = _failed_imported_outline_job(imported_bundle, repository, answer.text)
    review = repository.approve_imported_outline_replacement(
        imported_bundle.lecture_id, job.id, "operator@example.test", "Reviewed replacement"
    )
    _mark_imported_outline_job_running(imported_bundle, job.id)
    job = repository.get(job.id)
    with imported_bundle.database.session() as session:
        audit = session.get(ExistingArtifactImportModel, imported_bundle.import_id)
        assert audit is not None
        audit.outline_id = None
    before = imported_bundle.outline_path.read_bytes()
    lecture = LectureKey("Neuro", 1, 24, "Synapse")
    service = OutlineService(imported_bundle.settings, repository, OutlinePdfRenderer())

    with pytest.raises(ImportedOutlineReplacementRequired, match="no longer eligible"):
        repository.assert_outline_replacement_allowed(
            imported_bundle.lecture_id, job.id, replacement_review=review
        )
    with pytest.raises(ImportedOutlineReplacementRequired, match="no longer eligible"):
        service.file(job, lecture, answer, replacement_review=review)
    assert imported_bundle.outline_path.read_bytes() == before
    assert repository.current_outline(imported_bundle.lecture_id).id == imported_bundle.outline_id  # type: ignore[union-attr]


def test_review_trigger_rejects_corrupted_imported_outline_identity(
    imported_bundle: ImportedBundle,
) -> None:
    repository = GenerationRepository(imported_bundle.database)
    job = _failed_imported_outline_job(imported_bundle, repository, "# CORE CONCEPTS\n- One")
    with imported_bundle.database.session() as session:
        imported = session.get(OutlineOutputModel, imported_bundle.outline_id)
        assert imported is not None
        imported.slide_sha256 = "0" * 64
    with pytest.raises(Exception, match="outline replacement review identity is invalid"):
        with imported_bundle.database.session() as session:
            session.add(
                OutlineReplacementReviewModel(
                    generation_job_id=job.id,
                    lecture_id=imported_bundle.lecture_id,
                    import_id=imported_bundle.import_id,
                    operator="operator@example.test",
                    reason="Reviewed replacement",
                )
            )
            session.flush()


def test_historical_ordinary_revision_fingerprint_is_unchanged() -> None:
    historical = StudyRevision(
        id=9,
        upload_item_id="ordinary-transcript",
        lecture_id=7,
        kind=UploadKind.TRANSCRIPTS,
        source_sha256="a" * 64,
        immutable_source_path=Path("/tmp/ordinary-source.txt"),
        derived_sha256="b" * 64,
        immutable_derived_path=Path("/tmp/ordinary-cleaned.txt"),
        canonical_source_path=Path("/tmp/ordinary-source-current.txt"),
        canonical_derived_path=Path("/tmp/ordinary-cleaned-current.txt"),
        icloud_path=None,
        prompt_sha256="c" * 64,
        state="complete",
        current=True,
    )
    explicit_ordinary = replace(
        historical,
        provenance_kind="llm_cleaned",
        import_id=None,
    )
    expected_pre_v20_digest = "045b3bb6d8e02dc969571f4afe015a1b4885733252196408ef13ba6731df1142"
    assert revision_fingerprint(historical) == expected_pre_v20_digest
    assert revision_fingerprint(explicit_ordinary) == expected_pre_v20_digest
    assert revision_fingerprint(replace(historical, current=False)) == expected_pre_v20_digest
