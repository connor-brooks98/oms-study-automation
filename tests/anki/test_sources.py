import hashlib
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from oms_hub.anki.domain import SourceKind
from oms_hub.anki.sources import (
    LectureSourceExtractor,
    NotebookSummaryParser,
    SourcePassage,
    SummaryMalformedError,
)
from oms_hub.ingestion.domain import StudyRevision, UploadKind
from oms_hub.study_generation.domain import OutlineRecord
from oms_hub.study_generation.outline import OutlinePdfRenderer


class FakeRevisionRepository:
    def __init__(self, revisions: dict[int, StudyRevision]) -> None:
        self.revisions = revisions

    def get_study_revision(self, revision_id: int) -> StudyRevision:
        return self.revisions[revision_id]


class FakeOutlineRepository:
    def __init__(self, outlines: dict[int, OutlineRecord]) -> None:
        self.outlines = outlines

    def outline(self, outline_id: int) -> OutlineRecord | None:
        return self.outlines.get(outline_id)


def _revision(
    revision_id: int,
    kind: UploadKind,
    source: Path,
    *,
    derived: Path | None = None,
) -> StudyRevision:
    return StudyRevision(
        id=revision_id,
        upload_item_id=f"upload-{revision_id}",
        lecture_id=12,
        kind=kind,
        source_sha256=f"{revision_id:064x}",
        immutable_source_path=source,
        derived_sha256=None,
        immutable_derived_path=derived,
        canonical_source_path=None,
        canonical_derived_path=None,
        icloud_path=None,
        prompt_sha256=None,
        state="current",
        current=True,
    )


def _presentation(path: Path, image_path: Path) -> None:
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[1])
    first.shapes.title.text = "Iron deficiency anemia"
    first.placeholders[1].text = "Ferritin is low before serum iron."
    first.notes_slide.notes_text_frame.text = (
        "Remember that total iron-binding capacity rises."
    )
    second = deck.slides.add_slide(deck.slide_layouts[6])
    second.shapes.add_picture(
        str(image_path),
        Inches(1),
        Inches(1),
        width=Inches(2),
    )
    deck.save(path)


def test_extracts_stable_slide_notes_and_image_marker(
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "marker.png"
    Image.new("RGB", (8, 8), "white").save(image_path)
    slides_path = tmp_path / "lecture.pptx"
    _presentation(slides_path, image_path)
    repository = FakeRevisionRepository(
        {7: _revision(7, UploadKind.SLIDES, slides_path)}
    )
    extractor = LectureSourceExtractor(repository)

    first = extractor.extract([7])
    second = extractor.extract([7])

    assert [passage.passage_id for passage in first] == [
        passage.passage_id for passage in second
    ]
    assert [
        (passage.source_kind, passage.slide_number, passage.locator)
        for passage in first
    ] == [
        (SourceKind.SLIDE, 1, "slide:1"),
        (SourceKind.SPEAKER_NOTES, 1, "slide:1:notes"),
        (SourceKind.VISION, 2, "slide:2:image"),
    ]
    assert "Iron deficiency anemia" in first[0].text
    assert "total iron-binding capacity" in first[1].text
    assert first[2].text == ""
    assert first[2].extraction_status == "vision_unavailable"
    assert first[0].citation == "Lecture 12, slide 1"


def test_segments_timestamped_transcript_with_stable_overlap(
    tmp_path: Path,
) -> None:
    transcript = tmp_path / "cleaned.txt"
    transcript.write_text(
        "[00:00] Iron deficiency begins with depletion of iron stores. "
        "Ferritin falls early.\n"
        "[00:18] Total iron-binding capacity rises as deficiency progresses. "
        "Microcytosis appears later.\n"
        "[00:36] Treat the cause and replace iron when appropriate.\n",
        encoding="utf-8",
    )
    repository = FakeRevisionRepository(
        {
            8: _revision(
                8,
                UploadKind.TRANSCRIPTS,
                tmp_path / "raw.txt",
                derived=transcript,
            )
        }
    )
    extractor = LectureSourceExtractor(
        repository,
        transcript_max_chars=105,
        transcript_overlap_sentences=1,
    )

    passages = extractor.extract([8])

    assert len(passages) >= 2
    assert all(
        passage.source_kind is SourceKind.TRANSCRIPT
        for passage in passages
    )
    assert passages[0].start_seconds == 0
    assert passages[1].start_seconds is not None
    assert set(passages[0].text.split()) & set(passages[1].text.split())
    assert passages == extractor.extract([8])
    assert passages[0].citation.startswith("Lecture 12, transcript ")


def test_passage_rejects_blank_non_vision_evidence() -> None:
    try:
        SourcePassage.create(
            revision_id=9,
            lecture_id=12,
            artifact_id="upload-9",
            source_kind=SourceKind.SLIDE,
            locator="slide:1",
            text="",
            slide_number=1,
        )
    except ValueError as error:
        assert "blank" in str(error)
    else:
        raise AssertionError("blank source evidence was accepted")


def test_parses_notebook_outline_into_authority_labeled_passages(
    tmp_path: Path,
) -> None:
    payload = OutlinePdfRenderer().render(
        "Lecture 12 Outline",
        "# CORE CONCEPTS\n"
        "- Hereditary spherocytosis has increased MCHC [27, 28]\n\n"
        "# DEPTH MAP\n"
        "- DEEP: HS genes and EMA binding [31]\n\n"
        "# PROFESSOR EMPHASIS FLAGS\n"
        "- Repeated 3+ Times: distinguish HS from AIHA [42]",
    )
    path = tmp_path / "outline.pdf"
    path.write_bytes(payload)
    record = OutlineRecord(
        id=9,
        lecture_id=12,
        job_id="outline-job",
        path=path,
        sha256=hashlib.sha256(payload).hexdigest(),
        current=True,
    )

    passages = NotebookSummaryParser().parse(record)

    assert [passage.source_id for passage in passages] == [
        "SUM:12:CORE:01",
        "SUM:12:DEPTH:D1",
        "SUM:12:EMPH:E1",
    ]
    assert all(passage.source_kind is SourceKind.SUMMARY for passage in passages)
    assert passages[0].summary_backrefs == ("27", "28")
    assert passages[1].summary_backrefs == ("31",)
    assert passages[2].summary_backrefs == ("42",)
    assert passages[0].text.startswith("Hereditary spherocytosis")


def test_summary_requires_depth_and_professor_emphasis_sections() -> None:
    with pytest.raises(SummaryMalformedError, match="DEPTH MAP"):
        NotebookSummaryParser().parse_text(
            lecture_id=12,
            outline_id=9,
            text="# CORE CONCEPTS\n- Ferritin falls early",
        )


def test_summary_rejects_unreadable_pdf(tmp_path: Path) -> None:
    path = tmp_path / "outline.pdf"
    payload = b"not a PDF"
    path.write_bytes(payload)
    record = OutlineRecord(
        id=9,
        lecture_id=12,
        job_id="outline-job",
        path=path,
        sha256=hashlib.sha256(payload).hexdigest(),
        current=True,
    )

    with pytest.raises(SummaryMalformedError, match="readable PDF"):
        NotebookSummaryParser().parse(record)


def test_extractor_combines_pinned_revisions_and_notebook_summary(
    tmp_path: Path,
) -> None:
    transcript = tmp_path / "transcript.txt"
    transcript.write_text("[00:00] Ferritin falls early.", encoding="utf-8")
    summary = OutlinePdfRenderer().render(
        "Lecture 12 Outline",
        "# DEPTH MAP\n- DEEP: iron homeostasis [1]\n\n"
        "# PROFESSOR EMPHASIS FLAGS\n- Repeated: ferritin [2]",
    )
    summary_path = tmp_path / "outline.pdf"
    summary_path.write_bytes(summary)
    outline = OutlineRecord(
        id=9,
        lecture_id=12,
        job_id="outline-job",
        path=summary_path,
        sha256=hashlib.sha256(summary).hexdigest(),
        current=True,
    )
    extractor = LectureSourceExtractor(
        FakeRevisionRepository(
            {
                8: _revision(
                    8,
                    UploadKind.TRANSCRIPTS,
                    transcript,
                    derived=transcript,
                )
            }
        ),
        outlines=FakeOutlineRepository({9: outline}),
    )

    passages = extractor.extract([8], summary_outline_id=9)

    assert [passage.source_kind for passage in passages] == [
        SourceKind.TRANSCRIPT,
        SourceKind.SUMMARY,
        SourceKind.SUMMARY,
    ]
