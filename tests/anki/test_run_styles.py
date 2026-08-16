from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from oms_hub.anki.course_policy import CourseCurationPolicy, PolicyEmphasisColor
from oms_hub.anki.domain import SourceKind
from oms_hub.anki.fidelity_audit import R2FidelityDiagnostic, audit_fidelity
from oms_hub.anki.sources import (
    SourceEmphasisEvidence,
    SourcePassage,
    project_source_emphasis_evidence,
)
from oms_hub.document_processing.run_styles import (
    StyledTextRunSidecar,
    _theme_scheme,
    extract_styled_text_run_sidecar,
    normalized_text_sha256,
)


def _policy(
    *, fallback: str = "block", colors: tuple[PolicyEmphasisColor, ...] = ()
) -> CourseCurationPolicy:
    colors = colors or (PolicyEmphasisColor(rgb="FF0000", label="red"),)
    return CourseCurationPolicy(
        policy_id="course", revision=1, course_id="course", professor_label="Professor",
        scope_instruction="Red only.", emphasis_mode="colored_text", emphasis_colors=colors,
        missing_emphasis_fallback=fallback, tag_scope_mode="hard_filter",
        classification_strictness="strict", generation_style_profile="cloze",
        ordinary_cost_limit_microusd=1, hard_stop_cost_limit_microusd=1,
    )


def _sidecar(tmp_path: Path):
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    frame = slide.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text_frame
    frame.paragraphs[0].add_run().text = ""
    red = frame.paragraphs[0].add_run()
    red.text = "Red run"
    red.font.color.rgb = RGBColor(255, 0, 0)
    themed = frame.add_paragraph().add_run()
    themed.text = "Theme run"
    themed.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    themed.font.color.brightness = 0.2
    table = slide.shapes.add_table(1, 1, 0, Inches(2), Inches(2), Inches(1)).table
    table_frame = table.cell(0, 0).text_frame
    table_frame.paragraphs[0].add_run().text = "Table run one"
    table_frame.paragraphs[0].add_run().text = "Table run two"
    table_frame.add_paragraph().add_run().text = "Table run three"
    group = slide.shapes.add_group_shape()
    nested = group.shapes.add_group_shape()
    nested.shapes.add_textbox(0, 0, Inches(1), Inches(1)).text = "Nested run"
    slide.notes_slide.notes_text_frame.paragraphs[0].add_run().text = "Notes run"
    path = tmp_path / "styles.pptx"
    deck.save(path)
    return extract_styled_text_run_sidecar(path, source_id="slides")


def test_sidecar_is_deterministic_and_preserves_physical_run_locations(tmp_path: Path) -> None:
    first = _sidecar(tmp_path)
    second = _sidecar(tmp_path)
    assert first == second
    assert first.runs[0].locator == "slide:1:shape:1:p:1:r:1"
    assert first.runs[1].locator == "slide:1:shape:1:p:1:r:2"
    assert [
        run.locator for run in first.runs if "cell:1:1" in run.locator
    ] == [
        "slide:1:shape:2:cell:1:1:p:1:r:1",
        "slide:1:shape:2:cell:1:1:p:1:r:2",
        "slide:1:shape:2:cell:1:1:p:2:r:1",
    ]
    assert any(run.locator == "slide:1:shape:3.1.1:p:1:r:1" for run in first.runs)
    assert any(run.locator == "slide:1:notes:p:1:r:1" for run in first.runs)


def test_explicit_theme_brightness_and_policy_projection(tmp_path: Path) -> None:
    sidecar = _sidecar(tmp_path)
    red, themed = sidecar.runs[1:3]
    assert red.resolved_color == "FF0000"
    assert themed.theme_color == "accent1"
    assert themed.resolved_color == "729ACA"
    assert themed.brightness == 0.2
    evidence = project_source_emphasis_evidence(sidecar, _policy())
    assert [item.text for item in evidence] == ["Red run"]
    repeated = project_source_emphasis_evidence(sidecar, _policy())
    assert evidence[0].provenance_hash == repeated[0].provenance_hash


def test_rgb_delta_eight_is_inclusive_and_nine_is_not(tmp_path: Path) -> None:
    sidecar = _sidecar(tmp_path)
    at_eight = project_source_emphasis_evidence(
        sidecar, _policy(colors=(PolicyEmphasisColor(rgb="F70008", label="near"),))
    )
    at_nine = project_source_emphasis_evidence(
        sidecar, _policy(colors=(PolicyEmphasisColor(rgb="F60000", label="far"),))
    )
    assert [item.text for item in at_eight] == ["Red run"]
    assert at_nine == ()


def test_theme_clrmap_and_malformed_or_inherited_colors_fail_closed(tmp_path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    frame = slide.shapes.add_textbox(0, 0, Inches(2), Inches(1)).text_frame
    inherited = frame.paragraphs[0].add_run()
    inherited.text = "Inherited"
    themed = frame.add_paragraph().add_run()
    themed.text = "Mapped theme"
    themed.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
    malformed = frame.add_paragraph().add_run()
    malformed.text = "Malformed"
    malformed.font.color.rgb = RGBColor(255, 0, 0)
    malformed._r.rPr.solidFill.srgbClr.append(OxmlElement("a:tint"))
    system = frame.add_paragraph().add_run()
    system.text = "System"
    system.font.color.rgb = RGBColor(1, 2, 3)
    system_color = system._r.rPr.solidFill.srgbClr
    system_color.tag = qn("a:sysClr")
    system_color.attrib.clear()
    system_color.set("val", "windowText")
    system_color.set("lastClr", "112233")
    path = tmp_path / "unresolved.pptx"
    deck.save(path)

    sidecar = extract_styled_text_run_sidecar(path, source_id="slides")

    assert sidecar.runs[0].inherited_color_resolution is None
    assert (sidecar.runs[1].theme_color, sidecar.runs[1].resolved_color) == ("dk1", "000000")
    assert sidecar.runs[2].resolved_color is None
    assert sidecar.runs[2].inherited_color_resolution == "unresolved"
    assert sidecar.runs[3].resolved_color == "112233"
    evidence = project_source_emphasis_evidence(
        sidecar, _policy(colors=(PolicyEmphasisColor(theme_ref="dk1", label="dark"),))
    )
    assert [item.text for item in evidence] == ["Mapped theme"]
    audit = audit_fidelity(
        sidecar, _policy(colors=(PolicyEmphasisColor(theme_ref="dk1", label="dark"),))
    )
    assert (
        audit.matching_colored_count,
        audit.nonmatching_colored_count,
        audit.unresolved_color_count,
    ) == (1, 1, 1)


def test_r2_outcomes_are_hash_bound(tmp_path: Path) -> None:
    sidecar = _sidecar(tmp_path)
    continuing = audit_fidelity(sidecar, _policy())
    assert continuing.status == "continue"
    assert (
        continuing.matching_colored_count,
        continuing.nonmatching_colored_count,
        continuing.unresolved_color_count,
    ) == (1, 1, 0)
    no_match = _policy(colors=(PolicyEmphasisColor(rgb="00FF00", label="green"),))
    assert audit_fidelity(sidecar, no_match).status == "blocked"
    confirmation = _policy(
        fallback="require_confirmation", colors=no_match.emphasis_colors
    )
    assert audit_fidelity(sidecar, confirmation).status == "confirmation_required"
    transcript = SourcePassage.create(
        revision_id=1, lecture_id=1, artifact_id="upload", source_kind=SourceKind.TRANSCRIPT,
        locator="transcript:1:untimed", text="Fallback transcript",
    )
    degraded = audit_fidelity(
        sidecar, _policy(fallback="transcript_outline", colors=no_match.emphasis_colors),
        source_passages=(transcript,),
    )
    assert (degraded.status, degraded.may_advance, degraded.degraded_mode) == (
        "continue_degraded", True, "transcript_outline"
    )
    assert audit_fidelity(
        sidecar, _policy(fallback="transcript_outline", colors=no_match.emphasis_colors)
    ).status == "blocked_fallback_unavailable"
    assert degraded.diagnostic_sha256 == audit_fidelity(
        sidecar, _policy(fallback="transcript_outline", colors=no_match.emphasis_colors),
        source_passages=(transcript,),
    ).diagnostic_sha256


def test_hash_bound_contracts_reject_tampering_and_contradictory_outcomes(tmp_path: Path) -> None:
    sidecar = _sidecar(tmp_path)
    evidence = project_source_emphasis_evidence(sidecar, _policy())[0]
    with pytest.raises(ValueError, match="policy match"):
        SourceEmphasisEvidence.model_validate(
            {
                **evidence.model_dump(),
                "policy_match": False,
                "evidence_id": "",
                "provenance_hash": "",
            }
        )
    with pytest.raises(ValueError, match="evidence ID"):
        SourceEmphasisEvidence.model_validate({**evidence.model_dump(), "evidence_id": "0" * 64})
    with pytest.raises(ValueError, match="normalized text hash"):
        SourceEmphasisEvidence.model_validate(
            {
                **evidence.model_dump(),
                "normalized_text_sha256": "0" * 64,
                "evidence_id": "",
                "provenance_hash": "",
            }
        )
    with pytest.raises(ValueError, match="cannot be blank"):
        SourceEmphasisEvidence.model_validate(
            {
                **evidence.model_dump(),
                "text": " ",
                "normalized_text_sha256": normalized_text_sha256(" "),
                "evidence_id": "",
                "provenance_hash": "",
            }
        )
    sidecar_payload = sidecar.model_dump()
    sidecar_payload["runs"][0]["normalized_text_sha256"] = "0" * 64
    sidecar_payload["sidecar_sha256"] = ""
    with pytest.raises(ValueError, match="normalized text hash"):
        StyledTextRunSidecar.model_validate(sidecar_payload)
    with pytest.raises(ValueError, match="locators"):
        StyledTextRunSidecar.model_validate(
            {**sidecar.model_dump(), "runs": [*sidecar.runs, sidecar.runs[0]], "sidecar_sha256": ""}
        )
    continuing = audit_fidelity(sidecar, _policy())
    with pytest.raises(ValueError, match="must advance"):
        R2FidelityDiagnostic.model_validate(
            {**continuing.model_dump(), "may_advance": False, "diagnostic_sha256": ""}
        )
    with pytest.raises(ValueError, match="blocking"):
        R2FidelityDiagnostic.model_validate(
            {
                **continuing.model_dump(),
                "status": "blocked",
                "may_advance": False,
                "diagnostic_sha256": "",
            }
        )
    no_match = audit_fidelity(
        sidecar, _policy(colors=(PolicyEmphasisColor(rgb="00FF00", label="green"),))
    )
    with pytest.raises(ValueError, match="requires a matching"):
        R2FidelityDiagnostic.model_validate(
            {
                **no_match.model_dump(),
                "status": "continue",
                "may_advance": True,
                "diagnostic_sha256": "",
            }
        )
    with pytest.raises(ValueError, match="blocking"):
        R2FidelityDiagnostic.model_validate(
            {
                **no_match.model_dump(),
                "status": "confirmation_required",
                "may_advance": True,
                "diagnostic_sha256": "",
            }
        )
    with pytest.raises(ValueError, match="degraded"):
        R2FidelityDiagnostic.model_validate(
            {
                **continuing.model_dump(),
                "status": "continue_degraded",
                "may_advance": True,
                "degraded_mode": "transcript_outline",
                "diagnostic_sha256": "",
            }
        )
    with pytest.raises(ValueError, match="requires transcript or outline"):
        R2FidelityDiagnostic.model_validate(
            {
                **no_match.model_dump(),
                "status": "continue_degraded",
                "may_advance": True,
                "degraded_mode": "transcript_outline",
                "diagnostic_sha256": "",
            }
        )


def test_sidecar_rejects_source_mismatch_and_image_only_has_no_runs(tmp_path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    image_path = tmp_path / "image.png"
    Image.new("RGB", (1, 1), "white").save(image_path)
    slide.shapes.add_picture(str(image_path), 0, 0, Inches(1), Inches(1))
    path = tmp_path / "image-only.pptx"
    deck.save(path)

    assert extract_styled_text_run_sidecar(path, source_id="slides").runs == ()
    with pytest.raises(ValueError, match="SHA-256"):
        extract_styled_text_run_sidecar(path, source_id="slides", source_sha256="0" * 64)


def test_theme_scheme_skips_transformed_base_colors() -> None:
    theme = type(
        "Theme",
        (),
        {
            "blob": (
                b'<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                b"<a:themeElements><a:clrScheme name=\"test\">"
                b"<a:accent1><a:srgbClr val=\"FF0000\"><a:lumMod val=\"50000\"/>"
                b"</a:srgbClr></a:accent1>"
                b"<a:accent2><a:sysClr val=\"windowText\" lastClr=\"112233\"/>"
                b"</a:accent2></a:clrScheme></a:themeElements></a:theme>"
            )
        },
    )()
    relation = type("ThemeRelation", (), {"reltype": RT.THEME, "target_part": theme})()
    source = type(
        "SourcePart",
        (),
        {"partname": "/ppt/slides/slide1.xml", "rels": {"r": relation}},
    )()

    assert _theme_scheme(source) == {"accent2": "112233"}


def test_blank_colored_runs_do_not_satisfy_r2(tmp_path: Path) -> None:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    run = slide.shapes.add_textbox(0, 0, Inches(1), Inches(1)).text_frame.paragraphs[0].add_run()
    run.text = " "
    run.font.color.rgb = RGBColor(255, 0, 0)
    path = tmp_path / "blank-red.pptx"
    deck.save(path)

    sidecar = extract_styled_text_run_sidecar(path, source_id="slides")
    diagnostic = audit_fidelity(sidecar, _policy())

    assert (diagnostic.status, diagnostic.may_advance, diagnostic.matching_colored_count) == (
        "blocked", False, 0
    )
