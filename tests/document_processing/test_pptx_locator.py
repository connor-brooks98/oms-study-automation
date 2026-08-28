import sys
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from oms_hub.document_processing import presentation_render
from oms_hub.document_processing.anydoc_adapter import AnydocProcessor
from oms_hub.document_processing.domain import (
    DocumentLocator,
    ParsedAsset,
    ParsedDocument,
    ParsedSegment,
    SegmentKind,
)
from oms_hub.document_processing.ocr import LocalOcr
from oms_hub.document_processing.pptx_locator import PptxLocatorEnricher, walk_shapes_with_paths
from oms_hub.document_processing.presentation_render import PresentationRenderer
from oms_hub.files.office import OfficeUnavailableError
from tests.document_processing.pptx_factory import SlideFixture, build_pptx, snapshot_for


def test_enricher_uses_slide_order_for_text_notes_and_picture_locators(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(
            SlideFixture("Question 1", "Which structure?", note="Answer: A", image=True),
            SlideFixture("Question 2", "Which pathway?", note="Answer: B"),
        ),
    )

    parsed = AnydocProcessor(PptxLocatorEnricher()).parse(
        snapshot_for(source), tmp_path / "assets"
    )

    notes = tuple(segment for segment in parsed.segments if segment.kind is SegmentKind.NOTE)
    images = tuple(segment for segment in parsed.segments if segment.kind is SegmentKind.IMAGE)
    assert tuple(segment.locator.label for segment in notes) == ("slide 1 notes", "slide 2 notes")
    assert tuple(segment.text for segment in notes) == ("Answer: A", "Answer: B")
    assert tuple(segment.locator.label for segment in images) == ("slide 1 image 1",)
    assert images[0].asset_keys == (parsed.assets[0].key,)
    assert all(segment.locator.slide_number is not None for segment in parsed.segments)


def test_recursive_shape_walk_exposes_xml_order_group_paths(tmp_path: Path) -> None:
    source = build_pptx(tmp_path / "groups.pptx", slides=(SlideFixture("Title", "Body"),))
    presentation = Presentation(source)
    group = presentation.slides[0].shapes.add_group_shape()
    nested = group.shapes.add_group_shape()
    nested.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1)).text = "Nested"
    presentation.save(source)

    reopened = Presentation(source)
    located = tuple(walk_shapes_with_paths(reopened.slides[0].shapes))

    assert located[-1][1] == (3, 1, 1)
    assert located[-1][0].text == "Nested"


def test_enricher_retains_unmatched_anydoc_asset_without_auto_binding(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    snapshot = snapshot_for(source)
    parsed = ParsedDocument(
        source_id=snapshot.id,
        source_sha256=snapshot.sha256,
        source_format="pptx",
        parser_name="fixture",
        parser_version="1",
        segments=(),
        assets=(
            ParsedAsset(
                key="asset-0",
                path=None,
                media_type="image/png",
                sha256="a" * 64,
                locator=DocumentLocator("asset 0"),
                origin="ppt/media/not-in-slides.png",
            ),
        ),
        warnings=(),
    )

    enriched = PptxLocatorEnricher().enrich(snapshot, parsed)

    assert enriched.assets[0].locator.slide_number is None
    assert enriched.warnings == ("asset 'asset-0' could not be matched to a PowerPoint slide",)
    assert all(segment.kind is not SegmentKind.IMAGE for segment in enriched.segments)


def test_enricher_preserves_candidate_semantic_text_and_order_while_restoring_slides(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(
            SlideFixture("Question 1", "Original first slide"),
            SlideFixture("Question 2", "Original second slide"),
        ),
    )
    snapshot = snapshot_for(source)
    parsed = ParsedDocument(
        source_id=snapshot.id,
        source_sha256=snapshot.sha256,
        source_format="pptx",
        parser_name="anydoc",
        parser_version="1",
        segments=(
            ParsedSegment(
                "semantic-a",
                SegmentKind.PARAGRAPH,
                "Semantic first",
                DocumentLocator("block 1"),
            ),
            ParsedSegment(
                "semantic-b",
                SegmentKind.PARAGRAPH,
                "Semantic second",
                DocumentLocator("block 2"),
            ),
        ),
        assets=(),
        warnings=(),
    )

    enriched = PptxLocatorEnricher().enrich(snapshot, parsed)

    assert tuple(segment.key for segment in enriched.segments) == ("semantic-a", "semantic-b")
    assert tuple(segment.text for segment in enriched.segments) == (
        "Semantic first",
        "Semantic second",
    )
    assert tuple(segment.locator.slide_number for segment in enriched.segments) == (1, 2)


def test_enricher_restores_slide_style_when_anydoc_text_is_split_differently(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "split-style.pptx",
        slides=(SlideFixture("Answer", "C) Correct answer"),),
    )
    presentation = Presentation(source)
    run = presentation.slides[0].shapes[1].text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    presentation.save(source)
    snapshot = snapshot_for(source)
    parsed = ParsedDocument(
        snapshot.id,
        snapshot.sha256,
        "pptx",
        "fixture",
        "1",
        (
            ParsedSegment(
                "semantic-answer",
                SegmentKind.PARAGRAPH,
                "Semantically split answer block",
                DocumentLocator("block 1"),
            ),
        ),
        (),
        (),
    )

    enriched = PptxLocatorEnricher().enrich(snapshot, parsed)

    assert enriched.segments[0].locator.slide_number == 1
    assert enriched.segments[0].style_metadata == ("bold: C) Correct answer",)


def test_enricher_maps_split_duplicate_question_and_answer_slides_in_order(
    tmp_path: Path,
) -> None:
    body = "Question stem\nA) Wrong\nB) Correct"
    source = build_pptx(
        tmp_path / "split-duplicate.pptx",
        slides=(SlideFixture("Question", body), SlideFixture("Answer", body)),
    )
    presentation = Presentation(source)
    answer_run = presentation.slides[1].shapes[1].text_frame.paragraphs[2].runs[0]
    answer_run.font.bold = True
    presentation.save(source)
    snapshot = snapshot_for(source)
    parsed = ParsedDocument(
        snapshot.id,
        snapshot.sha256,
        "pptx",
        "fixture",
        "1",
        tuple(
            ParsedSegment(
                f"block-{index}",
                SegmentKind.PARAGRAPH,
                text,
                DocumentLocator(f"block {index}"),
            )
            for index, text in enumerate(
                ("Question stem", "A) Wrong\nB) Correct") * 2,
                start=1,
            )
        ),
        (),
        (),
    )

    enriched = PptxLocatorEnricher().enrich(snapshot, parsed)

    assert tuple(segment.locator.slide_number for segment in enriched.segments) == (
        1,
        1,
        2,
        2,
    )
    assert enriched.segments[3].style_metadata == ("bold: B) Correct",)


def test_enricher_adds_slide_style_marker_when_candidate_has_no_slide_locator(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "unlocated-style.pptx",
        slides=(
            SlideFixture("Question", "Which answer?"),
            SlideFixture("Answer", "B) Correct"),
        ),
    )
    presentation = Presentation(source)
    presentation.slides[1].shapes[1].text_frame.paragraphs[0].runs[0].font.bold = True
    presentation.save(source)
    snapshot = snapshot_for(source)
    parsed = ParsedDocument(
        snapshot.id,
        snapshot.sha256,
        "pptx",
        "fixture",
        "1",
        (
            ParsedSegment(
                "unlocated",
                SegmentKind.PARAGRAPH,
                "Candidate without matching location",
                DocumentLocator("block 1"),
            ),
        ),
        (),
        (),
    )

    enriched = PptxLocatorEnricher().enrich(snapshot, parsed)

    marker = next(
        segment for segment in enriched.segments
        if segment.key == "slide-2-style-metadata"
    )
    assert marker.text == ""
    assert marker.locator.slide_number == 2
    assert marker.style_metadata == ("bold: B) Correct",)


def test_enricher_binds_repeated_media_on_one_slide_to_every_picture_occurrence(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?", image=True),),
    )
    presentation = Presentation(source)
    presentation.slides[0].shapes.add_picture(
        str(source.with_suffix(".png")), Inches(3), Inches(3.5)
    )
    presentation.save(source)

    parsed = AnydocProcessor(PptxLocatorEnricher()).parse(
        snapshot_for(source), tmp_path / "assets"
    )

    images = tuple(segment for segment in parsed.segments if segment.kind is SegmentKind.IMAGE)
    assert parsed.warnings == ()
    assert parsed.assets[0].locator.slide_number == 1
    assert len(images) == 2
    assert all(segment.locator.slide_number == 1 for segment in images)
    assert all(segment.asset_keys == (parsed.assets[0].key,) for segment in images)


def test_enricher_leaves_media_reused_across_slides_unbound(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(
            SlideFixture("Question 1", "Which structure?", image=True),
            SlideFixture("Question 2", "Which pathway?", image=True),
        ),
    )

    parsed = AnydocProcessor(PptxLocatorEnricher()).parse(
        snapshot_for(source), tmp_path / "assets"
    )

    assert parsed.assets[0].locator.slide_number is None
    assert parsed.warnings == (
        "asset 'asset-0' occurs on multiple slides and was not automatically bound",
    )
    images = tuple(segment for segment in parsed.segments if segment.kind is SegmentKind.IMAGE)
    assert len(images) == 2
    assert tuple(segment.locator.slide_number for segment in images) == (1, 2)


def test_enricher_ocrs_image_only_slide(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "image-only.pptx",
        slides=(SlideFixture("", "", image=True),),
    )

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda _path: "Question screenshot text"))
    ).parse(snapshot_for(source), tmp_path / "assets")

    ocr = tuple(segment for segment in parsed.segments if segment.key.endswith("-ocr"))
    assert tuple(segment.text for segment in ocr) == ("Question screenshot text",)
    assert ocr[0].locator.slide_number == 1


def test_enricher_ocrs_large_picture_on_mixed_text_slide(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "mixed.pptx",
        slides=(SlideFixture("Question 1", "Tutor note", image=True),),
    )
    presentation = Presentation(source)
    picture = next(
        shape for shape in presentation.slides[0].shapes
        if shape.shape_type is MSO_SHAPE_TYPE.PICTURE
    )
    picture.width = Inches(9)
    picture.height = Inches(5)
    presentation.save(source)

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda _path: "Embedded Q-bank question"))
    ).parse(snapshot_for(source), tmp_path / "assets")

    assert "Embedded Q-bank question" in tuple(segment.text for segment in parsed.segments)


def test_enricher_skips_small_logo_on_text_slide(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "logo.pptx",
        slides=(SlideFixture("Question 1", "Visible question text", image=True),),
    )
    calls: list[Path] = []

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda path: calls.append(path) or "logo"))
    ).parse(snapshot_for(source), tmp_path / "assets")

    assert calls == []
    assert all(not segment.key.endswith("-ocr") for segment in parsed.segments)
    assert all(not warning.startswith("BLOCKER:") for warning in parsed.warnings)


def test_enricher_large_picture_selection_is_slide_specific_for_reused_media(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "reused.pptx",
        slides=(
            SlideFixture("Question 1", "Tutor note", image=True),
            SlideFixture("Question 2", "Tutor note", image=True),
        ),
    )
    presentation = Presentation(source)
    first = next(
        shape for shape in presentation.slides[0].shapes
        if shape.shape_type is MSO_SHAPE_TYPE.PICTURE
    )
    first.width = Inches(9)
    first.height = Inches(5)
    second = next(
        shape for shape in presentation.slides[1].shapes
        if shape.shape_type is MSO_SHAPE_TYPE.PICTURE
    )
    second.width = Inches(0.5)
    second.height = Inches(0.5)
    presentation.save(source)
    calls: list[Path] = []

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda path: calls.append(path) or "large screenshot"))
    ).parse(snapshot_for(source), tmp_path / "assets")

    assert len(calls) == 1
    assert tuple(
        segment.locator.slide_number
        for segment in parsed.segments
        if segment.key.endswith("-ocr")
    ) == (1,)


def test_enricher_reports_one_blocker_when_required_ocr_is_empty(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "unreadable.pptx",
        slides=(SlideFixture("", "", image=True),),
    )

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda _path: ""))
    ).parse(snapshot_for(source), tmp_path / "assets")

    assert tuple(warning for warning in parsed.warnings if warning.startswith("BLOCKER:")) == (
        "BLOCKER: OCR is required but unavailable or empty for slide 1",
    )


def test_enricher_does_not_block_native_question_when_diagnostic_image_has_no_text(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "radiograph-question.pptx",
        slides=(
            SlideFixture(
                "Question 1",
                "Which diagnosis best explains this patient's chronic painful swollen knee image?",
                image=True,
            ),
        ),
    )
    presentation = Presentation(source)
    picture = next(
        shape
        for shape in presentation.slides[0].shapes
        if shape.shape_type is MSO_SHAPE_TYPE.PICTURE
    )
    picture.width = Inches(9)
    picture.height = Inches(5)
    presentation.save(source)

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda _path: ""))
    ).parse(snapshot_for(source), tmp_path / "assets")

    assert all(not warning.startswith("BLOCKER:") for warning in parsed.warnings)
    assert all(not segment.key.endswith("-ocr") for segment in parsed.segments)


def test_enricher_does_not_count_speaker_notes_as_visible_screenshot_text(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "noted-screenshot.pptx",
        slides=(
            SlideFixture(
                "Question 1",
                "",
                note="Tutor notes explain this screenshot in more than eight hidden words.",
                image=True,
            ),
        ),
    )
    presentation = Presentation(source)
    picture = next(
        shape
        for shape in presentation.slides[0].shapes
        if shape.shape_type is MSO_SHAPE_TYPE.PICTURE
    )
    picture.width = Inches(9)
    picture.height = Inches(5)
    presentation.save(source)

    parsed = AnydocProcessor(
        PptxLocatorEnricher(LocalOcr(lambda _path: ""))
    ).parse(snapshot_for(source), tmp_path / "assets")

    assert tuple(warning for warning in parsed.warnings if warning.startswith("BLOCKER:")) == (
        "BLOCKER: OCR is required but unavailable or empty for slide 1",
    )


def test_enricher_preserves_answer_emphasis_as_sidecar_metadata(tmp_path: Path) -> None:
    source = build_pptx(
        tmp_path / "formatted-answer.pptx",
        slides=(SlideFixture("Answer", "C) Correct answer"),),
    )
    presentation = Presentation(source)
    run = presentation.slides[0].shapes[1].text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.italic = True
    run.font.underline = True
    run.font.color.rgb = RGBColor(255, 0, 0)
    highlight = OxmlElement("a:highlight")
    highlight.append(OxmlElement("a:srgbClr"))
    run._r.get_or_add_rPr().append(highlight)
    presentation.save(source)

    parsed = AnydocProcessor(PptxLocatorEnricher()).parse(
        snapshot_for(source), tmp_path / "assets"
    )

    formatted = next(segment for segment in parsed.segments if "Correct answer" in segment.text)
    assert formatted.style_metadata == (
        "bold: C) Correct answer",
        "italic: C) Correct answer",
        "underline: C) Correct answer",
        "highlighted: C) Correct answer",
        "color #FF0000: C) Correct answer",
    )


def test_enricher_restores_missing_style_cues_when_one_cue_already_matched(
    tmp_path: Path,
) -> None:
    source = build_pptx(
        tmp_path / "partial-style.pptx",
        slides=(SlideFixture("Styled title", "C) Correct answer"),),
    )
    presentation = Presentation(source)
    presentation.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].font.bold = True
    presentation.slides[0].shapes[1].text_frame.paragraphs[0].runs[0].font.color.rgb = (
        RGBColor(255, 0, 0)
    )
    presentation.save(source)
    snapshot = snapshot_for(source)
    parsed = ParsedDocument(
        snapshot.id,
        snapshot.sha256,
        "pptx",
        "fixture",
        "1",
        (
            ParsedSegment(
                "title-only",
                SegmentKind.HEADING,
                "Styled title",
                DocumentLocator("block 1"),
            ),
        ),
        (),
        (),
    )

    enriched = PptxLocatorEnricher().enrich(snapshot, parsed)

    assert enriched.segments[0].style_metadata == (
        "bold: Styled title",
        "color #FF0000: C) Correct answer",
    )


def test_renderer_persists_bounded_full_slide_candidates_with_slide_locators(
    tmp_path: Path, monkeypatch
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    monkeypatch.setitem(sys.modules, "fitz", _FitzFixture())

    rendered = PresentationRenderer(_PdfFixtureConverter()).render(
        snapshot_for(source), tmp_path / "assets"
    )

    assert rendered.warnings == ()
    assert tuple(asset.key for asset in rendered.assets) == ("slide-1-render", "slide-2-render")
    assert tuple(asset.locator.label for asset in rendered.assets) == (
        "slide 1 render",
        "slide 2 render",
    )
    assert all(asset.path is not None and asset.path.is_file() for asset in rendered.assets)


def test_renderer_returns_nonblocking_warning_for_corrupt_rasterizer(
    tmp_path: Path, monkeypatch
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    monkeypatch.setitem(sys.modules, "fitz", _BrokenFitzFixture())

    rendered = PresentationRenderer(_PdfFixtureConverter()).render(
        snapshot_for(source), tmp_path / "assets"
    )

    assert rendered.assets == ()
    assert rendered.warnings == ("slide renderer unavailable: corrupt PDF fixture",)


def test_renderer_returns_nonblocking_warning_for_sanitization_failure(
    tmp_path: Path, monkeypatch
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    monkeypatch.setitem(sys.modules, "fitz", _InvalidPngFitzFixture())

    rendered = PresentationRenderer(_PdfFixtureConverter()).render(
        snapshot_for(source), tmp_path / "assets"
    )

    assert rendered.assets == ()
    assert rendered.warnings[0].startswith("slide renderer unavailable: quiz image must be")


def test_renderer_preserves_unavailable_warning_when_temp_pdf_is_locked(
    tmp_path: Path, monkeypatch
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    original_unlink = Path.unlink

    def locked_unlink(path: Path, *, missing_ok: bool = False) -> None:
        if path.name == "slides.pdf":
            raise PermissionError("locked PDF")
        original_unlink(path, missing_ok=missing_ok)

    monkeypatch.setattr(Path, "unlink", locked_unlink)
    rendered = PresentationRenderer(_UnavailableConverter()).render(
        snapshot_for(source), tmp_path / "assets"
    )

    assert rendered.assets == ()
    assert rendered.warnings == ("slide renderer unavailable: Office is unavailable",)


def test_renderer_configures_temporary_directory_cleanup_while_preserving_unavailable_warning(
    tmp_path: Path, monkeypatch
) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    temporary_directories: list[tuple[str, bool]] = []
    original_temporary_directory = presentation_render.TemporaryDirectory

    def tracking_temporary_directory(
        *, prefix: str, ignore_cleanup_errors: bool
    ):
        temporary_directories.append((prefix, ignore_cleanup_errors))
        return original_temporary_directory(
            prefix=prefix, ignore_cleanup_errors=ignore_cleanup_errors
        )

    monkeypatch.setattr(
        presentation_render, "TemporaryDirectory", tracking_temporary_directory
    )
    rendered = PresentationRenderer(_UnavailableConverter()).render(
        snapshot_for(source), tmp_path / "assets"
    )

    assert temporary_directories == [("oms-slide-render-", True)]
    assert rendered.assets == ()
    assert rendered.warnings == ("slide renderer unavailable: Office is unavailable",)


def test_renderer_removes_temporary_pdf_after_success(tmp_path: Path, monkeypatch) -> None:
    source = build_pptx(
        tmp_path / "questions.pptx",
        slides=(SlideFixture("Question", "Which structure?"),),
    )
    deleted_paths: list[Path] = []
    original_unlink = Path.unlink

    def tracking_unlink(path: Path, *, missing_ok: bool = False) -> None:
        if path.name == "slides.pdf":
            deleted_paths.append(path)
        original_unlink(path, missing_ok=missing_ok)

    monkeypatch.setattr(Path, "unlink", tracking_unlink)
    monkeypatch.setitem(sys.modules, "fitz", _FitzFixture())

    rendered = PresentationRenderer(_PdfFixtureConverter()).render(
        snapshot_for(source), tmp_path / "assets"
    )

    assert rendered.warnings == ()
    assert len(deleted_paths) == 1
    assert not deleted_paths[0].exists()


class _PdfFixtureConverter:
    def convert(self, source: Path, destination: Path) -> None:
        destination.write_bytes(b"fixture pdf")


class _UnavailableConverter:
    def convert(self, source: Path, destination: Path) -> None:
        raise OfficeUnavailableError("Office is unavailable")


class _FitzFixture:
    def open(self, path: Path) -> "_FixtureDocument":
        return _FixtureDocument()


class _BrokenFitzFixture:
    def open(self, path: Path) -> "_FixtureDocument":
        raise RuntimeError("corrupt PDF fixture")


class _InvalidPngFitzFixture:
    def open(self, path: Path) -> "_InvalidPngFixtureDocument":
        return _InvalidPngFixtureDocument()


class _FixtureDocument:
    def __enter__(self) -> "_FixtureDocument":
        return self

    def __exit__(self, exc_type: object, exc_value: object, traceback: object) -> None:
        return None

    def __iter__(self):
        return iter((_FixturePage(), _FixturePage()))


class _FixturePage:
    def get_pixmap(self, *, alpha: bool) -> "_FixturePixmap":
        assert not alpha
        return _FixturePixmap()


class _FixturePixmap:
    def tobytes(self, format: str) -> bytes:
        assert format == "png"
        image = Image.new("RGB", (3, 2), (20, 40, 60))
        output = BytesIO()
        image.save(output, format="PNG")
        return output.getvalue()


class _InvalidPngFixtureDocument(_FixtureDocument):
    def __iter__(self):
        return iter((_InvalidPngFixturePage(),))


class _InvalidPngFixturePage:
    def get_pixmap(self, *, alpha: bool) -> "_InvalidPngFixturePixmap":
        assert not alpha
        return _InvalidPngFixturePixmap()


class _InvalidPngFixturePixmap:
    def tobytes(self, format: str) -> bytes:
        assert format == "png"
        return b"not an image"
