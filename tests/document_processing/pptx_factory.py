"""Small, deterministic PowerPoint fixtures for document-processing tests."""

import hashlib
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from oms_hub.document_processing.domain import SourceSnapshot


@dataclass(frozen=True, slots=True)
class SlideFixture:
    title: str
    body: str
    note: str = ""
    image: bool = False


def build_pptx(path: Path, *, slides: tuple[SlideFixture, ...]) -> Path:
    """Build a presentation whose slide text, notes, and image are explicit."""
    presentation = Presentation()
    image_path = path.with_suffix(".png")
    _write_image(image_path)
    for fixture in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(0.7)).text = (
            fixture.title
        )
        slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(1.5)).text = (
            fixture.body
        )
        if fixture.note:
            slide.notes_slide.notes_text_frame.text = fixture.note
        if fixture.image:
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(3.5))
    presentation.save(path)
    return path


def snapshot_for(path: Path) -> SourceSnapshot:
    payload = path.read_bytes()
    return SourceSnapshot(
        id="questions-pptx",
        title="Questions",
        path=path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        sha256=hashlib.sha256(payload).hexdigest(),
    )


def _write_image(path: Path) -> None:
    image = Image.new("RGB", (12, 8), (200, 20, 20))
    payload = BytesIO()
    image.save(payload, format="PNG")
    path.write_bytes(payload.getvalue())
