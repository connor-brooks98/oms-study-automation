from __future__ import annotations

import asyncio
import dataclasses
import hashlib
import importlib.metadata
import importlib.util
import inspect
import json
import os
import re
import stat
import sys
import tempfile
import time
from collections.abc import Iterator
from importlib import import_module
from pathlib import Path
from types import ModuleType, SimpleNamespace

import httpx
import pytest
from PIL import Image

# ruff: noqa: E501

ROOT = Path(__file__).resolve().parents[2]
SCRIPT = ROOT / "scripts" / "run-gemini-contract-smoke.py"
GATE_RECORD = (
    ROOT / "artifacts" / "acceptance" / "grounded-learning" / "gate-2b-gemini-indexing.json"
)
LIVE_ENABLED = os.getenv("RUN_LIVE_GEMINI_TESTS") == "1"


def _load_smoke() -> ModuleType:
    assert SCRIPT.is_file(), "Task 2.8 smoke scaffold is missing"
    spec = importlib.util.spec_from_file_location("task_2_8_gemini_smoke", SCRIPT)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)

    async def shared_public_matrix(
        session: object,
        *,
        clock: object = None,
        failure_evidence: dict[str, object] | None = None,
        diagnostic_sink: object = None,
    ) -> dict[str, object]:
        with tempfile.TemporaryDirectory() as directory:
            view = module._synthetic_index_input(Path(directory))
            return await module._run_shadow_sequence(
                session,
                view,
                module._private_shadow_preflight_from_view(view),
                mode="public_matrix",
                clock=clock or module.monotonic,
                failure_evidence=failure_evidence,
                diagnostic_sink=diagnostic_sink,
            )

    module.run_contract_smoke = shared_public_matrix
    return module


def _sdk_interactions() -> object:
    return import_module("google.genai.interactions")


class _FakeSession:
    def __init__(self, smoke: ModuleType, *, fail_import: bool = False) -> None:
        self.smoke = smoke
        self.fail_import = fail_import
        self.calls: list[tuple[str, object]] = []
        self.store_name = "fileSearchStores/raw-store-identity"
        self.file_name = "files/raw-file-identity"
        self.operation_name = "operations/raw-operation-identity"
        self.document_name = "fileSearchStores/raw-store-identity/documents/raw-document"
        self.model_contract = smoke.PRIVATE_SHADOW_MODEL_CONTRACT
        self.live_files: dict[str, str] = {}
        self.live_stores: dict[str, str] = {}

    async def create_store(self, display_name: str, embedding_model: str) -> str:
        self.calls.append(("create_store", (display_name, embedding_model)))
        self.live_stores[self.store_name] = display_name
        return self.store_name

    async def find_stores(self, display_name: str) -> tuple[str, ...]:
        return tuple(name for name, value in self.live_stores.items() if value == display_name)

    async def upload_input(self, display_name: str, path: Path, media_type: str) -> str:
        file_name = f"files/raw-file-identity-{len(self.live_files) + 1}"
        self.calls.append(
            ("upload_input", (display_name, path.name, media_type, path.stat().st_size))
        )
        self.live_files[file_name] = display_name
        return file_name

    async def import_input(
        self, store_name: str, file_name: str, metadata: tuple[tuple[str, str], ...], chunking: object | None
    ) -> str:
        self.calls.append(("import_input", (store_name, file_name, metadata, chunking)))
        if self.fail_import:
            self.fail_import = False
            raise self.smoke.SmokeTemporaryFailure("synthetic temporary failure")
        return f"operations/raw-operation-identity-{len(self.calls)}"

    async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]:
        return tuple(name for name, value in self.live_files.items() if value in display_names)

    async def query_private(self, store_name: str, prompt: str, scope: object, **kwargs: object) -> object:
        structured_no_result = kwargs.get("require_structured_no_result", False)
        structured_supported = kwargs.get("require_structured_supported", False)
        result = await self.query(
            store_name,
            prompt,
            scope,
            response_schema=(
                self.smoke.SmokeAnswer
                if structured_no_result or structured_supported
                else None
            ),
            omit_thinking=True,
        )
        try:
            answer = self.smoke.SmokeAnswer.model_validate(result.answer)
        except Exception:
            raise self.smoke.SmokeContractError(
                "structured output was invalid", reason="structured_output_invalid"
            ) from None
        if structured_no_result or structured_supported:
            citation = result.citations[0] if result.citations else None
            return self.smoke.PrivateShadowQueryAudit(
                len(result.citations),
                len(result.citations),
                result.input_tokens or 0,
                result.output_tokens or 0,
                answer.supported,
                answer.answer == "",
                answer.answer,
                citation.page_number if citation is not None else None,
                citation.excerpt if citation is not None else None,
            )
        if not answer.supported:
            raise self.smoke.SmokeContractError("positive answer was invalid")
        if self.smoke.SYNTHETIC_MARKER not in answer.answer:
            raise self.smoke.SmokeContractError(
                "positive marker was missing", reason="positive_answer_missing_marker"
            )
        if not result.citations:
            raise self.smoke.SmokeContractError(
                "positive citation was missing", reason="positive_citation_missing"
            )
        return self.smoke.PrivateShadowQueryAudit(
            len(result.citations), len(result.citations),
            result.input_tokens or 0, result.output_tokens or 0, True, None,
            answer.answer, result.citations[0].page_number, result.citations[0].excerpt,
        )

    async def wait_for_import(self, operation_name: str) -> str:
        self.calls.append(("wait_for_import", operation_name))
        return f"fileSearchStores/raw-store-identity/documents/{operation_name.rsplit('/', 1)[1]}"

    async def query(
        self,
        store_name: str,
        prompt: str,
        scope: object,
        *,
        response_schema: type[object],
        omit_thinking: bool,
    ) -> object:
        del prompt, response_schema
        self.calls.append(("query", (store_name, scope, omit_thinking)))
        if scope.lecture_id == self.smoke.SYNTHETIC_LECTURE_ID:  # type: ignore[attr-defined]
            return self.smoke.SmokeQueryResult(
                answer={"answer": self.smoke.SYNTHETIC_FACT, "supported": True},
                citations=(
                    self.smoke.SmokeCitation(
                        document_name=self.document_name,
                        page_number=1,
                        excerpt=self.smoke.SYNTHETIC_FACT,
                    ),
                ),
                input_tokens=11,
                output_tokens=7,
            )
        return self.smoke.SmokeQueryResult(
            answer={"answer": "", "supported": False},
            citations=(),
        )

    async def list_documents(self, store_name: str) -> tuple[str, ...]:
        self.calls.append(("list_documents", store_name))
        return (self.document_name,)

    async def delete_document(self, document_name: str) -> None:
        self.calls.append(("delete_document", document_name))

    async def delete_file(self, file_name: str) -> None:
        self.calls.append(("delete_file", file_name))
        self.live_files.pop(file_name, None)

    async def delete_store(self, store_name: str) -> None:
        self.calls.append(("delete_store", store_name))
        self.live_stores.pop(store_name, None)


class _FiveInputSession:
    def __init__(self, smoke: ModuleType, failed_input: str) -> None:
        self.smoke = smoke
        self.failed_input = failed_input
        self.calls: list[tuple[str, object]] = []
        self.model_contract = smoke.PRIVATE_SHADOW_MODEL_CONTRACT
        self.live_files: dict[str, str] = {}
        self.live_stores: dict[str, str] = {}
        self.file_media_types: dict[str, str] = {}
        self._next_file = 0

    async def create_store(self, display_name: str, embedding_model: str) -> str:
        self.calls.append(("create_store", (display_name, embedding_model)))
        self.live_stores["fileSearchStores/synthetic"] = display_name
        return "fileSearchStores/synthetic"

    async def find_stores(self, display_name: str) -> tuple[str, ...]:
        self.calls.append(("find_stores", display_name))
        return tuple(name for name, value in self.live_stores.items() if value == display_name)

    async def upload_input(self, display_name: str, path: Path, media_type: str) -> str:
        self._next_file += 1
        file_name = f"files/synthetic-{self._next_file}"
        self.calls.append(
            ("upload_input", (display_name, path.name, media_type, path.stat().st_size))
        )
        self.live_files[file_name] = display_name
        self.file_media_types[file_name] = media_type
        return file_name

    async def import_input(
        self,
        store_name: str,
        file_name: str,
        metadata: tuple[tuple[str, str], ...],
        chunking: object | None,
    ) -> str:
        self.calls.append(("import_input", (store_name, file_name, metadata, chunking)))
        media_type = self.file_media_types[file_name]
        if self.failed_input == media_type or (
            self.failed_input == "normalized_markdown"
            and ("input_key", "normalized_markdown") in metadata
        ):
            raise self.smoke.SmokeContractError("synthetic import failure")
        return f"operations/{file_name.rsplit('-', 1)[1]}"

    async def wait_for_import(self, operation_name: str) -> str:
        self.calls.append(("wait_for_import", operation_name))
        return f"fileSearchStores/synthetic/documents/{operation_name.rsplit('/', 1)[1]}"

    async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]:
        self.calls.append(("find_files", display_names))
        return tuple(name for name, value in self.live_files.items() if value in display_names)

    async def query_private(self, *args: object, **kwargs: object) -> object:
        del args, kwargs
        pytest.fail("queries must not run after an input import failure")

    async def delete_document(self, document_name: str) -> None:
        self.calls.append(("delete_document", document_name))

    async def delete_file(self, file_name: str) -> None:
        self.calls.append(("delete_file", file_name))
        self.live_files.pop(file_name, None)

    async def delete_store(self, store_name: str) -> None:
        self.calls.append(("delete_store", store_name))
        self.live_stores.pop(store_name, None)


class _OneTimeResponseLossSession(_FiveInputSession):
    def __init__(
        self,
        smoke: ModuleType,
        failure_stage: str,
        error: BaseException | None = None,
    ) -> None:
        super().__init__(smoke, "")
        self.failure_stage = failure_stage
        self.failed = False
        self.error = error or smoke.GeminiProviderError("synthetic response loss")

    def _fail_once(self) -> None:
        if not self.failed:
            self.failed = True
            raise self.error

    async def upload_input(self, *args: object, **kwargs: object) -> str:
        value = await super().upload_input(*args, **kwargs)
        if self.failure_stage == "upload_input":
            self._fail_once()
        return value

    async def import_input(self, *args: object, **kwargs: object) -> str:
        value = await super().import_input(*args, **kwargs)
        if self.failure_stage == "import_input":
            self._fail_once()
        return value

    async def wait_for_import(self, *args: object, **kwargs: object) -> str:
        value = await super().wait_for_import(*args, **kwargs)
        if self.failure_stage == "wait_for_import":
            self._fail_once()
        return value


class _SdkDocuments:
    def __init__(self) -> None:
        self.calls: list[tuple[str, object]] = []

    async def list(self, *, parent: str) -> tuple[object, ...]:
        self.calls.append(("list", parent))
        return (SimpleNamespace(name="fileSearchStores/sdk-store/documents/sdk-document"),)

    async def delete(self, *, name: str, config: object) -> None:
        self.calls.append(("delete", (name, config)))


class _SdkStores:
    def __init__(self) -> None:
        self.calls: list[tuple[str, object]] = []
        self.documents = _SdkDocuments()
        self.listed_stores: tuple[object, ...] = ()

    async def create(self, *, config: object) -> object:
        self.calls.append(("create", config))
        return SimpleNamespace(name="fileSearchStores/sdk-store")

    async def import_file(
        self,
        *,
        file_search_store_name: str,
        file_name: str,
        config: object,
    ) -> object:
        self.calls.append(("import_file", (file_search_store_name, file_name, config)))
        return SimpleNamespace(name="operations/sdk-operation")

    async def delete(self, *, name: str, config: object) -> None:
        self.calls.append(("delete", (name, config)))

    async def list(self, *, config: object) -> object:
        self.calls.append(("list", config))

        async def entries() -> object:
            for item in self.listed_stores:
                yield item

        return entries()


class _SdkFiles:
    def __init__(self) -> None:
        self.calls: list[tuple[str, object]] = []
        self.listed_files: tuple[object, ...] = ()

    async def upload(self, *, file: object, config: object) -> object:
        content = file.read_bytes() if isinstance(file, Path) else file.read()
        self.calls.append(("upload", (content, config)))
        assert isinstance(config, dict)
        return SimpleNamespace(
            name=(
                "files/sdk-file"
                if config["mime_type"] == "application/pdf"
                else f"files/sdk-file-{config['mime_type']}"
            )
        )

    async def delete(self, *, name: str) -> None:
        self.calls.append(("delete", name))

    async def list(self, *, config: object) -> object:
        self.calls.append(("list", config))

        async def entries() -> object:
            for item in self.listed_files:
                yield item

        return entries()


class _SdkOperations:
    def __init__(self, *, error_status: int | None = None, delay: float = 0.0) -> None:
        self.calls: list[object] = []
        self.error_status = error_status
        self.delay = delay

    async def get(self, operation: object) -> object:
        self.calls.append(operation)
        if self.delay:
            await asyncio.sleep(self.delay)
        return SimpleNamespace(
            name="operations/sdk-operation",
            done=True,
            error={"code": self.error_status} if self.error_status is not None else None,
            response=SimpleNamespace(
                parent="sdk-store",
                document_name="sdk-document",
            ),
        )


class _SdkInteractions:
    def __init__(self, smoke: ModuleType) -> None:
        self.smoke = smoke
        self.calls: list[dict[str, object]] = []

    async def create(self, **body: object) -> object:
        sdk = _sdk_interactions()
        self.calls.append(body)
        tools = body["tools"]
        assert isinstance(tools, list)
        file_search = tools[0]
        if self.smoke.WRONG_LECTURE_ID in file_search["metadata_filter"]:
            annotations: list[object] = []
            answer = self.smoke.SmokeAnswer(answer="", supported=False)
        else:
            metadata = {
                key: value
                for key, value in (
                    ("course_id", self.smoke.SYNTHETIC_COURSE_ID),
                    ("exam_id", self.smoke.SYNTHETIC_EXAM_ID),
                    ("lecture_id", self.smoke.SYNTHETIC_LECTURE_ID),
                    ("source_revision_id", self.smoke.SYNTHETIC_REVISION_ID),
                    ("authority_class", "course_material"),
                    ("input_key", "pdf"),
                    ("input_kind", "pdf"),
                    (
                        "input_sha256",
                        hashlib.sha256(self.smoke.synthetic_pdf_bytes()).hexdigest(),
                    ),
                )
            }
            annotations = [
                sdk.FileCitation(
                    custom_metadata=metadata,
                    document_uri="fileSearchStores/sdk-store",
                    file_name="sdk-file",
                    page_number=1,
                    source=f"\n{self.smoke.SYNTHETIC_FACT}\n",
                )
            ]
            answer = self.smoke.SmokeAnswer(
                answer=self.smoke.SYNTHETIC_FACT,
                supported=True,
            )
        output_text = answer.model_dump_json()
        return sdk.Interaction(
            status="completed",
            output_text=output_text,
            steps=[
                sdk.FileSearchResultStep(call_id="sdk-file-search-call"),
                sdk.ModelOutputStep(
                    content=[
                        sdk.TextContent(
                            text=output_text,
                            annotations=annotations,
                        )
                    ],
                )
            ],
            usage=sdk.Usage(total_input_tokens=13, total_output_tokens=8),
        )


def _transient_500() -> httpx.HTTPStatusError:
    return httpx.HTTPStatusError(
        "synthetic", request=httpx.Request("POST", "https://example.invalid"), response=httpx.Response(500)
    )


class _PrivateSdkInteractions:
    def __init__(self, smoke: ModuleType, view: object, manifest: object) -> None:
        self.smoke = smoke
        self.view = view
        self.manifest = manifest
        self.calls: list[dict[str, object]] = []

    async def create(self, **body: object) -> object:
        sdk = _sdk_interactions()
        self.calls.append(body)
        tools = body["tools"]
        assert isinstance(tools, list)
        file_search = tools[0]
        if self.smoke.PRIVATE_SHADOW_WRONG_LECTURE_ID in file_search["metadata_filter"]:
            annotations: list[object] = []
            output_text = self.smoke.SmokeAnswer(
                answer="",
                supported=False,
            ).model_dump_json()
        else:
            pdf = next(item for item in self.manifest.inputs if item.input_key == "pdf")
            annotations = [
                sdk.FileCitation(
                    custom_metadata={
                        "course_id": self.view.course_id,
                        "exam_id": self.view.exam_id,
                        "lecture_id": self.view.lecture_id,
                        "source_revision_id": self.view.source_revision_id,
                        "input_key": pdf.input_key,
                        "input_kind": pdf.input_kind,
                        "input_sha256": pdf.sha256,
                    },
                    document_uri="fileSearchStores/private-store",
                    file_name="sdk-file",
                    page_number=1,
                    source=self.view.evidence_units[0].normalized_text,
                )
            ]
            output_text = "private response discarded by adapter"
        return sdk.Interaction(
            status="completed",
            output_text=output_text,
            steps=[
                sdk.FileSearchResultStep(call_id="private-search"),
                sdk.ModelOutputStep(
                    content=[sdk.TextContent(text=output_text, annotations=annotations)]
                ),
            ],
            usage=sdk.Usage(total_input_tokens=19, total_output_tokens=6),
        )


class _SdkAio:
    def __init__(
        self,
        smoke: ModuleType,
        *,
        operation_error_status: int | None = None,
        operation_delay: float = 0.0,
    ) -> None:
        self.file_search_stores = _SdkStores()
        self.files = _SdkFiles()
        self.operations = _SdkOperations(
            error_status=operation_error_status,
            delay=operation_delay,
        )
        self.interactions = _SdkInteractions(smoke)
        self.closed = 0

    async def aclose(self) -> None:
        self.closed += 1


class _SdkClient:
    def __init__(
        self,
        smoke: ModuleType,
        *,
        operation_error_status: int | None = None,
        operation_delay: float = 0.0,
    ) -> None:
        self.aio = _SdkAio(
            smoke,
            operation_error_status=operation_error_status,
            operation_delay=operation_delay,
        )


class _FakeSecrets:
    def __init__(self, value: str | None) -> None:
        self.value = value
        self.calls: list[str] = []

    def get(self, key: str) -> str | None:
        self.calls.append(key)
        return self.value


def _real_citation(smoke: ModuleType, **overrides: object) -> object:
    sdk = _sdk_interactions()
    values: dict[str, object] = {
        "custom_metadata": {
            "course_id": smoke.SYNTHETIC_COURSE_ID,
            "exam_id": smoke.SYNTHETIC_EXAM_ID,
            "lecture_id": smoke.SYNTHETIC_LECTURE_ID,
            "source_revision_id": smoke.SYNTHETIC_REVISION_ID,
        },
        "document_uri": "fileSearchStores/sdk-store",
        "file_name": "task-2-8-synthetic.pdf",
        "page_number": 1,
        "source": smoke.SYNTHETIC_FACT,
    }
    values.update(overrides)
    return sdk.FileCitation(**values)


def _real_interaction(
    smoke: ModuleType,
    *,
    citation: object | None = None,
    steps: list[object] | None = None,
    usage: object | None = None,
) -> object:
    sdk = _sdk_interactions()
    if steps is None:
        steps = [
            sdk.FileSearchResultStep(call_id="sdk-file-search-call"),
            sdk.ModelOutputStep(
                content=[
                    sdk.TextContent(
                        text=smoke.SYNTHETIC_FACT,
                        annotations=[citation or _real_citation(smoke)],
                    )
                ]
            ),
        ]
    return sdk.Interaction(
        status="completed",
        output_text=smoke.SYNTHETIC_FACT,
        steps=steps,
        usage=usage or sdk.Usage(
            total_input_tokens=13,
            total_output_tokens=8,
        ),
    )


def test_google_genai_2_14_session_maps_exact_sdk_contract() -> None:
    smoke = _load_smoke()
    clients: list[_SdkClient] = []

    def sdk_factory(**kwargs: object) -> _SdkClient:
        assert kwargs == {
            "api_key": "synthetic-sdk-key",
            "http_options": {"api_version": "v1beta", "timeout": 120_000},
        }
        client = _SdkClient(smoke)
        clients.append(client)
        return client

    session = smoke.GoogleGenaiSmokeSession("synthetic-sdk-key", sdk_factory=sdk_factory)
    record = asyncio.run(smoke.run_contract_smoke(session, clock=lambda: 100.0))

    assert record["status"] == "passed"
    assert [item["input_kind"] for item in record["input_results"]] == [
        "pptx",
        "pdf",
        "markdown",
        "image",
        "image",
    ]
    assert record["aggregate"]["input_count"] == 5
    assert record["aggregate"]["indexed_bytes"] > len(smoke.synthetic_pdf_bytes())
    all_aio = [client.aio for client in clients]
    assert all(aio.closed == 1 for aio in all_aio)
    store_calls = [call for aio in all_aio for call in aio.file_search_stores.calls]
    assert ("list", {"page_size": 20}) in store_calls
    assert any(call[0] == "create" for call in store_calls)
    upload_calls = [call for aio in all_aio for call in aio.files.calls if call[0] == "upload"]
    assert [call[1][1]["mime_type"] for call in upload_calls] == [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/pdf",
        "text/markdown",
        "image/png",
        "image/jpeg",
    ]
    assert all(
        call[1][1]["display_name"].startswith("task-2-8-public_matrix-")
        for call in upload_calls
    )
    import_call = next(call for call in store_calls if call[0] == "import_file")
    import_config = import_call[1][2]
    assert import_config.http_options.extra_body["customMetadata"][0] == {
        "key": "authority_class",
        "stringValue": "course_material",
    }
    query_bodies = [
        aio.interactions.calls[0] for aio in all_aio if aio.interactions.calls
    ]
    assert len(query_bodies) == 2
    assert all("generation_config" not in body for body in query_bodies)
    assert query_bodies[0]["model"] == "gemini-3.7-flash"
    assert query_bodies[0]["store"] is False
    assert set(query_bodies[0]) == {
        "input",
        "model",
        "response_format",
        "store",
        "tools",
    }
    assert query_bodies[0]["response_format"] == {
        "type": "text",
        "mime_type": "application/json",
        "schema": smoke.SmokeAnswer.model_json_schema(),
    }
    assert query_bodies[1]["response_format"] == {
        "type": "text",
        "mime_type": "application/json",
        "schema": smoke.SmokeAnswer.model_json_schema(),
    }
    assert query_bodies[0]["tools"] == [
        {
            "type": "file_search",
            "file_search_store_names": ["fileSearchStores/sdk-store"],
            "metadata_filter": (
                'course_id="task-2-8-synthetic-course" AND '
                'exam_id="task-2-8-synthetic-exam" AND '
                'lecture_id="task-2-8-synthetic-lecture"'
            ),
        }
    ]
    assert any(aio.file_search_stores.documents.calls for aio in all_aio)
    assert any(call[0] == "delete" for call in store_calls)


def test_real_sdk_smoke_store_list_respects_provider_page_limit() -> None:
    smoke = _load_smoke()
    sdk = import_module("google.genai")
    requests: list[tuple[str, str, dict[str, str], bytes]] = []

    async def handler(request: httpx.Request) -> httpx.Response:
        requests.append(
            (
                request.method,
                request.url.path,
                dict(request.url.params),
                request.content,
            )
        )
        if request.url.params.get("pageToken") is None:
            return httpx.Response(
                200,
                json={
                    "fileSearchStores": [
                        {
                            "name": "fileSearchStores/one",
                            "displayName": "target-store",
                        }
                    ],
                    "nextPageToken": "synthetic-next-page",
                },
                request=request,
            )
        return httpx.Response(
            200,
            json={
                "fileSearchStores": [
                    {"name": "fileSearchStores/two", "displayName": "other-store"}
                ]
            },
            request=request,
        )

    async_client = httpx.AsyncClient(transport=httpx.MockTransport(handler))
    sdk_client = sdk.Client(
        api_key="synthetic-sdk-key",
        http_options={
            "api_version": "v1beta",
            "base_url": "https://unit.invalid",
            "httpx_async_client": async_client,
        },
    )
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key",
        sdk_factory=lambda **kwargs: sdk_client,
    )

    matched = asyncio.run(session.find_stores("target-store"))

    assert matched == ("fileSearchStores/one",)
    assert requests == [
        (
            "GET",
            "/v1beta/fileSearchStores",
            {"pageSize": "20"},
            b"",
        ),
        (
            "GET",
            "/v1beta/fileSearchStores",
            {"pageSize": "20", "pageToken": "synthetic-next-page"},
            b"",
        ),
    ]


@pytest.mark.parametrize(
    ("filename", "media_type", "input_key", "chunking"),
    (
        (
            "lecture.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
            None,
        ),
        ("lecture.pdf", "application/pdf", "pdf", None),
        (
            "normalized.md",
            "text/markdown",
            "normalized_markdown",
            {
                "white_space_config": {
                    "max_tokens_per_chunk": 700,
                    "max_overlap_tokens": 100,
                }
            },
        ),
        ("visual.png", "image/png", "image_png", None),
        ("visual.jpg", "image/jpeg", "image_jpeg", None),
    ),
)
def test_real_sdk_public_matrix_upload_and_import_wire_contract(
    tmp_path: Path,
    filename: str,
    media_type: str,
    input_key: str,
    chunking: object | None,
) -> None:
    smoke = _load_smoke()
    sdk = import_module("google.genai")
    requests: list[httpx.Request] = []

    async def handler(request: httpx.Request) -> httpx.Response:
        requests.append(request)
        if request.url.path == "/upload/v1beta/files":
            return httpx.Response(
                200,
                headers={"X-Goog-Upload-URL": "https://unit.invalid/upload-session"},
                request=request,
            )
        if request.url.path == "/upload-session":
            return httpx.Response(
                200,
                json={
                    "file": {
                        "name": f"files/{input_key}",
                        "mimeType": media_type,
                    }
                },
                headers={"X-Goog-Upload-Status": "final"},
                request=request,
            )
        return httpx.Response(
            200,
            json={"name": "operations/import-1", "done": False},
            request=request,
        )

    def sdk_factory(**kwargs: object) -> object:
        assert kwargs["api_key"] == "synthetic-sdk-key"
        return sdk.Client(
            api_key="synthetic-sdk-key",
            http_options={
                "api_version": "v1beta",
                "base_url": "https://unit.invalid",
                "httpx_async_client": httpx.AsyncClient(
                    transport=httpx.MockTransport(handler)
                ),
            },
        )

    fixture = tmp_path / filename
    fixture.write_bytes(b"synthetic-public-matrix")
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key",
        sdk_factory=sdk_factory,
    )

    file_name = asyncio.run(session.upload_input(filename, fixture, media_type))
    operation = asyncio.run(
        session.import_input(
            "fileSearchStores/synthetic",
            file_name,
            (("input_key", input_key),),
            chunking,
        )
    )

    assert operation == "operations/import-1"
    upload_init = requests[0]
    upload_body = json.loads(upload_init.content)
    assert upload_init.headers["x-goog-upload-header-content-type"] == media_type
    assert upload_body["file"]["mime_type"] == media_type
    expected_import = {
        "fileName": f"files/{input_key}",
        "customMetadata": [
            {"key": "input_key", "stringValue": input_key}
        ],
    }
    if chunking is not None:
        expected_import["chunkingConfig"] = {
            "whiteSpaceConfig": {
                "maxTokensPerChunk": 700,
                "maxOverlapTokens": 100,
            }
        }
    assert json.loads(requests[-1].content) == expected_import


def test_real_sdk_pdf_import_wire_contract() -> None:
    smoke = _load_smoke()
    sdk = import_module("google.genai")
    captured: list[dict[str, object]] = []

    async def handler(request: httpx.Request) -> httpx.Response:
        captured.append(json.loads(request.content))
        return httpx.Response(
            200,
            json={"name": "operations/import-1", "done": False},
            request=request,
        )

    def sdk_factory(**kwargs: object) -> object:
        assert kwargs["api_key"] == "synthetic-sdk-key"
        return sdk.Client(
            api_key="synthetic-sdk-key",
            http_options={
                "api_version": "v1beta",
                "base_url": "https://unit.invalid",
                "httpx_async_client": httpx.AsyncClient(
                    transport=httpx.MockTransport(handler)
                ),
            },
        )

    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key",
        sdk_factory=sdk_factory,
    )

    assert not hasattr(session, "upload_pdf")
    assert not hasattr(session, "import_file")
    operation = asyncio.run(
        session.import_input(
            "fileSearchStores/synthetic",
            "files/lecture-pdf",
            (("input_key", "lecture_pdf"),),
            None,
        )
    )

    assert operation == "operations/import-1"
    assert captured == [
        {
            "fileName": "files/lecture-pdf",
            "customMetadata": [
                {"key": "input_key", "stringValue": "lecture_pdf"}
            ],
        }
    ]


def test_public_synthetic_index_input_contains_the_ordered_five_media_matrix(
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = smoke._synthetic_index_input(tmp_path)
    manifest = smoke._private_shadow_manifest(view)

    assert [(item.input_kind, item.media_type) for item in manifest.inputs] == [
        (
            "pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        ("pdf", "application/pdf"),
        ("markdown", "text/markdown"),
        ("image", "image/png"),
        ("image", "image/jpeg"),
    ]
    from pptx import Presentation

    presentation = Presentation(view.pptx.path)
    assert len(presentation.slides) == 1
    assert presentation.slides[0].shapes.title.text == smoke.SYNTHETIC_FACT
    assert smoke.SYNTHETIC_FACT.encode() in view.pdf.path.read_bytes()
    assert smoke.SYNTHETIC_FACT in view.markdown.path.read_text(encoding="utf-8")
    for asset in view.assets:
        assert asset.width is not None and asset.height is not None
        assert asset.width < 4096 and asset.height < 4096
        with Image.open(asset.path) as image:
            assert image.size == (asset.width, asset.height)
            image.verify()


def test_public_synthetic_fixtures_are_byte_identical_across_clock_seconds(
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    first_root = tmp_path / "first"
    second_root = tmp_path / "second"
    first_root.mkdir()
    second_root.mkdir()
    first = smoke._synthetic_index_input(first_root)
    time.sleep(2)
    second = smoke._synthetic_index_input(second_root)

    assert (
        first.pptx.sha256,
        first.pdf.sha256,
        first.markdown.sha256,
        *(asset.sha256 for asset in first.assets),
    ) == (
        second.pptx.sha256,
        second.pdf.sha256,
        second.markdown.sha256,
        *(asset.sha256 for asset in second.assets),
    )


@pytest.mark.parametrize(
    ("failure_stage", "unknown_resource", "cleanup", "reconciliation"),
    (
        ("upload_input", "file", "completed", "empty"),
        ("import_input", "document", "unknown", "unknown"),
        ("wait_for_import", "document", "unknown", "unknown"),
    ),
)
def test_public_matrix_one_time_response_loss_keeps_resource_truth_sticky(
    failure_stage: str,
    unknown_resource: str,
    cleanup: str,
    reconciliation: str,
) -> None:
    smoke = _load_smoke()
    session = _OneTimeResponseLossSession(smoke, failure_stage)
    evidence: dict[str, object] = {}

    with pytest.raises(smoke.GeminiProviderError, match="synthetic response loss"):
        asyncio.run(smoke.run_contract_smoke(session, failure_evidence=evidence))

    names = [name for name, _ in session.calls]
    assert names.count("upload_input") == 5
    expected_imports = 4 if failure_stage == "upload_input" else 5
    assert names.count("import_input") == expected_imports
    expected_waits = 4 if failure_stage in {"upload_input", "import_input"} else 5
    assert names.count("wait_for_import") == expected_waits
    assert names.count("query_private") == 0
    assert evidence["resources_created"][unknown_resource] == "unknown"
    assert evidence["cleanup"]["status"] == cleanup
    assert evidence["reconciliation"] == reconciliation
    assert len(evidence["input_results"]) == 5
    assert json.dumps(evidence, sort_keys=True).find("files/synthetic-") == -1
    assert json.dumps(evidence, sort_keys=True).find("fileSearchStores/synthetic") == -1


@pytest.mark.parametrize(
    ("failure_stage", "error_kind"),
    (("wait_for_import", "temporary"), ("import_input", "contract")),
)
def test_public_matrix_non_provider_lifecycle_failure_keeps_document_unknown(
    failure_stage: str,
    error_kind: str,
) -> None:
    smoke = _load_smoke()
    error: BaseException = (
        smoke.SmokeTemporaryFailure("synthetic temporary loss")
        if error_kind == "temporary"
        else smoke.SmokeContractError("synthetic completion contract failure")
    )
    session = _OneTimeResponseLossSession(smoke, failure_stage, error)
    evidence: dict[str, object] = {}

    with pytest.raises(type(error)):
        asyncio.run(smoke.run_contract_smoke(session, failure_evidence=evidence))

    names = [name for name, _ in session.calls]
    assert names.count("upload_input") == 5
    assert names.count("import_input") == 5
    assert names.count("wait_for_import") == (
        4 if failure_stage == "import_input" else 5
    )
    assert names.count("query_private") == 0
    assert evidence["resources_created"]["document"] == "unknown"
    assert evidence["cleanup"]["status"] == "unknown"
    assert evidence["reconciliation"] == "unknown"
    assert len(evidence["input_results"]) == 5
    encoded = json.dumps(evidence, sort_keys=True)
    assert "files/synthetic-" not in encoded
    assert "fileSearchStores/synthetic" not in encoded


def test_private_shadow_query_uses_real_sdk_models_and_maps_direct_evidence(
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    manifest = smoke._private_shadow_manifest(view)
    interactions = _PrivateSdkInteractions(smoke, view, manifest)
    clients: list[_SdkClient] = []

    def sdk_factory(**kwargs: object) -> _SdkClient:
        assert kwargs["api_key"] == "synthetic-sdk-key"
        client = _SdkClient(smoke)
        client.aio.interactions = interactions
        client.aio.file_search_stores.listed_stores = (
            SimpleNamespace(
                name="fileSearchStores/reconciled",
                display_name="task-private-store",
            ),
            SimpleNamespace(name="fileSearchStores/ignored", display_name="other-run"),
        )
        client.aio.files.listed_files = (
            SimpleNamespace(name="files/reconciled", display_name="task-private-upload"),
            SimpleNamespace(name="files/ignored", display_name="other-run"),
        )
        clients.append(client)
        return client

    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key",
        sdk_factory=sdk_factory,
    )
    uploaded = asyncio.run(session.upload_input("lecture.pdf", view.pdf.path, "application/pdf"))
    operation = asyncio.run(
        session.import_input(
            "fileSearchStores/private-store",
            uploaded,
            (("input_key", "normalized_markdown"),),
            {
                "white_space_config": {
                    "max_tokens_per_chunk": 700,
                    "max_overlap_tokens": 100,
                }
            },
        )
    )
    reconciled = asyncio.run(session.find_files(("task-private-upload",)))
    reconciled_stores = asyncio.run(session.find_stores("task-private-store"))
    positive = asyncio.run(
        session.query_private(
            "fileSearchStores/private-store",
            "safe private query",
            smoke.SmokeScope(view.course_id, view.exam_id, view.lecture_id),
            source_revision_id=view.source_revision_id,
            manifest=manifest,
            file_bindings=(("files/sdk-file", "pdf"),),
        )
    )
    negative = asyncio.run(
        session.query_private(
            "fileSearchStores/private-store",
            "safe negative query",
            smoke.SmokeScope(
                view.course_id,
                view.exam_id,
                smoke.PRIVATE_SHADOW_WRONG_LECTURE_ID,
            ),
            source_revision_id=view.source_revision_id,
            manifest=manifest,
            file_bindings=(("files/sdk-file", "pdf"),),
            require_structured_no_result=True,
        )
    )

    assert positive == smoke.PrivateShadowQueryAudit(
        1, 1, 19, 6, True, False, "private response discarded by adapter", 1, "Do not emit this text."
    )
    assert negative == smoke.PrivateShadowQueryAudit(0, 0, 19, 6, False, True, "")
    assert operation == "operations/sdk-operation"
    assert reconciled == ("files/reconciled",)
    assert reconciled_stores == ("fileSearchStores/reconciled",)
    list_call = next(
        client.aio.file_search_stores.calls[0]
        for client in clients
        if client.aio.file_search_stores.calls
        and client.aio.file_search_stores.calls[0][0] == "list"
    )
    assert list_call == ("list", {"page_size": 20})
    upload_call = next(client.aio.files.calls[0] for client in clients if client.aio.files.calls)
    assert upload_call == (
        "upload",
        (
            view.pdf.path.read_bytes(),
            {"display_name": "lecture.pdf", "mime_type": "application/pdf"},
        ),
    )
    import_call = next(
        client.aio.file_search_stores.calls[0]
        for client in clients
        if client.aio.file_search_stores.calls
    )
    assert import_call[0] == "import_file"
    import_args = import_call[1]
    assert import_args[:2] == ("fileSearchStores/private-store", "files/sdk-file")
    assert import_args[2].http_options.extra_body == {
        "customMetadata": [
            {"key": "input_key", "stringValue": "normalized_markdown"}
        ],
        "chunkingConfig": {
            "whiteSpaceConfig": {
                "maxTokensPerChunk": 700,
                "maxOverlapTokens": 100,
            }
        },
    }
    assert set(interactions.calls[0]) == {"input", "model", "store", "tools"}
    assert set(interactions.calls[1]) == {
        "input",
        "model",
        "response_format",
        "store",
        "tools",
    }
    assert interactions.calls[1]["response_format"] == {
        "type": "text",
        "mime_type": "application/json",
        "schema": smoke.SmokeAnswer.model_json_schema(),
    }
    assert all(body["model"] == "gemini-3.7-flash" for body in interactions.calls)
    assert all(body["store"] is False for body in interactions.calls)


def test_interaction_citation_must_bind_to_uploaded_file() -> None:
    smoke = _load_smoke()
    response = SimpleNamespace(
        steps=[
            SimpleNamespace(
                type="model_output",
                content=[
                    SimpleNamespace(
                        type="text",
                        text=smoke.SYNTHETIC_FACT,
                        annotations=[
                            SimpleNamespace(
                                type="file_citation",
                                custom_metadata={
                                    "course_id": smoke.SYNTHETIC_COURSE_ID,
                                    "exam_id": smoke.SYNTHETIC_EXAM_ID,
                                    "lecture_id": smoke.SYNTHETIC_LECTURE_ID,
                                    "source_revision_id": smoke.SYNTHETIC_REVISION_ID,
                                },
                                file_name="other.pdf",
                                page_number=1,
                                source=smoke.SYNTHETIC_FACT,
                            )
                        ],
                    )
                ],
            )
        ]
    )

    with pytest.raises(smoke.SmokeContractError, match="wrong file") as raised:
        smoke._citations(
            response,
            "fileSearchStores/sdk-store",
            smoke.SmokeScope(
                smoke.SYNTHETIC_COURSE_ID,
                smoke.SYNTHETIC_EXAM_ID,
                smoke.SYNTHETIC_LECTURE_ID,
            ),
            "fileSearchStores/sdk-store/documents/sdk-document",
            "task-2-8-synthetic.pdf",
        )
    assert raised.value.reason == "citation_wrong_file"

    annotation = response.steps[0].content[0].annotations[0]
    annotation.file_name = "task-2-8-synthetic.pdf"
    annotation.custom_metadata["lecture_id"] = smoke.WRONG_LECTURE_ID
    with pytest.raises(smoke.SmokeContractError, match="requested scope") as scope:
        smoke._citations(
            response,
            "fileSearchStores/sdk-store",
            smoke.SmokeScope(
                smoke.SYNTHETIC_COURSE_ID,
                smoke.SYNTHETIC_EXAM_ID,
                smoke.SYNTHETIC_LECTURE_ID,
            ),
            "fileSearchStores/sdk-store/documents/sdk-document",
            "task-2-8-synthetic.pdf",
        )
    assert scope.value.reason == "citation_scope_mismatch"


def test_interaction_citation_rejects_conflicting_document_uri() -> None:
    smoke = _load_smoke()
    response = SimpleNamespace(
        steps=[
            SimpleNamespace(
                type="model_output",
                content=[
                    SimpleNamespace(
                        type="text",
                        text=smoke.SYNTHETIC_FACT,
                        annotations=[
                            SimpleNamespace(
                                type="file_citation",
                                custom_metadata={
                                    "course_id": smoke.SYNTHETIC_COURSE_ID,
                                    "exam_id": smoke.SYNTHETIC_EXAM_ID,
                                    "lecture_id": smoke.SYNTHETIC_LECTURE_ID,
                                    "source_revision_id": smoke.SYNTHETIC_REVISION_ID,
                                },
                                document_uri="fileSearchStores/other/documents/other",
                                file_name="task-2-8-synthetic.pdf",
                                page_number=1,
                                source=smoke.SYNTHETIC_FACT,
                            )
                        ],
                    )
                ],
            )
        ]
    )

    with pytest.raises(smoke.SmokeContractError, match="wrong document") as raised:
        smoke._citations(
            response,
            "fileSearchStores/sdk-store",
            smoke.SmokeScope(
                smoke.SYNTHETIC_COURSE_ID,
                smoke.SYNTHETIC_EXAM_ID,
                smoke.SYNTHETIC_LECTURE_ID,
            ),
            "fileSearchStores/sdk-store/documents/sdk-document",
            "task-2-8-synthetic.pdf",
        )
    assert raised.value.reason == "citation_wrong_document"


def test_interaction_citation_accepts_exact_document_uri() -> None:
    smoke = _load_smoke()
    document_name = "fileSearchStores/sdk-store/documents/sdk-document"
    audit = smoke._audit_citations(
        _real_interaction(
            smoke,
            citation=_real_citation(smoke, document_uri=document_name),
        ),
        smoke.SmokeScope(
            smoke.SYNTHETIC_COURSE_ID,
            smoke.SYNTHETIC_EXAM_ID,
            smoke.SYNTHETIC_LECTURE_ID,
        ),
        document_name,
        "task-2-8-synthetic.pdf",
    )

    assert audit.checks["citation_document_binding"] == "passed"


def test_contract_failure_record_retains_only_allowlisted_reason() -> None:
    smoke = _load_smoke()

    record = smoke._failure_record(
        smoke.SmokeContractError(
            "raw provider response must not be retained",
            reason="structured_output_invalid",
        ),
        {"failure_stage": "positive_query"},
    )

    assert record["contract_reason"] == "structured_output_invalid"
    assert "raw provider response" not in json.dumps(record, sort_keys=True)

    provider = smoke.GeminiProviderError(
        "fixed redacted message",
        diagnostic_code="transport_error",
    )
    provider_record = smoke._failure_record(provider, {"failure_stage": "positive_query"})
    assert provider_record["provider_reason"] == "transport_error"


def test_interaction_citation_metadata_and_excerpt_are_bounded() -> None:
    smoke = _load_smoke()

    with pytest.raises(smoke.SmokeContractError, match="metadata was invalid") as metadata:
        smoke._string_metadata({f"key-{index}": "value" for index in range(17)})
    assert metadata.value.reason == "citation_metadata_invalid"
    with pytest.raises(smoke.SmokeContractError, match="excerpt was invalid") as long_excerpt:
        smoke._citation_excerpt(
            SimpleNamespace(source="x" * 4097),
            smoke.SYNTHETIC_FACT,
        )
    assert long_excerpt.value.reason == "citation_excerpt_invalid"
    assert (
        smoke._citation_excerpt(
            SimpleNamespace(source=f"\n{smoke.SYNTHETIC_FACT}\n"),
            smoke.SYNTHETIC_FACT,
        )
        == smoke.SYNTHETIC_FACT
    )
    with pytest.raises(smoke.SmokeContractError, match="excerpt was invalid") as control:
        smoke._citation_excerpt(SimpleNamespace(source="invalid\x00excerpt"), "")
    assert control.value.reason == "citation_excerpt_invalid"


def test_positive_query_parser_failures_have_fixed_redacted_reasons() -> None:
    smoke = _load_smoke()

    with pytest.raises(smoke.SmokeContractError) as missing_excerpt:
        smoke._citation_excerpt(SimpleNamespace(), None)
    with pytest.raises(smoke.SmokeContractError) as invalid_page:
        smoke._optional_page("one")
    with pytest.raises(smoke.SmokeContractError) as invalid_usage:
        smoke._optional_count("one")

    assert missing_excerpt.value.reason == "citation_excerpt_unavailable"
    assert invalid_page.value.reason == "citation_page_invalid"
    assert invalid_usage.value.reason == "usage_count_invalid"


def test_authorized_entrypoint_reads_stored_key_once_without_retaining_it(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets("stored-synthetic-key")
    created: list[str] = []

    def session_factory(api_key: str) -> _FakeSession:
        created.append(api_key)
        return _FakeSession(smoke)

    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")
    record = asyncio.run(
        smoke.run_authorized_live_smoke(
            secret_store=secrets,
            session_factory=session_factory,
        )
    )

    assert secrets.calls == ["gemini-api-key"]
    assert created == ["stored-synthetic-key"]
    assert "stored-synthetic-key" not in json.dumps(record, sort_keys=True)


def test_authorized_entrypoint_uses_shared_manifest_import_path(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()

    class LegacyPdfSession(_FakeSession):
        async def import_input(self, *args: object, **kwargs: object) -> str:
            del args, kwargs
            raise RuntimeError("shared-import-sentinel")

    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")
    with pytest.raises(RuntimeError, match="shared-import-sentinel"):
        asyncio.run(
            smoke.run_authorized_live_smoke(
                secret_store=_FakeSecrets("stored-synthetic-key"),
                session_factory=lambda _: LegacyPdfSession(smoke),
            )
        )


@pytest.mark.parametrize(
    "failed_input",
    (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/pdf",
        "normalized_markdown",
        "image/png",
        "image/jpeg",
    ),
)
def test_public_matrix_collects_each_input_failure_before_cleanup(
    failed_input: str,
) -> None:
    smoke = _load_smoke()
    session = _FiveInputSession(smoke, failed_input)
    evidence: dict[str, object] = {}

    with pytest.raises(smoke.SmokeContractError, match="synthetic import failure"):
        asyncio.run(smoke.run_contract_smoke(session, failure_evidence=evidence))

    uploads = [value for name, value in session.calls if name == "upload_input"]
    imports = [value for name, value in session.calls if name == "import_input"]
    assert [value[2] for value in uploads] == [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/pdf",
        "text/markdown",
        "image/png",
        "image/jpeg",
    ]
    assert len(imports) == 5
    assert not [value for name, value in session.calls if name == "query_private"]
    assert evidence["input_results"] == [
        {
            "input_kind": input_kind,
            "stage": "import_input",
            "outcome": "failed" if index == [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/pdf",
                "normalized_markdown",
                "image/png",
                "image/jpeg",
            ].index(failed_input) else "passed",
            "error_category": "contract" if index == [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/pdf",
                "normalized_markdown",
                "image/png",
                "image/jpeg",
            ].index(failed_input) else "none",
        }
        for index, input_kind in enumerate(("pptx", "pdf", "markdown", "image", "image"))
    ]
    assert evidence["aggregate"] == {
        "input_count": 5,
        "indexed_bytes": sum(value[3] for value in uploads),
        "transient_attempts": 0,
    }
    assert evidence["cleanup"] == {"attempted": 10, "status": "unknown"}
    assert evidence["reconciliation"] == "unknown"
    assert len([value for name, value in session.calls if name == "delete_document"]) == 4
    assert len([value for name, value in session.calls if name == "delete_file"]) == 5
    assert len([value for name, value in session.calls if name == "delete_store"]) == 1
    assert session.live_files == {}
    assert session.live_stores == {}


def test_public_runner_rejects_actual_wrong_marker_without_test_adapter_validation(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()

    class WrongMarkerAudit(_FakeSession):
        async def query_private(self, *args: object, **kwargs: object) -> object:
            result = await super().query_private(*args, **kwargs)
            if kwargs.get("require_structured_no_result"):
                return result
            return SimpleNamespace(
                citation_count=1,
                resolved_citation_count=1,
                input_tokens=11,
                output_tokens=7,
                supported=True,
                answer_empty=None,
                answer="wrong-marker",
            )

    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")
    with pytest.raises(smoke.SmokeContractError, match="marker"):
        asyncio.run(
            smoke.run_authorized_live_smoke(
                secret_store=_FakeSecrets("stored-synthetic-key"),
                session_factory=lambda _: WrongMarkerAudit(smoke),
            )
        )


def test_public_success_omits_provider_identity_projection() -> None:
    smoke = _load_smoke()
    session = _FakeSession(smoke)
    session.store_name = "fileSearchStores/nondefault-store"
    session.file_name = "files/nondefault-file"
    session.operation_name = "operations/nondefault-operation"
    session.document_name = "fileSearchStores/nondefault-store/documents/nondefault-document"

    record = asyncio.run(smoke.run_contract_smoke(session))

    assert "provider_ids" not in record


def test_public_success_commits_only_safe_matrix_fields() -> None:
    smoke = _load_smoke()

    class ActualAudit(_FakeSession):
        async def query_private(self, *args: object, **kwargs: object) -> object:
            result = await super().query_private(*args, **kwargs)
            if kwargs.get("require_structured_no_result"):
                return result
            return SimpleNamespace(
                citation_count=1,
                resolved_citation_count=1,
                input_tokens=11,
                output_tokens=7,
                supported=True,
                answer_empty=None,
                answer=smoke.SYNTHETIC_FACT,
                citation_page=7,
                citation_excerpt="actual synthetic excerpt",
            )

    record = asyncio.run(smoke.run_contract_smoke(ActualAudit(smoke)))

    assert set(record) == {
        "schema_version",
        "status",
        "input_results",
        "aggregate",
        "cleanup",
        "reconciliation",
        "error_category",
    }


def test_public_failure_checks_do_not_claim_unobserved_passes() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class FailsBeforeQuery(_FakeSession):
        async def create_store(self, *args: object, **kwargs: object) -> str:
            del args, kwargs
            raise smoke.GeminiProviderError("synthetic provider failure")

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(smoke.run_contract_smoke(FailsBeforeQuery(smoke), failure_evidence=evidence))

    checks = evidence["checks"]
    assert checks["document_listing"] != "passed"
    assert checks["cleanup_document"] != "passed"


def test_public_runner_rejects_provider_unsupported_positive_audit() -> None:
    smoke = _load_smoke()

    class UnsupportedAudit(_FakeSession):
        async def query_private(self, *args: object, **kwargs: object) -> object:
            audit = await super().query_private(*args, **kwargs)
            if kwargs.get("require_structured_no_result"):
                return audit
            return dataclasses.replace(audit, supported=False)

    with pytest.raises(smoke.SmokeContractError, match="unsupported"):
        asyncio.run(smoke.run_contract_smoke(UnsupportedAudit(smoke)))


def test_public_failure_does_not_mark_negative_checks_passed_before_query() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class PositiveFailure(_FakeSession):
        async def query_private(self, *args: object, **kwargs: object) -> object:
            if kwargs.get("require_structured_no_result"):
                pytest.fail("negative query must not run")
            raise smoke.SmokeContractError("positive failure", reason="positive_answer_invalid")

    with pytest.raises(smoke.SmokeContractError):
        asyncio.run(smoke.run_contract_smoke(PositiveFailure(smoke), failure_evidence=evidence))

    assert evidence["checks"]["negative_structured_output"] != "passed"
    assert evidence["checks"]["wrong_lecture_filtering"] != "passed"


def test_public_cleanup_only_failure_has_public_observed_evidence() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class CleanupFailure(_FakeSession):
        async def delete_file(self, file_name: str) -> None:
            await super().delete_file(file_name)
            raise RuntimeError("cleanup failure")

    with pytest.raises(smoke.SmokeContractError, match="cleanup"):
        asyncio.run(smoke.run_contract_smoke(CleanupFailure(smoke), failure_evidence=evidence))

    assert evidence["resources_created"] == {
        "document": "confirmed", "file": "confirmed", "store": "confirmed"
    }
    assert evidence["cleanup"]["status"] == "failed"
    assert evidence["checks"]["cleanup_document"] == "passed"
    assert evidence["checks"]["cleanup_file"] == "cleanup_delete_failed"
    assert evidence["checks"]["cleanup_store"] == "passed"


@pytest.mark.parametrize(
    ("failure_stage", "uncertain_resource"),
    (
        ("upload_input", "file"),
        ("import_input", "document"),
        ("wait_for_import", "document"),
    ),
)
def test_public_response_loss_retains_unknown_resource_and_cleanup_outcomes(
    failure_stage: str,
    uncertain_resource: str,
) -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class ResponseLoss(_FakeSession):
        def __init__(self) -> None:
            super().__init__(smoke)
            self.file_reconciliations = 0

        async def upload_input(self, *args: object, **kwargs: object) -> str:
            value = await super().upload_input(*args, **kwargs)
            if failure_stage == "upload_input":
                raise smoke.GeminiProviderError("synthetic upload response loss")
            return value

        async def import_input(self, *args: object, **kwargs: object) -> str:
            value = await super().import_input(*args, **kwargs)
            if failure_stage == "import_input":
                raise smoke.GeminiProviderError("synthetic import response loss")
            return value

        async def wait_for_import(self, *args: object, **kwargs: object) -> str:
            value = await super().wait_for_import(*args, **kwargs)
            if failure_stage == "wait_for_import":
                raise smoke.GeminiProviderError("synthetic wait response loss")
            return value

        async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]:
            self.file_reconciliations += 1
            if self.file_reconciliations == 2:
                raise RuntimeError("synthetic reconciliation loss")
            return await super().find_files(display_names)

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(smoke.run_contract_smoke(ResponseLoss(), failure_evidence=evidence))

    assert {
        key: evidence[key]
        for key in ("failure_stage", "resources_created", "cleanup")
    } == {
        "failure_stage": failure_stage,
        "resources_created": {
            "document": "unknown" if uncertain_resource == "document" else "not_started",
            "file": "unknown" if uncertain_resource == "file" else "confirmed",
            "store": "confirmed",
        },
        "cleanup": {
            "attempted": 1 if failure_stage == "upload_input" else 6,
            "status": "unknown",
        },
    }


@pytest.mark.parametrize(
    ("answer", "supported", "citation", "expected_reason"),
    (
        ("cobalt-otter-28", False, "complete", "positive_answer_unsupported"),
        ("wrong-marker", True, "complete", "positive_answer_missing_marker"),
        ("cobalt-otter-28", True, "missing", "positive_citation_missing"),
        ("cobalt-otter-28", True, "page_missing", "citation_page_absent"),
        ("cobalt-otter-28", True, "excerpt_missing", "citation_excerpt_absent"),
    ),
)
def test_public_positive_failures_keep_distinct_safe_reasons(
    answer: str,
    supported: bool,
    citation: str,
    expected_reason: str,
) -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class ProviderShape(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id != smoke.SYNTHETIC_LECTURE_ID:
                return result
            citations = result.citations
            if citation == "missing":
                citations = ()
            elif citation == "page_missing":
                citations = (dataclasses.replace(citations[0], page_number=None),)
            elif citation == "excerpt_missing":
                citations = (dataclasses.replace(citations[0], excerpt=""),)
            return smoke.SmokeQueryResult(
                answer={"answer": answer, "supported": supported},
                citations=citations,
                input_tokens=11,
                output_tokens=7,
            )

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_contract_smoke(
                ProviderShape(smoke), failure_evidence=evidence
            )
        )

    assert smoke._failure_record(raised.value, evidence)["contract_reason"] == expected_reason


def test_public_positive_malformed_provider_shape_is_not_a_negative_failure() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class MalformedPositive(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                return smoke.SmokeQueryResult(answer={}, citations=result.citations)
            return result

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_contract_smoke(
                MalformedPositive(smoke), failure_evidence=evidence
            )
        )

    assert smoke._failure_record(raised.value, evidence)["contract_reason"] == (
        "structured_output_invalid"
    )


def test_public_runner_forwards_contract_diagnostic_lifecycle_events() -> None:
    smoke = _load_smoke()
    captured: dict[str, object] = {}

    class Sink:
        def capture(self, label: str, value: object) -> None:
            captured[label] = value

        def capture_exception(self, label: str, error: BaseException) -> None:
            del label, error

    asyncio.run(smoke.run_contract_smoke(_FakeSession(smoke), diagnostic_sink=Sink()))

    assert captured["contract.expected"] == {
        "course_id": smoke.SYNTHETIC_COURSE_ID,
        "exam_id": smoke.SYNTHETIC_EXAM_ID,
        "lecture_id": smoke.SYNTHETIC_LECTURE_ID,
        "source_revision_id": smoke.SYNTHETIC_REVISION_ID,
        "fixture_sha256": hashlib.sha256(smoke.synthetic_pdf_bytes()).hexdigest(),
    }
    assert captured["contract.check_matrix"] == {
        "positive_answer": "passed",
        "citation_presence": "passed",
        "negative_structured_output": "passed",
        "create_store": "passed",
        "document_listing": "not_run",
        "cleanup_store": "passed",
        "cleanup_document": "passed",
        "cleanup_file": "passed",
        "wrong_lecture_filtering": "passed",
    }


def test_shared_sdk_paths_label_diagnostic_provider_failures(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    smoke = _load_smoke()
    captures: list[str] = []

    class Sink:
        def capture_exception(self, label: str, error: BaseException) -> None:
            del error
            captures.append(label)

    async def fail_provider(
        request: object, *, diagnostic_sink: object | None = None, label: str = "", **kwargs: object
    ) -> object:
        del request, kwargs
        if diagnostic_sink is not None:
            diagnostic_sink.capture_exception(label, RuntimeError("synthetic"))
        raise smoke.GeminiProviderError("synthetic")

    monkeypatch.setattr(smoke, "_provider_call", fail_provider)
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key", sdk_factory=lambda **kwargs: _SdkClient(smoke), diagnostic_sink=Sink()
    )
    path = tmp_path / "input.pdf"
    path.write_bytes(smoke.synthetic_pdf_bytes())
    for call in (
        lambda: session.upload_input("input.pdf", path, "application/pdf"),
        lambda: session.import_input("fileSearchStores/test", "files/test", (), None),
        lambda: session.query_private(
            "fileSearchStores/test", "prompt", smoke.SmokeScope("course", "exam", "lecture"),
            source_revision_id="revision", manifest=SimpleNamespace(inputs=()), file_bindings=(),
        ),
    ):
        with pytest.raises(smoke.GeminiProviderError):
            asyncio.run(call())

    assert captures == ["upload_input", "import_input", "query_private"]


@pytest.mark.parametrize(
    ("resource", "failure", "label"),
    (
        ("stores", "request", "find_stores.request.failed"),
        ("stores", "iteration", "find_stores.iteration"),
        ("files", "request", "find_files.request.failed"),
        ("files", "iteration", "find_files.iteration"),
    ),
)
def test_reconciliation_diagnostics_capture_request_and_pager_failures(
    resource: str,
    failure: str,
    label: str,
) -> None:
    smoke = _load_smoke()

    class Sink:
        def __init__(self) -> None:
            self.labels: list[str] = []

        def capture(self, label: str, value: object) -> None:
            del label, value

        def capture_exception(self, label: str, error: BaseException) -> None:
            del error
            self.labels.append(label)

    class RequestFailure:
        async def list(self, *, config: object) -> object:
            del config
            raise RuntimeError("synthetic request failure")

    class IterationFailure:
        async def list(self, *, config: object) -> object:
            del config

            async def entries() -> object:
                raise RuntimeError("synthetic pager failure")
                yield SimpleNamespace()

            return entries()

    def session_with(resource_api: object, sink: Sink) -> object:
        client = _SdkClient(smoke)
        if resource == "stores":
            client.aio.file_search_stores = resource_api
        else:
            client.aio.files = resource_api
        return smoke.GoogleGenaiSmokeSession(
            "synthetic-sdk-key",
            sdk_factory=lambda **kwargs: client,
            diagnostic_sink=sink,
        )

    sink = Sink()
    session = session_with(
        RequestFailure() if failure == "request" else IterationFailure(), sink
    )
    expected_error = smoke.GeminiProviderError
    with pytest.raises(expected_error):
        asyncio.run(
            session.find_stores("target")
            if resource == "stores"
            else session.find_files(("target",))
        )
    assert sink.labels == [label]


def test_reconciliation_pagers_never_capture_successful_resource_values() -> None:
    smoke = _load_smoke()
    store_id = "fileSearchStores/raw-store-sentinel"
    store_display_name = "raw-store-display-sentinel"
    file_id = "files/raw-file-sentinel"
    file_display_name = "raw-file-display-sentinel"

    class Sink:
        def __init__(self) -> None:
            self.values: list[object] = []

        def capture(self, label: str, value: object) -> None:
            del label
            self.values.append(value)

        def capture_exception(self, label: str, error: BaseException) -> None:
            del label, error

    class Pager:
        def __init__(self, item: object, raw: str) -> None:
            self.item = item
            self.raw = raw

        def __repr__(self) -> str:
            return self.raw

        def __aiter__(self) -> object:
            async def entries() -> object:
                yield self.item

            return entries()

    class Stores:
        async def list(self, *, config: object) -> object:
            del config
            return Pager(
                SimpleNamespace(name=store_id, display_name=store_display_name),
                f"{store_id}:{store_display_name}",
            )

    class Files:
        async def list(self, *, config: object) -> object:
            del config
            return Pager(
                SimpleNamespace(name=file_id, display_name=file_display_name),
                f"{file_id}:{file_display_name}",
            )

    client = _SdkClient(smoke)
    client.aio.file_search_stores = Stores()
    client.aio.files = Files()
    sink = Sink()
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key",
        sdk_factory=lambda **kwargs: client,
        diagnostic_sink=sink,
    )

    assert asyncio.run(session.find_stores(store_display_name)) == (store_id,)
    assert asyncio.run(session.find_files((file_display_name,))) == (file_id,)
    assert sink.values == []
    assert all(
        sentinel not in repr(value)
        for sentinel in (store_id, store_display_name, file_id, file_display_name)
        for value in sink.values
    )


def test_authorized_entrypoint_fails_closed_when_stored_key_is_missing(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets(None)
    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")

    with pytest.raises(smoke.LiveSmokeBlocked, match="stored Gemini credential is unavailable"):
        asyncio.run(smoke.run_authorized_live_smoke(secret_store=secrets))

    assert secrets.calls == ["gemini-api-key"]


def test_completed_operation_error_keeps_safe_provider_classification() -> None:
    smoke = _load_smoke()

    def sdk_factory(**kwargs: object) -> _SdkClient:
        del kwargs
        return _SdkClient(smoke, operation_error_status=429)

    session = smoke.GoogleGenaiSmokeSession("synthetic-sdk-key", sdk_factory=sdk_factory)

    with pytest.raises(smoke.GeminiProviderError) as raised:
        asyncio.run(session.wait_for_import("operations/sdk-operation"))

    assert type(raised.value).__name__ == "GeminiQuotaError"
    assert "429" not in str(raised.value)


def test_operation_poll_is_bounded_by_remaining_deadline(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()
    ticks = iter((0.0, 899.999))
    monkeypatch.setattr(smoke, "monotonic", lambda: next(ticks))

    def sdk_factory(**kwargs: object) -> _SdkClient:
        del kwargs
        return _SdkClient(smoke, operation_delay=0.02)

    session = smoke.GoogleGenaiSmokeSession("synthetic-sdk-key", sdk_factory=sdk_factory)

    with pytest.raises(smoke.SmokeTemporaryFailure, match="timed out"):
        asyncio.run(session.wait_for_import("operations/sdk-operation"))


def test_import_poll_retry_does_not_sleep_past_its_existing_deadline() -> None:
    smoke = _load_smoke()
    delays: list[float] = []

    class TransientOperations(_SdkOperations):
        async def get(self, operation: object) -> object:
            self.calls.append(operation)
            raise _transient_500()

    async def sleep(delay: float) -> None:
        delays.append(delay)

    ticks = iter((0.0, 0.0, 899.5))
    client = _SdkClient(smoke)
    client.aio.operations = TransientOperations()
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key",
        sdk_factory=lambda **kwargs: client,
        sleep=sleep,
        clock=lambda: next(ticks),
    )

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(session.wait_for_import("operations/sdk-operation"))

    assert len(client.aio.operations.calls) == 1
    assert delays == []


def test_idempotent_query_retries_transient_failures_without_mutations() -> None:
    smoke = _load_smoke()
    delays: list[float] = []

    class TransientInteractions(_SdkInteractions):
        async def create(self, **body: object) -> object:
            if len(self.calls) < 2:
                self.calls.append(body)
                raise _transient_500()
            return await super().create(**body)

    async def sleep(delay: float) -> None:
        delays.append(delay)

    client = _SdkClient(smoke)
    client.aio.interactions = TransientInteractions(smoke)
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key", sdk_factory=lambda **kwargs: client, sleep=sleep
    )
    session._store_name = "fileSearchStores/sdk-store"
    session._document_name = "fileSearchStores/sdk-store/documents/sdk-document"
    session._file_name = "files/sdk-file"

    result = asyncio.run(
        session.query_private(
            "fileSearchStores/sdk-store",
            "prompt",
            smoke.SmokeScope(
                smoke.SYNTHETIC_COURSE_ID,
                smoke.SYNTHETIC_EXAM_ID,
                smoke.SYNTHETIC_LECTURE_ID,
            ),
            source_revision_id=smoke.SYNTHETIC_REVISION_ID,
            manifest=SimpleNamespace(inputs=()),
            file_bindings=(),
            require_structured_supported=True,
        )
    )

    assert result.supported is True
    assert len(client.aio.interactions.calls) == 3
    assert delays == [1.0, 2.0]
    assert session.transient_attempts == 2
    assert not client.aio.file_search_stores.calls
    assert not client.aio.files.calls


def test_create_store_transient_failure_is_not_retried() -> None:
    smoke = _load_smoke()
    delays: list[float] = []

    class TransientStores(_SdkStores):
        async def create(self, *, config: object) -> object:
            self.calls.append(("create", config))
            raise _transient_500()

    async def sleep(delay: float) -> None:
        delays.append(delay)

    client = _SdkClient(smoke)
    client.aio.file_search_stores = TransientStores()
    session = smoke.GoogleGenaiSmokeSession(
        "synthetic-sdk-key", sdk_factory=lambda **kwargs: client, sleep=sleep
    )

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(session.create_store("synthetic", "models/gemini-embedding-2"))

    assert [call[0] for call in client.aio.file_search_stores.calls] == ["create"]
    assert delays == []
    assert session.transient_attempts == 0


def test_completed_import_rejects_full_document_from_another_store() -> None:
    smoke = _load_smoke()

    class CrossStoreOperations(_SdkOperations):
        async def get(self, operation: object) -> object:
            self.calls.append(operation)
            return SimpleNamespace(
                name="operations/sdk-operation",
                done=True,
                error=None,
                response=SimpleNamespace(
                    parent="other-store",
                    document_name="fileSearchStores/other-store/documents/sdk-document",
                ),
            )

    def sdk_factory(**kwargs: object) -> _SdkClient:
        del kwargs
        client = _SdkClient(smoke)
        client.aio.operations = CrossStoreOperations()
        return client

    session = smoke.GoogleGenaiSmokeSession("synthetic-sdk-key", sdk_factory=sdk_factory)
    session._store_name = "fileSearchStores/sdk-store"

    with pytest.raises(smoke.SmokeContractError, match="did not match the store"):
        asyncio.run(session.wait_for_import("operations/sdk-operation"))


def test_primary_provider_failure_wins_when_cleanup_also_fails() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class BodyAndCleanupFailure(_FakeSession):
        async def delete_file(self, file_name: str) -> None:
            self.calls.append(("delete_file", file_name))
            raise smoke.SmokeContractError("synthetic cleanup failure")

    session = BodyAndCleanupFailure(smoke, fail_import=True)

    with pytest.raises(smoke.SmokeTemporaryFailure, match="temporary"):
        asyncio.run(smoke.run_contract_smoke(session, failure_evidence=evidence))

    assert [name for name, _ in session.calls][-2:] == ["delete_file", "delete_store"]
    assert evidence["cleanup"] == {"attempted": 10, "status": "unknown"}
    assert evidence["reconciliation"] == "unknown"


def test_terminal_query_transient_is_classified_after_cleanup_without_duplicate_mutations() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}
    delays: list[float] = []

    class AlwaysTransientInteractions(_PrivateSdkInteractions):
        async def create(self, **body: object) -> object:
            self.calls.append(body)
            raise _transient_500()

    async def sleep(delay: float) -> None:
        delays.append(delay)

    with tempfile.TemporaryDirectory() as directory:
        view = smoke._synthetic_index_input(Path(directory))
        manifest = smoke._private_shadow_manifest(view)
        client = _SdkClient(smoke)
        client.aio.interactions = AlwaysTransientInteractions(smoke, view, manifest)
        session = smoke.GoogleGenaiSmokeSession(
            "synthetic-sdk-key", sdk_factory=lambda **kwargs: client, sleep=sleep
        )

        with pytest.raises(smoke.GeminiTransientError):
            asyncio.run(
                smoke._run_shadow_sequence(
                    session,
                    view,
                    smoke._private_shadow_preflight_from_view(view),
                    mode="private_acceptance",
                    clock=lambda: 100.0,
                    failure_evidence=evidence,
                )
            )

    assert len(client.aio.interactions.calls) == 3
    assert delays == [1.0, 2.0]
    assert evidence["failure_class"] == "infrastructure_transient"
    assert evidence["transient_attempts"] == 2
    assert evidence["provider_cleanup_outcome"] == "complete"
    assert [call[0] for call in client.aio.file_search_stores.calls].count("create") == 1
    assert [call[0] for call in client.aio.file_search_stores.calls].count("import_file") == 5
    assert [call[0] for call in client.aio.files.calls].count("upload") == 5
    assert [call[0] for call in client.aio.file_search_stores.documents.calls].count("delete") == 5
    assert [call[0] for call in client.aio.files.calls].count("delete") == 5
    assert [call[0] for call in client.aio.file_search_stores.calls].count("delete") == 1


def test_failed_smoke_emits_only_redacted_stage_error_and_cleanup_evidence() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class QueryFailure(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            del args, kwargs
            self.calls.append(("query", "private-query-payload"))
            raise smoke.GeminiProviderError(
                "Gemini provider request failed.",
                provider_status_code=400,
            )

    session = QueryFailure(smoke)
    with pytest.raises(smoke.GeminiProviderError) as raised:
        asyncio.run(smoke.run_contract_smoke(session, failure_evidence=evidence))

    record = smoke._failure_record(raised.value, evidence)
    assert {key: record[key] for key in record if key not in {"checks", "input_results", "aggregate", "reconciliation"}} == {
        "schema_version": 1,
        "status": "failed",
        "failure_stage": "positive_query",
        "error_category": "provider",
        "provider_status_code": 400,
        "retryable": False,
        "resources_created": {
            "document": "confirmed",
            "file": "confirmed",
            "store": "confirmed",
        },
        "cleanup": {"attempted": 11, "status": "completed"},
    }
    assert record["aggregate"]["input_count"] == 5
    assert record["reconciliation"] == "empty"
    assert record["checks"]["positive_answer"] == "positive_query_failed"
    assert record["checks"]["negative_structured_output"] == "not_run"
    assert record["checks"]["document_listing"] == "not_run"
    encoded = json.dumps(record, sort_keys=True)
    assert "private-query-payload" not in encoded
    assert smoke.SYNTHETIC_FACT not in encoded
    for raw_identity in (session.store_name, session.file_name, session.document_name):
        assert raw_identity not in encoded


def test_live_cli_failure_prints_redacted_json_without_traceback(
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    smoke = _load_smoke()

    async def fail_live(*, failure_evidence: dict[str, object]) -> dict[str, object]:
        failure_evidence.update(
            {
                "failure_stage": "positive_query",
                "resources_created": {
                    "document": "confirmed",
                    "file": "confirmed",
                    "store": "confirmed",
                },
                "cleanup": {"attempted": 3, "status": "completed"},
            }
        )
        raise smoke.GeminiProviderError(
            "Gemini provider request failed.",
            provider_status_code=400,
        )

    monkeypatch.setattr(smoke, "run_authorized_live_smoke", fail_live)

    assert smoke.main(["--execute-live"]) == 1
    record = json.loads(capsys.readouterr().out)
    assert record["failure_stage"] == "positive_query"
    assert record["provider_status_code"] == 400
    assert record["cleanup"] == {"attempted": 3, "status": "completed"}


def test_malformed_structured_output_is_redacted_at_cli_boundary(
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
) -> None:
    smoke = _load_smoke()

    class MalformedAnswer(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            del args, kwargs
            return smoke.SmokeQueryResult(
                answer={"private-provider-body": "must-not-escape"},
                citations=(),
            )

    async def fail_live(*, failure_evidence: dict[str, object]) -> dict[str, object]:
        return await smoke.run_contract_smoke(
            MalformedAnswer(smoke),
            failure_evidence=failure_evidence,
        )

    monkeypatch.setattr(smoke, "run_authorized_live_smoke", fail_live)

    assert smoke.main(["--execute-live"]) == 1
    output = capsys.readouterr()
    assert json.loads(output.out)["error_category"] == "contract"
    assert output.err == ""
    assert "private-provider-body" not in output.out
    assert "must-not-escape" not in output.out


def test_response_loss_records_unknown_resource_and_cleanup_outcome() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class StoreResponseLoss(_FakeSession):
        async def create_store(self, display_name: str, embedding_model: str) -> str:
            del display_name, embedding_model
            raise smoke.translate_gemini_error(TimeoutError("private response loss"))

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(
            smoke.run_contract_smoke(
                StoreResponseLoss(smoke),
                failure_evidence=evidence,
            )
        )

    assert {key: evidence[key] for key in evidence if key not in {"input_results", "aggregate", "reconciliation", "checks"}} == {
        "failure_stage": "create_store",
        "resources_created": {
            "document": "not_started",
            "file": "not_started",
            "store": "unknown",
        },
        "cleanup": {"attempted": 0, "status": "unknown"},
    }
    assert evidence["aggregate"]["input_count"] == 5
    assert evidence["input_results"] == []
    assert evidence["reconciliation"] == "empty"
    assert evidence["checks"]["create_store"] == "create_store_failed"
    assert evidence["checks"]["cleanup_store"] == "not_available"


def _clock() -> Iterator[float]:
    yield 100.0
    yield 101.25


def test_positive_validation_accepts_multiple_matching_citations() -> None:
    smoke = _load_smoke()

    class MultipleCitations(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer=result.answer,
                    citations=result.citations * 2,
                    input_tokens=result.input_tokens,
                    output_tokens=result.output_tokens,
                )
            return result

    record = asyncio.run(smoke.run_contract_smoke(MultipleCitations(smoke)))

    assert record["status"] == "passed"
    assert all(item["outcome"] == "passed" for item in record["input_results"])


def test_positive_validation_failure_retains_fixed_redacted_reason() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class NoCitations(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                return smoke.SmokeQueryResult(answer=result.answer, citations=())
            return result

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_contract_smoke(
                NoCitations(smoke),
                failure_evidence=evidence,
            )
        )

    record = smoke._failure_record(raised.value, evidence)
    assert record["contract_reason"] == "positive_citation_missing"
    assert smoke.SYNTHETIC_FACT not in json.dumps(record, sort_keys=True)


def test_positive_query_does_not_disclose_the_expected_answer() -> None:
    smoke = _load_smoke()

    class PromptCapture(_FakeSession):
        def __init__(self) -> None:
            super().__init__(smoke)
            self.prompts: list[str] = []

        async def query(self, *args: object, **kwargs: object) -> object:
            self.prompts.append(args[1])
            return await super().query(*args, **kwargs)

    session = PromptCapture()
    asyncio.run(smoke.run_contract_smoke(session))

    assert smoke.SYNTHETIC_FACT not in session.prompts[0]
    assert "cobalt-otter-28" not in session.prompts[0]


def test_positive_citation_and_negative_schema_use_separate_existing_queries() -> None:
    smoke = _load_smoke()

    class SchemaCapture(_FakeSession):
        def __init__(self) -> None:
            super().__init__(smoke)
            self.schemas: list[object] = []

        async def query(self, *args: object, **kwargs: object) -> object:
            self.schemas.append(kwargs["response_schema"])
            return await super().query(*args, **kwargs)

    session = SchemaCapture()
    asyncio.run(smoke.run_contract_smoke(session))

    assert session.schemas == [smoke.SmokeAnswer, smoke.SmokeAnswer]


def test_public_positive_adapter_returns_provider_audit_for_runner_validation() -> None:
    smoke = _load_smoke()

    class ProviderShape(_FakeSession):
        def __init__(self) -> None:
            super().__init__(smoke)
            self.schemas: list[object] = []

        async def query(self, *args: object, **kwargs: object) -> object:
            self.schemas.append(kwargs["response_schema"])
            return smoke.SmokeQueryResult(
                answer={"answer": "wrong-marker", "supported": False},
                citations=(),
                input_tokens=11,
                output_tokens=7,
            )

    session = ProviderShape()
    audit = asyncio.run(
        session.query_private(
            session.store_name,
            "safe public prompt",
            smoke.SmokeScope(
                smoke.SYNTHETIC_COURSE_ID,
                smoke.SYNTHETIC_EXAM_ID,
                smoke.SYNTHETIC_LECTURE_ID,
            ),
            source_revision_id=smoke.SYNTHETIC_REVISION_ID,
            manifest=SimpleNamespace(),
            file_bindings=(),
            require_structured_supported=True,
        )
    )

    assert session.schemas == [smoke.SmokeAnswer]
    assert audit == smoke.PrivateShadowQueryAudit(
        0, 0, 11, 7, False, False, "wrong-marker", None, None
    )


def test_negative_adapter_returns_provider_audit_for_runner_validation() -> None:
    smoke = _load_smoke()

    class ProviderShape(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            del args, kwargs
            return smoke.SmokeQueryResult(
                answer={"answer": "unexpected", "supported": True},
                citations=(),
                input_tokens=11,
                output_tokens=7,
            )

    session = ProviderShape(smoke)
    audit = asyncio.run(
        session.query_private(
            session.store_name,
            "safe negative prompt",
            smoke.SmokeScope(
                smoke.SYNTHETIC_COURSE_ID,
                smoke.SYNTHETIC_EXAM_ID,
                smoke.WRONG_LECTURE_ID,
            ),
            source_revision_id=smoke.SYNTHETIC_REVISION_ID,
            manifest=SimpleNamespace(),
            file_bindings=(),
            require_structured_no_result=True,
        )
    )

    assert audit == smoke.PrivateShadowQueryAudit(
        0, 0, 11, 7, True, False, "unexpected", None, None
    )


def test_public_runner_rejects_negative_policy_after_adapter_audit() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class InvalidNegative(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.WRONG_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer={"answer": "unexpected", "supported": True},
                    citations=(),
                    input_tokens=11,
                    output_tokens=7,
                )
            return result

    with pytest.raises(smoke.SmokeContractError) as raised:
        with tempfile.TemporaryDirectory() as directory:
            view = smoke._synthetic_index_input(Path(directory))
            asyncio.run(
                smoke._run_shadow_sequence(
                    InvalidNegative(smoke),
                    view,
                    smoke._private_shadow_preflight_from_view(view),
                    mode="public_matrix",
                    clock=smoke.monotonic,
                    failure_evidence=evidence,
                )
            )

    assert smoke._failure_record(raised.value, evidence)["contract_reason"] == (
        "private_wrong_scope_retrieved"
    )


def test_negative_structured_answer_must_report_unsupported() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class SupportedNegative(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.WRONG_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer={"answer": "unsupported answer", "supported": True},
                    citations=(),
                )
            return result

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_contract_smoke(
                SupportedNegative(smoke),
                failure_evidence=evidence,
            )
        )

    assert smoke._failure_record(raised.value, evidence)["contract_reason"] == (
        "private_wrong_scope_retrieved"
    )


def test_negative_structured_answer_must_be_empty() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class NonemptyNegative(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.WRONG_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer={"answer": "unrelated text", "supported": False},
                    citations=(),
                )
            return result

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_contract_smoke(
                NonemptyNegative(smoke),
                failure_evidence=evidence,
            )
        )

    assert smoke._failure_record(raised.value, evidence)["contract_reason"] == (
        "private_wrong_scope_retrieved"
    )


def test_positive_validation_accepts_the_retrieved_marker_value() -> None:
    smoke = _load_smoke()

    class MarkerOnlyAnswer(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer={"answer": "cobalt-otter-28", "supported": True},
                    citations=result.citations,
                )
            return result

    record = asyncio.run(smoke.run_contract_smoke(MarkerOnlyAnswer(smoke)))

    assert record["status"] == "passed"


def test_positive_answer_without_marker_retains_fixed_redacted_reason() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class WrongAnswer(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer={"answer": "wrong-marker", "supported": True},
                    citations=result.citations,
                )
            return result

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_contract_smoke(WrongAnswer(smoke), failure_evidence=evidence)
        )

    record = smoke._failure_record(raised.value, evidence)
    assert record["contract_reason"] == "positive_answer_missing_marker"
    assert smoke.SYNTHETIC_FACT not in json.dumps(record, sort_keys=True)


def test_offline_fake_proves_full_smoke_sequence_and_redacted_record() -> None:
    smoke = _load_smoke()
    session = _FakeSession(smoke)
    clock = _clock()

    record = asyncio.run(smoke.run_contract_smoke(session, clock=lambda: next(clock)))

    encoded = json.dumps(record, sort_keys=True)
    assert record["status"] == "passed"
    assert [item["input_kind"] for item in record["input_results"]] == [
        "pptx",
        "pdf",
        "markdown",
        "image",
        "image",
    ]
    assert all(item["outcome"] == "passed" for item in record["input_results"])
    assert smoke.SYNTHETIC_FACT not in encoded
    assert record["aggregate"]["input_count"] == 5
    assert record["aggregate"]["indexed_bytes"] > len(smoke.synthetic_pdf_bytes())
    assert record["cleanup"] == {"attempted": 11, "status": "completed"}
    assert record["reconciliation"] == "empty"
    for raw_identity in (
        session.store_name,
        session.file_name,
        session.operation_name,
        session.document_name,
    ):
        assert raw_identity not in encoded
    names = [name for name, _ in session.calls]
    assert names.count("upload_input") == 5
    assert names.count("import_input") == 5
    assert names.count("wait_for_import") == 5
    assert names.count("query") == 2
    assert names.count("delete_document") == 5
    assert names.count("delete_file") == 5
    assert names.count("delete_store") == 1


def test_temporary_failure_fake_cleans_up_without_a_live_outage() -> None:
    smoke = _load_smoke()
    session = _FakeSession(smoke, fail_import=True)

    with pytest.raises(smoke.SmokeTemporaryFailure, match="temporary"):
        asyncio.run(smoke.run_contract_smoke(session))

    names = [name for name, _ in session.calls]
    assert names.count("upload_input") == 5
    assert names.count("import_input") == 5
    assert names.count("wait_for_import") == 4
    assert names.count("query") == 0
    assert names.count("delete_document") == 4
    assert names.count("delete_file") == 5
    assert names.count("delete_store") == 1


def test_temporary_failure_fixture_persists_retry_state_then_resumes() -> None:
    smoke = _load_smoke()

    record = smoke.run_temporary_failure_fixture()

    assert record == {
        "first_state": "retryable_failure",
        "retry_count": 1,
        "error_category": "transient",
        "backoff_seconds": 5,
        "resumed_state": "ready",
        "service_calls": 2,
    }


def test_locked_environment_exposes_exact_google_genai_models() -> None:
    sdk = _sdk_interactions()
    assert importlib.metadata.version("google-genai") == "2.14.0"
    assert sdk.Interaction.model_fields["steps"].default is None
    assert sdk.ModelOutputStep.model_fields["content"].default is None
    assert sdk.TextContent.model_fields["annotations"].default is None
    assert set(sdk.FileCitation.model_fields) >= {
        "custom_metadata",
        "document_uri",
        "file_name",
        "page_number",
        "source",
        "start_index",
        "end_index",
        "type",
    }
    assert sdk.FileCitation(type="file_citation").type == "file_citation"
    assert set(sdk.Usage.model_fields) >= {
        "total_input_tokens",
        "total_output_tokens",
    }


@pytest.mark.parametrize(
    ("case", "check", "diagnosis"),
    [
        ("steps", "citation_presence", "citation_steps_absent"),
        ("content", "citation_presence", "citation_content_absent"),
        ("annotations", "citation_presence", "citation_annotations_absent"),
        ("metadata", "citation_scope_binding", "citation_metadata_absent"),
    ],
)
def test_real_sdk_optional_citation_containers_have_distinct_diagnoses(
    case: str,
    check: str,
    diagnosis: str,
) -> None:
    smoke = _load_smoke()
    sdk = _sdk_interactions()
    responses = {
        "steps": sdk.Interaction(status="completed", steps=None),
        "content": sdk.Interaction(
            status="completed",
            steps=[sdk.ModelOutputStep(content=None)],
        ),
        "annotations": sdk.Interaction(
            status="completed",
            steps=[
                sdk.ModelOutputStep(
                    content=[sdk.TextContent(text="synthetic", annotations=None)]
                )
            ],
        ),
        "metadata": sdk.Interaction(
            status="completed",
            steps=[
                sdk.ModelOutputStep(
                    content=[
                        sdk.TextContent(
                            text="synthetic",
                            annotations=[sdk.FileCitation(type="file_citation")],
                        )
                    ]
                )
            ],
        ),
    }

    audit = smoke._audit_citations(
        responses[case],
        smoke.SmokeScope(
            smoke.SYNTHETIC_COURSE_ID,
            smoke.SYNTHETIC_EXAM_ID,
            smoke.SYNTHETIC_LECTURE_ID,
        ),
        "fileSearchStores/sdk-store/documents/sdk-document",
        "task-2-8-synthetic.pdf",
    )

    assert audit.checks[check] == diagnosis


@pytest.mark.parametrize(
    ("case", "diagnosis"),
    [
        ("content", "citation_content_invalid"),
        ("annotations", "citation_annotations_invalid"),
    ],
)
def test_real_sdk_malformed_citation_containers_are_not_absent(
    case: str,
    diagnosis: str,
) -> None:
    smoke = _load_smoke()
    sdk = _sdk_interactions()
    if case == "content":
        step = sdk.ModelOutputStep.model_construct(content="malformed")
    else:
        step = sdk.ModelOutputStep(
            content=[
                sdk.TextContent.model_construct(
                    type="text",
                    text="synthetic",
                    annotations="malformed",
                )
            ]
        )
    response = sdk.Interaction(status="completed", steps=[step])

    audit = smoke._audit_citations(
        response,
        smoke.SmokeScope(
            smoke.SYNTHETIC_COURSE_ID,
            smoke.SYNTHETIC_EXAM_ID,
            smoke.SYNTHETIC_LECTURE_ID,
        ),
        "fileSearchStores/sdk-store/documents/sdk-document",
        "task-2-8-synthetic.pdf",
    )

    assert audit.checks["citation_presence"] == diagnosis


@pytest.mark.parametrize(
    ("field", "check", "diagnosis"),
    [
        ("document_uri", "citation_document_binding", "citation_document_uri_absent"),
        ("file_name", "citation_file_binding", "citation_file_absent"),
        ("page_number", "citation_page_binding", "citation_page_absent"),
        ("source", "citation_excerpt_binding", "citation_excerpt_absent"),
    ],
)
def test_real_file_citation_optional_none_is_not_a_scope_mismatch(
    field: str,
    check: str,
    diagnosis: str,
) -> None:
    smoke = _load_smoke()
    values: dict[str, object] = {field: None}
    if field == "source":
        values.update(start_index=None, end_index=None)
    response = _real_interaction(smoke, citation=_real_citation(smoke, **values))

    audit = smoke._audit_citations(
        response,
        smoke.SmokeScope(
            smoke.SYNTHETIC_COURSE_ID,
            smoke.SYNTHETIC_EXAM_ID,
            smoke.SYNTHETIC_LECTURE_ID,
        ),
        "fileSearchStores/sdk-store/documents/sdk-document",
        "task-2-8-synthetic.pdf",
    )

    assert audit.checks[check] == diagnosis
    if field == "document_uri":
        assert len(audit.citations) == 1


def test_real_sdk_usage_uses_current_names_and_preserves_zero() -> None:
    smoke = _load_smoke()
    sdk = _sdk_interactions()

    audit = smoke._audit_usage(sdk.Usage(total_input_tokens=0, total_output_tokens=0))

    assert audit.input_tokens == 0
    assert audit.output_tokens == 0
    assert audit.checks == {"usage_input": "passed", "usage_output": "passed"}

    missing = smoke._audit_usage(sdk.Usage())
    assert missing.checks == {
        "usage_input": "usage_input_absent",
        "usage_output": "usage_output_absent",
    }

    malformed = smoke._audit_usage(
        sdk.Usage.model_construct(
            total_input_tokens="legacy",
            total_output_tokens=-1,
        )
    )
    assert malformed.checks == {
        "usage_input": "usage_input_invalid",
        "usage_output": "usage_output_invalid",
    }


def test_real_sdk_optional_values_distinguish_absent_malformed_and_mismatch() -> None:
    smoke = _load_smoke()
    sdk = _sdk_interactions()
    scope = smoke.SmokeScope(
        smoke.SYNTHETIC_COURSE_ID,
        smoke.SYNTHETIC_EXAM_ID,
        smoke.SYNTHETIC_LECTURE_ID,
    )

    malformed = _real_citation(smoke)
    malformed = sdk.FileCitation.model_construct(
        **{
            **malformed.model_dump(),
            "custom_metadata": "malformed",
        }
    )
    malformed_audit = smoke._audit_citations(
        _real_interaction(smoke, citation=malformed),
        scope,
        "fileSearchStores/sdk-store/documents/sdk-document",
        "task-2-8-synthetic.pdf",
    )
    assert (
        malformed_audit.checks["citation_scope_binding"]
        == "citation_metadata_invalid"
    )

    mismatched = _real_citation(
        smoke,
        custom_metadata={
            "course_id": smoke.SYNTHETIC_COURSE_ID,
            "exam_id": smoke.SYNTHETIC_EXAM_ID,
            "lecture_id": smoke.WRONG_LECTURE_ID,
            "source_revision_id": smoke.SYNTHETIC_REVISION_ID,
        },
    )
    mismatch_audit = smoke._audit_citations(
        _real_interaction(smoke, citation=mismatched),
        scope,
        "fileSearchStores/sdk-store/documents/sdk-document",
        "task-2-8-synthetic.pdf",
    )
    assert (
        mismatch_audit.checks["citation_scope_binding"]
        == "citation_scope_mismatch"
    )

    with pytest.raises(smoke.SmokeContractError) as absent:
        smoke._interaction_output(
            sdk.Interaction.model_construct(
                status="completed",
                output_text=None,
            )
        )
    assert absent.value.reason == "structured_output_absent"

    malformed_output = sdk.Interaction.model_construct(
        status="completed",
        output_text=7,
    )
    with pytest.raises(smoke.SmokeContractError) as invalid:
        smoke._interaction_output(malformed_output)
    assert invalid.value.reason == "structured_output_invalid"


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("course_id", "private-course"),
        ("exam_id", "private-exam"),
        ("lecture_id", "Lecture-13"),
        ("source_revision_id", "sr_private"),
        ("fixture_sha256", "0" * 64),
    ],
)
def test_synthetic_diagnostic_mismatch_blocks_before_secret_read(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    field: str,
    value: str,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets("must-not-be-read")
    request = dataclasses.replace(
        smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json"),
        **{field: value},
    )
    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic scope mismatch"):
        asyncio.run(
            smoke.run_authorized_live_smoke(
                secret_store=secrets,
                diagnostic_request=request,
            )
        )

    assert secrets.calls == []
    assert not request.output_path.exists()


def test_synthetic_diagnostic_sink_is_atomic_private_redacted_and_deletable(
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    sink = smoke._SyntheticDiagnosticSink.open(request)
    sink.add_secret("synthetic-secret-value")
    sink.capture(
        "provider_response",
        {
            "status": 400,
            "message": "full synthetic provider message",
            "body": {"synthetic": "full body"},
            "binary_body": b"prefix-synthetic-secret-value-suffix",
            "api_key": "synthetic-secret-value",
            "headers": {
                "Authorization": "Bearer synthetic-secret-value",
                "Cookie": "session=synthetic-secret-value",
                "X-Synthetic": "retained",
            },
        },
    )
    try:
        raise RuntimeError("full synthetic exception message")
    except RuntimeError as error:
        sink.capture_exception("provider_exception", error)
    sink.close()

    if os.name != "nt":
        assert stat.S_IMODE(request.output_path.stat().st_mode) == 0o600
    payload = json.loads(request.output_path.read_text(encoding="utf-8"))
    encoded = json.dumps(payload, sort_keys=True)
    assert "full synthetic provider message" in encoded
    assert "full body" in encoded
    assert "full synthetic exception message" in encoded
    assert "Traceback" in encoded
    assert "retained" not in encoded
    assert "synthetic-secret-value" not in encoded
    sink.delete()
    assert not request.output_path.exists()


def test_synthetic_diagnostic_path_inside_git_blocks_before_secret_read(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets("must-not-be-read")
    request = smoke._synthetic_diagnostic_request(ROOT / "diagnostic.json")
    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")

    with pytest.raises(smoke.LiveSmokeBlocked, match="outside the repository"):
        asyncio.run(
            smoke.run_authorized_live_smoke(
                secret_store=secrets,
                diagnostic_request=request,
            )
        )

    assert secrets.calls == []


def test_synthetic_diagnostic_rejects_injected_session_before_secret_read(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets("must-not-be-read")
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")

    with pytest.raises(smoke.LiveSmokeBlocked, match="default synthetic session"):
        asyncio.run(
            smoke.run_authorized_live_smoke(
                secret_store=secrets,
                session_factory=lambda api_key: _FakeSession(smoke),
                diagnostic_request=request,
            )
        )

    assert secrets.calls == []
    assert not request.output_path.exists()


def test_synthetic_diagnostic_overflow_leaves_no_partial_file(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    monkeypatch.setattr(smoke, "_MAX_DIAGNOSTIC_BYTES", 64)
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    sink = smoke._SyntheticDiagnosticSink.open(request)
    sink.capture("provider_response", {"body": "x" * 100})

    with pytest.raises(smoke.SmokeContractError, match="diagnostic overflow"):
        sink.close()

    assert not request.output_path.exists()


def test_synthetic_diagnostic_permissions_fall_back_without_fchmod(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    if not hasattr(smoke.os, "fchmod"):
        pytest.skip("POSIX fchmod fallback is unavailable")
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    sink = smoke._SyntheticDiagnosticSink.open(request)
    chmod_calls: list[tuple[object, int]] = []
    real_chmod = smoke.os.chmod

    def chmod(path: object, mode: int) -> None:
        chmod_calls.append((path, mode))
        real_chmod(path, mode)

    monkeypatch.delattr(smoke.os, "fchmod")
    monkeypatch.setattr(smoke.os, "chmod", chmod)

    sink.capture("synthetic", {"status": "ready"})
    sink.close()

    assert len(chmod_calls) == 1
    assert Path(chmod_calls[0][0]).parent == request.output_path.parent
    assert chmod_calls[0][1] == 0o600
    assert request.output_path.exists()
    sink.delete()


def _install_fake_windows_acl(
    monkeypatch: pytest.MonkeyPatch,
    smoke: ModuleType,
    *,
    extra_ace: bool = False,
) -> list[str]:
    full_access = 0x1F01FF
    sid = "S-1-5-21-1000"
    calls: list[str] = []

    class Token:
        def Close(self) -> None:
            calls.append("token.close")

    class Acl:
        def __init__(self) -> None:
            self.aces: list[tuple[tuple[int, int], int, str]] = []

        def AddAccessAllowedAce(self, revision: int, mask: int, trustee: str) -> None:
            assert revision == 2
            self.aces.append(((0, 0), mask, trustee))

        def GetAceCount(self) -> int:
            return len(self.aces) + int(extra_ace)

        def GetAce(self, index: int) -> tuple[tuple[int, int], int, str]:
            if index < len(self.aces):
                return self.aces[index]
            return ((0, 0), full_access, "S-1-1-0")

    class Attributes:
        def __init__(self) -> None:
            self.dacl: Acl | None = None
            self.control = 0
            self.bInheritHandle = True

        def SetSecurityDescriptorDacl(
            self, present: bool, dacl: Acl, defaulted: bool
        ) -> None:
            assert present and not defaulted
            self.dacl = dacl

        def SetSecurityDescriptorControl(self, mask: int, value: int) -> None:
            assert mask == 0x1000
            self.control = value

    class Handle:
        def __init__(self, path: str, attributes: Attributes) -> None:
            assert attributes.dacl is not None
            self.fd = os.open(path, os.O_CREAT | os.O_EXCL | os.O_RDWR, 0o600)
            self.dacl = attributes.dacl
            self.control = attributes.control

        def Detach(self) -> int:
            descriptor = self.fd
            self.fd = -1
            calls.append("handle.detach")
            return descriptor

        def Close(self) -> None:
            if self.fd != -1:
                os.close(self.fd)
                self.fd = -1
            calls.append("handle.close")

    class Descriptor:
        def __init__(self, handle: Handle) -> None:
            self.handle = handle

        def GetSecurityDescriptorDacl(self) -> Acl:
            return self.handle.dacl

        def GetSecurityDescriptorControl(self) -> tuple[int, int]:
            return self.handle.control, 1

    def create_file(
        path: str,
        access: int,
        share: int,
        attributes: Attributes,
        creation: int,
        flags: int,
        template: object,
    ) -> Handle:
        del access, flags, template
        assert share == 0 and creation == 1 and not attributes.bInheritHandle
        calls.append("create_file")
        return Handle(path, attributes)

    modules = {
        "msvcrt": SimpleNamespace(open_osfhandle=lambda handle, flags: handle),
        "ntsecuritycon": SimpleNamespace(FILE_ALL_ACCESS=full_access),
        "win32api": SimpleNamespace(GetCurrentProcess=lambda: object(), CloseHandle=os.close),
        "win32con": SimpleNamespace(
            GENERIC_READ=1,
            GENERIC_WRITE=2,
            CREATE_NEW=1,
            FILE_ATTRIBUTE_NORMAL=0,
        ),
        "win32file": SimpleNamespace(CreateFile=create_file),
        "win32security": SimpleNamespace(
            TOKEN_QUERY=1,
            TokenUser=1,
            ACL_REVISION=2,
            SE_DACL_PROTECTED=0x1000,
            SE_FILE_OBJECT=1,
            DACL_SECURITY_INFORMATION=4,
            ACCESS_ALLOWED_ACE_TYPE=0,
            INHERITED_ACE=0x10,
            OpenProcessToken=lambda process, access: Token(),
            GetTokenInformation=lambda token, kind: (sid, 0),
            ACL=Acl,
            SECURITY_ATTRIBUTES=Attributes,
            GetSecurityInfo=lambda handle, kind, info: Descriptor(handle),
            EqualSid=lambda left, right: left == right,
        ),
    }
    monkeypatch.setattr(smoke, "_IS_WINDOWS", True)
    monkeypatch.setattr(smoke, "import_module", modules.__getitem__)
    return calls


def test_windows_synthetic_diagnostic_requires_verified_current_user_dacl(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    sink = smoke._SyntheticDiagnosticSink.open(request)
    calls = _install_fake_windows_acl(monkeypatch, smoke)

    sink.capture("synthetic", {"status": "ready"})
    sink.close()

    assert calls == ["token.close", "create_file", "handle.detach"]
    assert request.output_path.exists()
    sink.delete()


def test_windows_synthetic_diagnostic_fails_closed_without_verified_dacl(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    sink = smoke._SyntheticDiagnosticSink.open(request)

    monkeypatch.setattr(smoke, "_IS_WINDOWS", True)
    monkeypatch.setattr(
        smoke, "import_module", lambda name: (_ for _ in ()).throw(ImportError(name))
    )
    sink.capture("synthetic", {"status": "ready"})

    with pytest.raises(smoke.SmokeContractError) as raised:
        sink.close()

    assert raised.value.reason == "diagnostic_permissions_unavailable"
    assert not request.output_path.exists()


def test_windows_synthetic_diagnostic_rejects_any_additional_ace(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    request = smoke._synthetic_diagnostic_request(tmp_path / "diagnostic.json")
    sink = smoke._SyntheticDiagnosticSink.open(request)
    _install_fake_windows_acl(monkeypatch, smoke, extra_ace=True)
    sink.capture("synthetic", {"status": "ready"})

    with pytest.raises(smoke.SmokeContractError) as raised:
        sink.close()

    assert raised.value.reason == "diagnostic_permissions_unavailable"
    assert not request.output_path.exists()


def test_check_matrix_continues_after_independent_positive_failures() -> None:
    smoke = _load_smoke()
    evidence: dict[str, object] = {}

    class IndependentFailures(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                return smoke.SmokeQueryResult(
                    answer={"answer": "wrong", "supported": True},
                    citations=(),
                )
            return result

    session = IndependentFailures(smoke)
    with pytest.raises(smoke.SmokeContractError):
        asyncio.run(smoke.run_contract_smoke(session, failure_evidence=evidence))

    checks = evidence["checks"]
    assert checks["positive_answer"] == "positive_answer_missing_marker"
    assert checks["citation_presence"] == "positive_citation_missing"
    assert checks["negative_structured_output"] == "not_run"
    assert checks["wrong_lecture_filtering"] == "not_run"
    assert checks["document_listing"] == "not_run"
    assert checks["cleanup_document"] == "passed"
    assert checks["cleanup_file"] == "passed"
    assert checks["cleanup_store"] == "passed"
    names = [name for name, _ in session.calls]
    assert names.count("query") == 1
    assert names.count("delete_document") == 5
    assert names.count("delete_file") == 5
    assert names.count("delete_store") == 1


def test_optional_document_uri_diagnosis_does_not_fail_bound_citation() -> None:
    smoke = _load_smoke()

    class OptionalDocumentUri(_FakeSession):
        async def query(self, *args: object, **kwargs: object) -> object:
            result = await super().query(*args, **kwargs)
            scope = args[2]
            if scope.lecture_id == smoke.SYNTHETIC_LECTURE_ID:
                checks = {
                    name: "passed" for name in smoke._CITATION_CHECKS
                }
                checks["citation_document_binding"] = (
                    "citation_document_uri_absent"
                )
                return dataclasses.replace(
                    result,
                    citation_checks=tuple(checks.items()),
                )
            return result

    record = asyncio.run(smoke.run_contract_smoke(OptionalDocumentUri(smoke)))

    assert record["status"] == "passed"


def test_cleanup_diagnostics_distinguish_request_from_context_close() -> None:
    smoke = _load_smoke()

    class Capture:
        def __init__(self) -> None:
            self.events: list[str] = []

        def capture(self, label: str, value: object) -> None:
            del value
            self.events.append(label)

        def capture_exception(self, label: str, error: BaseException) -> None:
            del error
            self.events.append(label)

    class DeleteDocuments:
        def __init__(self, *, fail_request: bool) -> None:
            self.fail_request = fail_request

        async def delete(self, **kwargs: object) -> None:
            del kwargs
            if self.fail_request:
                raise RuntimeError("synthetic delete request failure")

    class DeleteAio:
        def __init__(self, *, fail_request: bool, fail_close: bool) -> None:
            self.file_search_stores = SimpleNamespace(
                documents=DeleteDocuments(fail_request=fail_request)
            )
            self.fail_close = fail_close

        async def aclose(self) -> None:
            if self.fail_close:
                raise RuntimeError("synthetic context close failure")

    for fail_request, fail_close, expected in (
        (True, False, "cleanup.document.delete_request_failed"),
        (False, True, "cleanup.document.context_close_failed"),
    ):
        capture = Capture()

        def sdk_factory(
            *,
            request_failure: bool = fail_request,
            close_failure: bool = fail_close,
            **kwargs: object,
        ) -> object:
            del kwargs
            return SimpleNamespace(
                aio=DeleteAio(
                    fail_request=request_failure,
                    fail_close=close_failure,
                )
            )

        session = smoke.GoogleGenaiSmokeSession(
            "synthetic-sdk-key",
            sdk_factory=sdk_factory,
            diagnostic_sink=capture,
        )
        with pytest.raises(smoke.GeminiProviderError):
            asyncio.run(session.delete_document("documents/synthetic"))
        assert expected in capture.events


def test_google_genai_serializes_exact_interactions_wire_body() -> None:
    sdk = import_module("google.genai")
    captured: list[dict[str, object]] = []

    async def handler(request: httpx.Request) -> httpx.Response:
        captured.append(json.loads(request.content))
        return httpx.Response(
            200,
            json={"status": "completed", "output_text": "{}", "steps": []},
            request=request,
        )

    async def send() -> None:
        client = sdk.Client(
            api_key="synthetic-sdk-key",
            http_options={
                "api_version": "v1beta",
                "async_client_args": {"transport": httpx.MockTransport(handler)},
            },
        )
        try:
            await client.aio.interactions.create(
                model="gemini-3.7-flash",
                input="synthetic",
                store=False,
                tools=[
                    {
                        "type": "file_search",
                        "file_search_store_names": ["fileSearchStores/synthetic"],
                        "metadata_filter": 'course_id="synthetic"',
                    }
                ],
                response_format={
                    "type": "text",
                    "mime_type": "application/json",
                    "schema": {"type": "object"},
                },
            )
        finally:
            await client.aio.aclose()

    asyncio.run(send())

    assert len(captured) == 1
    assert set(captured[0]) == {
        "input",
        "model",
        "response_format",
        "store",
        "tools",
    }
    assert "response_mime_type" not in captured[0]
    assert captured[0]["response_format"] == {
        "type": "text",
        "mime_type": "application/json",
        "schema": {"type": "object"},
    }


def test_plan_names_the_required_later_live_wiring_change() -> None:
    smoke = _load_smoke()

    plan = smoke._plan()

    assert "run_authorized_live_smoke" in plan["required_owner_action"]
    assert plan["calls_provider"] is False
    assert plan["reads_secrets"] is False


def test_private_preflight_consumes_schema29_projection_without_source_trust_repository(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    ingestion = object()
    catalog = object()
    artifacts = SimpleNamespace(repository=ingestion, catalog=catalog)
    parser = object()
    projected = object()
    captured: dict[str, object] = {}

    def project(slide_revision_id: str, **dependencies: object) -> object:
        captured["slide_revision_id"] = slide_revision_id
        captured.update(dependencies)
        return projected

    monkeypatch.setattr(smoke, "project_schema29_index_input", project, raising=False)

    result = smoke.prepare_private_shadow_index_input(
        "29",
        schema_version=29,
        artifacts=artifacts,
        materialization_root=tmp_path,
        parser=parser,
    )

    assert result is projected
    assert captured == {
        "slide_revision_id": "29",
        "schema_version": 29,
        "ingestion": ingestion,
        "catalog": catalog,
        "artifacts": artifacts,
        "materialization_root": tmp_path,
        "parser": parser,
    }


def test_private_preflight_flow_emits_only_allowlisted_schema29_evidence(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    from oms_hub.artifacts import ArtifactRole
    from oms_hub.files.atomic import sha256_file
    from oms_hub.knowledge.models import (
        EvidenceLocator,
        EvidenceLocatorKind,
        EvidenceUnit,
        SourceRevisionState,
    )
    from oms_hub.knowledge.service import CanonicalInputArtifact, IndexInputView
    from oms_hub.providers.contracts import AuthorityClass

    smoke = _load_smoke()
    pptx = tmp_path / "synthetic.pptx"
    pdf = tmp_path / "synthetic.pdf"
    markdown = tmp_path / "normalized.md"
    pptx.write_bytes(b"offline synthetic pptx")
    pdf.write_bytes(smoke.synthetic_pdf_bytes())
    markdown.write_text("# Synthetic\n\nOffline evidence.\n", encoding="utf-8")
    evidence = EvidenceUnit(
        evidence_id="ev_schema29",
        source_revision_id="sr_schema29",
        authority_class=AuthorityClass.COURSE_MATERIAL,
        course_id="course",
        exam_id="exam",
        lecture_id="lecture",
        locator=EvidenceLocator(EvidenceLocatorKind.SLIDE, "1"),
        normalized_text="Offline evidence.",
        content_sha256=hashlib.sha256(b"Offline evidence.").hexdigest(),
    )
    second_block = dataclasses.replace(
        evidence,
        evidence_id="ev_schema29_second_block",
        locator=EvidenceLocator(EvidenceLocatorKind.SLIDE, "1:2"),
    )
    next_slide = dataclasses.replace(
        evidence,
        evidence_id="ev_schema29_next_slide",
        locator=EvidenceLocator(EvidenceLocatorKind.SLIDE, "slide 2:1"),
    )
    view = IndexInputView(
        source_document_id="legacy-study-revision:29",
        source_revision_id="sr_schema29",
        source_family="legacy_slides",
        revision_state=SourceRevisionState.READY,
        authority_class=AuthorityClass.COURSE_MATERIAL,
        course_id="course",
        exam_id="exam",
        lecture_id="lecture",
        pptx=CanonicalInputArtifact(
            "29:pptx",
            ArtifactRole.PPTX,
            pptx,
            sha256_file(pptx),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        pdf=CanonicalInputArtifact(
            "29:pdf",
            ArtifactRole.PDF,
            pdf,
            sha256_file(pdf),
            "application/pdf",
        ),
        markdown=CanonicalInputArtifact(
            "sr_schema29:markdown",
            ArtifactRole.CLEANED,
            markdown,
            sha256_file(markdown),
            "text/markdown",
        ),
        evidence_units=(evidence, second_block, next_slide),
        assets=(),
    )
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)

    record = smoke.run_private_shadow_preflight(
        "29",
        schema_version=29,
        artifacts=SimpleNamespace(),
        materialization_root=tmp_path,
    )

    assert set(record) == {
        "status",
        "source_revision_hash",
        "document_types",
        "page_count",
        "slide_count",
        "provider_operation_states",
        "byte_usage",
        "warnings",
    }
    assert record["status"] == "ready"
    assert record["document_types"] == ["markdown", "pdf", "pptx"]
    assert record["page_count"] == 1
    assert record["slide_count"] == 2
    assert record["provider_operation_states"] == ["private_preflight_ready"]
    assert record["warnings"] == []
    script_source = SCRIPT.read_text(encoding="utf-8")
    assert "source_revisions" not in script_source
    assert "KnowledgeRepository" not in script_source


@pytest.mark.parametrize(
    ("locator", "expected"),
    (("1", 1), ("1:9", 1), ("slide 2", 2), ("slide 2:7", 2)),
)
def test_private_preflight_extracts_canonical_slide_number(
    locator: str,
    expected: int,
) -> None:
    smoke = _load_smoke()

    assert smoke._canonical_slide_number(locator) == expected


@pytest.mark.parametrize(
    "locator",
    (
        "",
        "0",
        "01",
        "1:0",
        "1:02",
        "slide 1 notes",
        "page 1",
        "1:2:3",
        "1-2",
        " 1",
        "1 ",
        "slide\t1",
        "slide\n1",
    ),
)
def test_private_preflight_rejects_noncanonical_slide_locator(locator: str) -> None:
    smoke = _load_smoke()

    with pytest.raises(smoke.LiveSmokeBlocked, match="invalid slide evidence"):
        smoke._canonical_slide_number(locator)


def _private_shadow_view(smoke: ModuleType, tmp_path: Path) -> object:
    from oms_hub.artifacts import ArtifactRole
    from oms_hub.document_processing.domain import DocumentLocator
    from oms_hub.files.atomic import sha256_file
    from oms_hub.knowledge.models import (
        EvidenceLocator,
        EvidenceLocatorKind,
        EvidenceUnit,
        SourceRevisionState,
    )
    from oms_hub.knowledge.service import (
        CanonicalInputArtifact,
        IndexAssetView,
        IndexInputView,
    )
    from oms_hub.providers.contracts import AuthorityClass

    revision_id = "sr_private_shadow_aaaaaaaaaaaa"
    pptx = tmp_path / "lecture.pptx"
    pdf = tmp_path / "lecture.pdf"
    markdown = tmp_path / "normalized.md"
    image = tmp_path / "diagram.png"
    pptx.write_bytes(b"private fixture pptx")
    pdf.write_bytes(smoke.synthetic_pdf_bytes())
    markdown.write_text("# Private fixture\n\nDo not emit this text.\n", encoding="utf-8")
    Image.new("RGB", (32, 24), "white").save(image, format="PNG")
    evidence = EvidenceUnit(
        evidence_id="ev_private_1",
        source_revision_id=revision_id,
        authority_class=AuthorityClass.COURSE_MATERIAL,
        course_id="course-private",
        exam_id="exam-private",
        lecture_id="lecture-private",
        locator=EvidenceLocator(EvidenceLocatorKind.SLIDE, "1"),
        normalized_text="Do not emit this text.",
        content_sha256=hashlib.sha256(b"Do not emit this text.").hexdigest(),
    )
    return IndexInputView(
        source_document_id="opaque-private-document",
        source_revision_id=revision_id,
        source_family="legacy_slides",
        revision_state=SourceRevisionState.READY,
        authority_class=AuthorityClass.COURSE_MATERIAL,
        course_id=evidence.course_id,
        exam_id=evidence.exam_id,
        lecture_id=evidence.lecture_id,
        pptx=CanonicalInputArtifact(
            f"{revision_id}:pptx",
            ArtifactRole.PPTX,
            pptx,
            sha256_file(pptx),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        pdf=CanonicalInputArtifact(
            f"{revision_id}:pdf",
            ArtifactRole.PDF,
            pdf,
            sha256_file(pdf),
            "application/pdf",
        ),
        markdown=CanonicalInputArtifact(
            f"{revision_id}:markdown",
            ArtifactRole.CLEANED,
            markdown,
            sha256_file(markdown),
            "text/markdown",
        ),
        evidence_units=(evidence,),
        assets=(
            IndexAssetView(
                asset_id="asset-private-diagram",
                path=image,
                media_type="image/png",
                sha256=sha256_file(image),
                locator=DocumentLocator("slide 1", slide_number=1),
                width=32,
                height=24,
                visual_semantic=True,
                evidence_ids=(evidence.evidence_id,),
            ),
        ),
    )


def _private_diagnostic_path(root: Path, run_id: str = "0" * 32) -> Path:
    return root / "oms-task28-runs" / run_id / "diagnostic" / "provider-diagnostic.json"


def _private_diagnostic_capability(
    monkeypatch: pytest.MonkeyPatch,
    root: Path,
    run_id: str = "0" * 32,
) -> Path:
    path = _private_diagnostic_path(root, run_id)
    path.parent.mkdir(parents=True, exist_ok=True)
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(path))
    return path


class _PrivateShadowSession:
    def __init__(
        self,
        smoke: ModuleType,
        *,
        fail_cleanup: bool = False,
        fail_positive: bool = False,
        invalid_negative: bool = False,
        uncertain_upload: bool = False,
        uncertain_store: bool = False,
        fail_reconciliation: bool = False,
        unknown_primary: bool = False,
        fail_markdown_import: bool = False,
    ) -> None:
        self.smoke = smoke
        self.fail_cleanup = fail_cleanup
        self.fail_positive = fail_positive
        self.invalid_negative = invalid_negative
        self.uncertain_upload = uncertain_upload
        self.uncertain_store = uncertain_store
        self.fail_reconciliation = fail_reconciliation
        self.unknown_primary = unknown_primary
        self.fail_markdown_import = fail_markdown_import
        self.calls: list[tuple[str, object]] = []
        self.live_files: dict[str, str] = {}
        self.live_stores: dict[str, str] = {}
        self.model_contract = smoke.PRIVATE_SHADOW_MODEL_CONTRACT

    async def create_store(self, display_name: str, embedding_model: str) -> str:
        self.calls.append(("create_store", (display_name, embedding_model)))
        store_name = "fileSearchStores/private-store"
        self.live_stores[store_name] = display_name
        if self.uncertain_store:
            self.uncertain_store = False
            raise self.smoke.SmokeContractError("uncertain store response")
        return store_name

    async def find_stores(self, display_name: str) -> tuple[str, ...]:
        self.calls.append(("find_stores", display_name))
        return tuple(
            name for name, current in self.live_stores.items() if current == display_name
        )

    async def upload_input(
        self,
        display_name: str,
        path: Path,
        media_type: str,
    ) -> str:
        self.calls.append(("upload_input", (display_name, path.name, media_type)))
        file_name = f"files/{display_name}"
        self.live_files[file_name] = display_name
        if self.uncertain_upload:
            self.uncertain_upload = False
            raise self.smoke.SmokeContractError("uncertain upload response")
        return file_name

    async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]:
        self.calls.append(("find_files", display_names))
        if self.fail_reconciliation and len(
            [call for call in self.calls if call[0] == "find_files"]
        ) == 3:
            raise RuntimeError("raw reconciliation failure must not escape")
        return tuple(
            name
            for name, display_name in self.live_files.items()
            if display_name in display_names
        )

    async def import_input(
        self,
        store_name: str,
        file_name: str,
        metadata: tuple[tuple[str, str], ...],
        chunking: object | None,
    ) -> str:
        self.calls.append(("import_input", (store_name, file_name, metadata, chunking)))
        if self.fail_markdown_import and ("input_key", "normalized_markdown") in metadata:
            raise self.smoke.GeminiProviderError(
                "redacted provider failure",
                provider_status_code=400,
                category="provider",
                diagnostic_code="unknown_provider",
            )
        return f"operations/{len(self.calls)}"

    async def wait_for_import(self, operation_name: str) -> str:
        self.calls.append(("wait_for_import", operation_name))
        return f"fileSearchStores/private-store/documents/{len(self.calls)}"

    async def query_private(
        self,
        store_name: str,
        prompt: str,
        scope: object,
        *,
        source_revision_id: str,
        manifest: object,
        file_bindings: tuple[tuple[str, str], ...],
        require_structured_no_result: bool = False,
        require_structured_supported: bool = False,
    ) -> object:
        del source_revision_id, manifest, file_bindings, require_structured_supported
        self.calls.append(("query_private", (store_name, prompt, scope)))
        if scope.lecture_id == "lecture-private":
            if self.unknown_primary:
                raise RuntimeError(
                    "raw provider identity fileSearchStores/private-store must not escape"
                )
            if self.fail_positive:
                raise self.smoke.SmokeContractError(
                    "primary private query failed",
                    reason="private_citation_unresolved",
                )
            assert require_structured_no_result is False
            return self.smoke.PrivateShadowQueryAudit(2, 2, 17, 9, None, None)
        assert require_structured_no_result is True
        return self.smoke.PrivateShadowQueryAudit(
            0,
            0,
            11,
            4,
            True if self.invalid_negative else False,
            False if self.invalid_negative else True,
        )

    async def delete_document(self, document_name: str) -> None:
        self.calls.append(("delete_document", document_name))
        if self.fail_cleanup and len(
            [call for call in self.calls if call[0] == "delete_document"]
        ) == 1:
            raise RuntimeError("raw cleanup failure must not escape")

    async def delete_file(self, file_name: str) -> None:
        self.calls.append(("delete_file", file_name))
        self.live_files.pop(file_name, None)

    async def delete_store(self, store_name: str) -> None:
        self.calls.append(("delete_store", store_name))
        self.live_stores.pop(store_name, None)


def test_private_shadow_requires_opt_in_before_projection_or_secret(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets("stored-private-key")
    projected = False

    def project(*args: object, **kwargs: object) -> object:
        del args, kwargs
        nonlocal projected
        projected = True
        return object()

    monkeypatch.delenv("RUN_PRIVATE_GEMINI_SHADOW", raising=False)
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", project)

    with pytest.raises(smoke.LiveSmokeBlocked, match="RUN_PRIVATE_GEMINI_SHADOW"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert projected is False
    assert secrets.calls == []


def test_private_shadow_requires_terminal_diagnostic_capability_before_preflight_or_secret(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    assert (
        inspect.signature(smoke.run_authorized_private_shadow)
        .parameters["diagnostic_path"]
        .default
        is inspect.Parameter.empty
    )
    diagnostic_path = _private_diagnostic_path(tmp_path)
    diagnostic_path.parent.mkdir(parents=True)
    projected = False
    secrets = _FakeSecrets("must-not-be-read")

    def project(*args: object, **kwargs: object) -> object:
        del args, kwargs
        nonlocal projected
        projected = True
        return object()

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.delenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", raising=False)
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", project)

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=diagnostic_path,
            )
        )

    assert projected is False
    assert secrets.calls == []


@pytest.mark.parametrize("invalid_path", (None, "not-a-path", object()))
def test_private_shadow_rejects_non_path_diagnostic_capability_before_preflight_or_secret(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    invalid_path: object,
) -> None:
    smoke = _load_smoke()
    supplied = _private_diagnostic_capability(monkeypatch, tmp_path)
    projected = False
    secrets = _FakeSecrets("must-not-be-read")

    def project(*args: object, **kwargs: object) -> object:
        del args, kwargs
        nonlocal projected
        projected = True
        return object()

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", project)

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=invalid_path,
            )
        )

    assert supplied.parent.is_dir()
    assert projected is False
    assert secrets.calls == []


def test_private_shadow_mismatch_fails_before_secret_and_provider(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    secrets = _FakeSecrets("stored-private-key")
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)
    approved["page_count"] = 2

    with pytest.raises(smoke.LiveSmokeBlocked, match="preflight mismatch"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=secrets,
                session_factory=lambda key: pytest.fail(f"provider received {key}"),
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert secrets.calls == []


def test_private_diagnostic_lexical_alias_blocks_before_preflight_or_secret(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    canonical = _private_diagnostic_path(tmp_path)
    canonical.parent.mkdir(parents=True)
    alias = canonical.parent / ".." / "diagnostic" / canonical.name
    projected = False
    secrets = _FakeSecrets("must-not-be-read")

    def project(*args: object, **kwargs: object) -> object:
        del args, kwargs
        nonlocal projected
        projected = True
        return object()

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(alias))
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", project)

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=alias,
            )
        )

    assert projected is False
    assert secrets.calls == []


def test_private_diagnostic_env_alias_must_lexically_match_argument(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    canonical = _private_diagnostic_path(tmp_path)
    canonical.parent.mkdir(parents=True)
    alias = canonical.parent / ".." / "diagnostic" / canonical.name
    secrets = _FakeSecrets("must-not-be-read")
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(alias))

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=canonical,
            )
        )

    assert secrets.calls == []


def test_private_shadow_indexes_every_input_queries_and_returns_only_aggregates(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    secrets = _FakeSecrets("stored-private-key")
    session = _PrivateShadowSession(smoke)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    record = asyncio.run(
        smoke.run_authorized_private_shadow(
            "29",
            schema_version=29,
            artifacts=SimpleNamespace(),
            materialization_root=tmp_path,
            approved_preflight=approved,
            secret_store=secrets,
            session_factory=lambda key: session if key == "stored-private-key" else None,
            clock=iter((100.0, 100.25)).__next__,
            diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
        )
    )

    assert set(record) == {
        "status",
        "source_revision_hash",
        "document_types",
        "page_count",
        "slide_count",
        "provider_operation_states",
        "citation_resolution_rate",
        "duration_ms",
        "byte_usage",
        "token_usage",
        "transient_attempts",
        "failure_class",
        "warnings",
    }
    assert record["status"] == "passed"
    assert record["document_types"] == ["image", "markdown", "pdf", "pptx"]
    assert record["citation_resolution_rate"] == 1.0
    assert record["duration_ms"] == 250
    assert record["token_usage"] == {"input": 28, "output": 13}
    assert record["warnings"] == []
    assert secrets.calls == ["gemini-api-key"]
    assert len([call for call in session.calls if call[0] == "upload_input"]) == 4
    assert len([call for call in session.calls if call[0] == "import_input"]) == 4
    assert len([call for call in session.calls if call[0] == "query_private"]) == 2
    negative_prompt = [call for call in session.calls if call[0] == "query_private"][1][1][1]
    assert negative_prompt == (
        "Use only files matching the requested lecture scope. If none match, return an "
        'empty answer and supported=false.'
    )
    assert len([call for call in session.calls if call[0] == "delete_document"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_file"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_store"]) == 1
    assert session.live_stores == {}
    serialized = json.dumps(record, sort_keys=True)
    assert "Do not emit this text" not in serialized
    assert "fileSearchStores/private-store" not in serialized


def test_private_shadow_cleanup_failure_still_attempts_every_delete(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, fail_cleanup=True)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError, match="cleanup failed") as raised:
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert raised.value.reason == "private_cleanup_failed"
    assert len([call for call in session.calls if call[0] == "delete_document"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_file"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_store"]) == 1


def test_private_shadow_rejects_uncited_but_supported_wrong_scope_answer(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, invalid_negative=True)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert raised.value.reason == "private_wrong_scope_retrieved"


def test_private_shadow_primary_failure_precedes_cleanup_failure(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, fail_cleanup=True, fail_positive=True)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert raised.value.reason == "private_citation_unresolved"
    assert len([call for call in session.calls if call[0] == "delete_document"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_file"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_store"]) == 1


def test_private_shadow_primary_failure_retains_safe_cleanup_evidence(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, fail_cleanup=True, fail_positive=True)
    evidence: dict[str, object] = {}
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                failure_evidence=evidence,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert raised.value.reason == "private_citation_unresolved"
    assert set(evidence) == {
        "status",
        "source_revision_hash",
        "document_types",
        "page_count",
        "slide_count",
        "provider_operation_states",
        "byte_usage",
        "transient_attempts",
        "failure_class",
        "failure_stage",
        "failure_input_identity",
        "provider_error_category",
        "provider_status_code",
        "provider_reason",
        "provider_cleanup_outcome",
        "provider_reconciliation_outcome",
        "diagnostic_sha256",
        "warnings",
    }
    assert evidence["status"] == "blocked"
    assert evidence["failure_stage"] == "positive_query"
    assert evidence["failure_input_identity"] == "none"
    assert evidence["provider_error_category"] == "none"
    assert evidence["provider_status_code"] is None
    assert evidence["provider_reason"] == "none"
    assert evidence["provider_cleanup_outcome"] == "failed"
    assert evidence["provider_reconciliation_outcome"] == "unknown"
    assert re.fullmatch(r"[0-9a-f]{64}", str(evidence["diagnostic_sha256"]))
    assert evidence["warnings"] == [
        "private_citation_unresolved",
        "private_cleanup_failed",
    ]
    states = evidence["provider_operation_states"]
    assert isinstance(states, list)
    assert states[-1] == "private_shadow_failed"
    assert "file_reconciliation_empty" in states
    assert "store_reconciliation_empty" in states
    assert len([call for call in session.calls if call[0] == "delete_document"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_file"]) == 4
    assert len([call for call in session.calls if call[0] == "delete_store"]) == 1
    serialized = json.dumps(evidence, sort_keys=True)
    assert "primary private query failed" not in serialized
    assert "raw cleanup failure" not in serialized
    assert "fileSearchStores/private-store" not in serialized
    assert "stored-private-key" not in serialized


def test_private_shadow_markdown_import_failure_retains_only_safe_diagnostics(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, fail_markdown_import=True)
    evidence: dict[str, object] = {}
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                failure_evidence=evidence,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert evidence["failure_stage"] == "import_input"
    assert evidence["failure_input_identity"] == "normalized_markdown"
    assert evidence["provider_error_category"] == "provider"
    assert evidence["provider_status_code"] == 400
    assert evidence["provider_reason"] == "provider_bad_request"
    assert evidence["provider_cleanup_outcome"] == "complete"
    assert evidence["provider_reconciliation_outcome"] == "empty"
    serialized = json.dumps(evidence, sort_keys=True)
    assert "redacted provider failure" not in serialized
    assert "stored-private-key" not in serialized
    assert "fileSearchStores/" not in serialized


def test_private_shadow_terminal_failure_writes_one_safe_diagnostic(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    evidence: dict[str, object] = {}
    diagnostic_path = _private_diagnostic_path(tmp_path)
    diagnostic_path.parent.mkdir(parents=True)

    class UnsafeMimeSession(_PrivateShadowSession):
        async def import_input(
            self,
            store_name: str,
            file_name: str,
            metadata: tuple[tuple[str, str], ...],
            chunking: object | None,
        ) -> str:
            if ("input_key", "normalized_markdown") not in metadata:
                return await super().import_input(store_name, file_name, metadata, chunking)

            class SdkError(Exception):
                status = "INVALID_ARGUMENT"
                message = (
                    "Unsupported MIME type: fake-secret Authorization: Bearer fake-secret "
                    "C:\\private\\Lecture-13\\normalized.md"
                )
                response = SimpleNamespace(
                    status_code=400,
                    headers={"Authorization": "Bearer fake-secret"},
                    text="fake-secret private provider body",
                )

            raise smoke.translate_gemini_error(SdkError())

    session = UnsafeMimeSession(smoke)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(diagnostic_path))
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                failure_evidence=evidence,
                diagnostic_path=diagnostic_path,
            )
        )

    assert evidence["provider_reason"] == "unsupported_mime_type"
    assert evidence["diagnostic_sha256"] == hashlib.sha256(diagnostic_path.read_bytes()).hexdigest()
    assert len(diagnostic_path.read_bytes()) <= 16 * 1024
    committed = json.dumps(evidence, sort_keys=True)
    local_record = json.loads(diagnostic_path.read_text(encoding="utf-8"))
    assert set(local_record) == {
        "exception_kind",
        "failure_input_identity",
        "failure_stage",
        "provider_message",
        "provider_reason",
        "provider_status_code",
        "schema_version",
    }
    assert local_record == {
        "schema_version": 1,
        "exception_kind": "gemini_provider_error",
        "provider_status_code": 400,
        "provider_reason": "unsupported_mime_type",
        "provider_message": "Unsupported MIME type.",
        "failure_stage": "import_input",
        "failure_input_identity": "normalized_markdown",
    }
    local = json.dumps(local_record, sort_keys=True)
    for secret in ("fake-secret", "Authorization", "Lecture-13", "private provider body"):
        assert secret not in committed
        assert secret not in local


def test_private_shadow_prior_check_failure_retains_only_safe_terminal_diagnostic(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    diagnostic_path = _private_diagnostic_capability(monkeypatch, tmp_path)
    evidence: dict[str, object] = {}

    class PriorCheckFailureSession(_PrivateShadowSession):
        async def find_stores(self, display_name: str) -> tuple[str, ...]:
            del display_name
            raise smoke.GeminiProviderError(
                "fake-secret Authorization: Bearer fake-secret C:\\private\\Lecture-13",
                provider_status_code=400,
                diagnostic_code="unsupported_mime_type",
            )

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: PriorCheckFailureSession(smoke),
                failure_evidence=evidence,
                diagnostic_path=diagnostic_path,
            )
        )

    assert evidence["failure_stage"] == "prior_state_check"
    assert evidence["provider_reconciliation_outcome"] == "unknown"
    assert evidence["diagnostic_sha256"] == hashlib.sha256(diagnostic_path.read_bytes()).hexdigest()
    local = diagnostic_path.read_text(encoding="utf-8")
    assert json.loads(local)["failure_stage"] == "prior_state_check"
    for secret in ("fake-secret", "Authorization", "Lecture-13", "stored-private-key"):
        assert secret not in json.dumps(evidence, sort_keys=True)
        assert secret not in local


def test_private_shadow_prior_state_mismatch_retains_terminal_diagnostic(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    diagnostic_path = _private_diagnostic_capability(monkeypatch, tmp_path)
    evidence: dict[str, object] = {}

    class PriorStateSession(_PrivateShadowSession):
        async def find_stores(self, display_name: str) -> tuple[str, ...]:
            del display_name
            return ("fileSearchStores/private-store",)

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.LiveSmokeBlocked, match="prior operator state mismatch"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: PriorStateSession(smoke),
                failure_evidence=evidence,
                diagnostic_path=diagnostic_path,
            )
        )

    assert evidence["provider_reconciliation_outcome"] == "not_empty"
    assert evidence["diagnostic_sha256"] == hashlib.sha256(diagnostic_path.read_bytes()).hexdigest()
    assert json.loads(diagnostic_path.read_text(encoding="utf-8")) == {
        "schema_version": 1,
        "exception_kind": "unclassified_error",
        "provider_status_code": None,
        "provider_reason": "none",
        "provider_message": "Provider failure classification was unavailable.",
        "failure_stage": "prior_state_check",
        "failure_input_identity": "none",
    }


def test_private_terminal_finalization_revalidates_swapped_diagnostic_parent(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    diagnostic_path = _private_diagnostic_capability(monkeypatch, tmp_path)
    external = tmp_path / "external"
    external.mkdir()

    class SwappedDiagnosticSession(_PrivateShadowSession):
        async def import_input(
            self,
            store_name: str,
            file_name: str,
            metadata: tuple[tuple[str, str], ...],
            chunking: object | None,
        ) -> str:
            del store_name, file_name, metadata, chunking
            diagnostic_path.parent.rmdir()
            diagnostic_path.parent.symlink_to(external, target_is_directory=True)
            raise smoke.GeminiProviderError(
                "safe failure", provider_status_code=400, diagnostic_code="unsupported_mime_type"
            )

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: SwappedDiagnosticSession(smoke),
                diagnostic_path=diagnostic_path,
            )
        )

    assert not external.exists() or not list(external.iterdir())
    assert not list(external.glob(".provider-diagnostic.json.*.tmp"))


def test_private_terminal_finalization_revalidates_after_exception_property_swap(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    diagnostic_path = _private_diagnostic_capability(monkeypatch, tmp_path)
    external = tmp_path / "external"
    external.mkdir()

    class SwappingProviderError(smoke.GeminiProviderError):
        def __init__(self) -> None:
            super().__init__(
                "safe failure", provider_status_code=400, diagnostic_code="unsupported_mime_type"
            )
            self.swapped = False

        def __getattribute__(self, name: str) -> object:
            if name == "provider_status_code" and not object.__getattribute__(self, "swapped"):
                object.__setattr__(self, "swapped", True)
                diagnostic_path.parent.rmdir()
                diagnostic_path.parent.symlink_to(external, target_is_directory=True)
            return super().__getattribute__(name)

    class GetterSwapSession(_PrivateShadowSession):
        async def import_input(
            self,
            store_name: str,
            file_name: str,
            metadata: tuple[tuple[str, str], ...],
            chunking: object | None,
        ) -> str:
            del store_name, file_name, metadata, chunking
            raise SwappingProviderError()

    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: GetterSwapSession(smoke),
                diagnostic_path=diagnostic_path,
            )
        )

    assert not external.exists() or not list(external.iterdir())
    assert not list(external.glob(".provider-diagnostic.json.*.tmp"))


def test_private_terminal_diagnostic_rejects_preexisting_and_overflow_without_partial(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    diagnostic_path = _private_diagnostic_path(tmp_path)
    diagnostic_path.parent.mkdir(parents=True)
    error = smoke.GeminiProviderError(
        "safe", provider_status_code=400, diagnostic_code="unsupported_mime_type"
    )
    diagnostic_path.write_text("preexisting", encoding="utf-8")

    with pytest.raises(smoke.LiveSmokeBlocked, match="already exists"):
        smoke._write_private_terminal_diagnostic(
            diagnostic_path,
            error,
            failure_stage="import_input",
            input_identity="normalized_markdown",
        )

    diagnostic_path.unlink()
    monkeypatch.setattr(smoke, "_MAX_PRIVATE_DIAGNOSTIC_BYTES", 1)
    with pytest.raises(smoke.SmokeContractError, match="overflow"):
        smoke._write_private_terminal_diagnostic(
            diagnostic_path,
            error,
            failure_stage="import_input",
            input_identity="normalized_markdown",
        )

    assert not diagnostic_path.exists()
    assert not list(diagnostic_path.parent.glob(".provider-diagnostic.json.*.tmp"))


def test_private_terminal_diagnostic_is_one_write_and_success_leaves_none(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    diagnostic_path = _private_diagnostic_path(tmp_path)
    diagnostic_path.parent.mkdir(parents=True)
    error = smoke.GeminiProviderError(
        "safe", provider_status_code=400, diagnostic_code="unsupported_mime_type"
    )
    digest = smoke._write_private_terminal_diagnostic(
        diagnostic_path,
        error,
        failure_stage="import_input",
        input_identity="normalized_markdown",
    )

    with pytest.raises(smoke.LiveSmokeBlocked, match="already exists"):
        smoke._write_private_terminal_diagnostic(
            diagnostic_path,
            error,
            failure_stage="import_input",
            input_identity="normalized_markdown",
        )
    assert digest == hashlib.sha256(diagnostic_path.read_bytes()).hexdigest()

    success_path = _private_diagnostic_path(tmp_path, "1" * 32)
    success_path.parent.mkdir(parents=True)
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(success_path))
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    asyncio.run(
        smoke.run_authorized_private_shadow(
            "29",
            schema_version=29,
            artifacts=SimpleNamespace(),
            materialization_root=tmp_path,
            approved_preflight=approved,
            secret_store=_FakeSecrets("stored-private-key"),
            session_factory=lambda key: session,
            diagnostic_path=success_path,
        )
    )

    assert not success_path.exists()


def test_private_diagnostic_capability_mismatch_blocks_before_secret_read(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    supplied = _private_diagnostic_path(tmp_path)
    supplied.parent.mkdir(parents=True)
    launcher_path = _private_diagnostic_path(tmp_path, "1" * 32)
    secrets = _FakeSecrets("must-not-be-read")
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(launcher_path))

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=supplied,
            )
        )

    assert secrets.calls == []


def test_private_diagnostic_relative_capability_blocks_before_secret_read(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    secrets = _FakeSecrets("must-not-be-read")
    relative = Path("oms-task28-runs") / ("0" * 32) / "diagnostic" / "provider-diagnostic.json"
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(relative))

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=relative,
            )
        )

    assert secrets.calls == []


@pytest.mark.parametrize("inside_repo", (True, False))
def test_private_diagnostic_unsafe_path_blocks_before_secret_read(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    inside_repo: bool,
) -> None:
    smoke = _load_smoke()
    if inside_repo:
        supplied = _private_diagnostic_path(ROOT)
    else:
        real_root = tmp_path / "real"
        _private_diagnostic_path(real_root).parent.mkdir(parents=True)
        linked_root = tmp_path / "linked"
        linked_root.symlink_to(real_root, target_is_directory=True)
        supplied = _private_diagnostic_path(linked_root)
    secrets = _FakeSecrets("must-not-be-read")
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(supplied))

    with pytest.raises(smoke.LiveSmokeBlocked, match="diagnostic capability"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight={},
                secret_store=secrets,
                diagnostic_path=supplied,
            )
        )

    assert secrets.calls == []


def test_private_success_after_transient_retries_leaves_no_diagnostic(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    diagnostic_path = _private_diagnostic_path(tmp_path)
    diagnostic_path.parent.mkdir(parents=True)
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke)
    session.transient_attempts = 2
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH", str(diagnostic_path))
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    record = asyncio.run(
        smoke.run_authorized_private_shadow(
            "29",
            schema_version=29,
            artifacts=SimpleNamespace(),
            materialization_root=tmp_path,
            approved_preflight=approved,
            secret_store=_FakeSecrets("stored-private-key"),
            session_factory=lambda key: session,
            diagnostic_path=diagnostic_path,
        )
    )

    assert record["transient_attempts"] == 2
    assert not diagnostic_path.exists()


def test_private_shadow_generic_bad_request_is_preserved_as_safe_diagnostic(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, fail_markdown_import=True)
    evidence: dict[str, object] = {}
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.GeminiProviderError):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                failure_evidence=evidence,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert evidence["provider_error_category"] == "provider"
    assert evidence["provider_status_code"] == 400
    assert evidence["provider_reason"] == "provider_bad_request"


@pytest.mark.parametrize(
    ("session_kwargs", "cleanup", "reconciliation", "warning"),
    (
        ({"fail_positive": True}, "complete", "empty", "private_citation_unresolved"),
        (
            {"fail_positive": True, "fail_reconciliation": True},
            "unknown",
            "unknown",
            "private_citation_unresolved",
        ),
        ({"unknown_primary": True}, "complete", "empty", "private_shadow_failed"),
    ),
)
def test_private_shadow_failure_evidence_is_conservative_and_redacted(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    session_kwargs: dict[str, bool],
    cleanup: str,
    reconciliation: str,
    warning: str,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, **session_kwargs)
    evidence: dict[str, object] = {}
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)
    expected_error = (
        RuntimeError if session_kwargs.get("unknown_primary") else smoke.SmokeContractError
    )

    with pytest.raises(expected_error):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                failure_evidence=evidence,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert evidence["failure_stage"] == "positive_query"
    assert evidence["provider_cleanup_outcome"] == cleanup
    assert evidence["provider_reconciliation_outcome"] == reconciliation
    assert evidence["warnings"][0] == warning
    serialized = json.dumps(evidence, sort_keys=True)
    assert "raw provider identity" not in serialized
    assert "raw reconciliation failure" not in serialized
    assert "fileSearchStores/private-store" not in serialized
    assert "stored-private-key" not in serialized


def test_private_shadow_cleanup_only_failure_retains_failed_outcome(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, fail_cleanup=True)
    evidence: dict[str, object] = {}
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError) as raised:
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                failure_evidence=evidence,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert raised.value.reason == "private_cleanup_failed"
    assert evidence["failure_stage"] == "cleanup"
    assert evidence["provider_cleanup_outcome"] == "failed"
    assert evidence["provider_reconciliation_outcome"] == "unknown"
    assert evidence["warnings"] == ["private_cleanup_failed"]


def test_private_shadow_reconciles_uncertain_upload_before_primary_failure(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, uncertain_upload=True)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError, match="uncertain upload response"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert len([call for call in session.calls if call[0] == "delete_file"]) == 1
    assert session.live_files == {}
    assert len([call for call in session.calls if call[0] == "find_files"]) >= 3


def test_private_shadow_reconciles_uncertain_store_before_primary_failure(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    smoke = _load_smoke()
    view = _private_shadow_view(smoke, tmp_path)
    session = _PrivateShadowSession(smoke, uncertain_store=True)
    monkeypatch.setenv("RUN_PRIVATE_GEMINI_SHADOW", "1")
    monkeypatch.setattr(smoke, "prepare_private_shadow_index_input", lambda *a, **k: view)
    approved = smoke._private_shadow_preflight_from_view(view)

    with pytest.raises(smoke.SmokeContractError, match="uncertain store response"):
        asyncio.run(
            smoke.run_authorized_private_shadow(
                "29",
                schema_version=29,
                artifacts=SimpleNamespace(),
                materialization_root=tmp_path,
                approved_preflight=approved,
                secret_store=_FakeSecrets("stored-private-key"),
                session_factory=lambda key: session,
                diagnostic_path=_private_diagnostic_capability(monkeypatch, tmp_path),
            )
        )

    assert len([call for call in session.calls if call[0] == "delete_store"]) == 1
    assert session.live_stores == {}
    assert len([call for call in session.calls if call[0] == "find_stores"]) >= 3


def test_opt_in_without_stored_credential_fails_before_provider(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    smoke = _load_smoke()
    monkeypatch.setenv("RUN_LIVE_GEMINI_TESTS", "1")

    with pytest.raises(smoke.LiveSmokeBlocked) as raised:
        asyncio.run(smoke.run_authorized_live_smoke())

    assert "stored Gemini credential is unavailable" in str(raised.value)


def test_synthetic_pdf_is_deterministic_and_contains_no_private_source() -> None:
    smoke = _load_smoke()

    first = smoke.synthetic_pdf_bytes()
    second = smoke.synthetic_pdf_bytes()

    assert first == second
    assert first.startswith(b"%PDF-")
    assert b"Lecture 13" not in first


def test_gate_2b_record_is_explicitly_blocked_until_live_acceptance() -> None:
    record = json.loads(GATE_RECORD.read_text(encoding="utf-8"))

    assert record["schema_version"] == 1
    assert record["gate"] == "gate_2b_gemini_provider"
    assert record["state"] == "blocked"
    assert record["claim"] == "gate_not_open"
    assert record["assessed_input"]["sha"] == "37c1b6ac572ecdc4549f4688ab5868c4a9a85240"
    assert {item["id"] for item in record["blockers"]} == {
        "live_provider_authorization",
        "official_sdk_dependency_and_adapter",
        "private_lecture_shadow_authorization",
        "live_acceptance_and_independent_review",
    }
    assert all(item["status"] == "blocked" for item in record["live_verification"])
    assert record["review"]["gate_open_approval"] == "not_granted"
    assert "provider_document_id" not in json.dumps(record)


@pytest.mark.skipif(not LIVE_ENABLED, reason="live Gemini contract tests are opt-in")
def test_live_gemini_file_search_contract() -> None:
    smoke = _load_smoke()
    asyncio.run(smoke.run_authorized_live_smoke())
