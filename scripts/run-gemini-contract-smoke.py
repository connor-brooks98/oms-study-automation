#!/usr/bin/env python3
# ruff: noqa: E501
"""Offline-tested orchestration for the explicitly authorized Task 2.8 live smoke."""

from __future__ import annotations

import argparse
import asyncio
import base64
import csv
import hashlib
import importlib.metadata
import json
import os
import re
import stat
import subprocess
import tempfile
import traceback
from collections.abc import AsyncIterable, Awaitable, Callable, Iterable, Mapping
from dataclasses import dataclass, field
from datetime import UTC, datetime
from functools import partial
from importlib import import_module
from io import BytesIO
from pathlib import Path
from time import monotonic
from typing import TYPE_CHECKING, Any, Literal, Protocol
from uuid import uuid4
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from pydantic import BaseModel, ConfigDict, SecretStr, ValidationError
from reportlab.lib.pagesizes import letter  # type: ignore[import-untyped]
from reportlab.pdfgen.canvas import Canvas  # type: ignore[import-untyped]

from oms_hub.providers.gemini.client import (
    GeminiClientFactory,
    SdkFactory,
    translate_gemini_error,
)
from oms_hub.providers.gemini.errors import GeminiProviderError, GeminiTransientError
from oms_hub.providers.gemini.evidence import (
    failure_record as private_shadow_failure_record,
)
from oms_hub.providers.gemini.evidence import (
    validate_private_shadow_record,
)
from oms_hub.providers.gemini.file_search import build_import_file_config
from oms_hub.providers.gemini.models import GeminiConfig
from oms_hub.source_trust_schema29 import project_schema29_index_input

if TYPE_CHECKING:
    from oms_hub.artifacts import ArtifactService
    from oms_hub.indexing.service import IndexResult
    from oms_hub.knowledge.service import IndexInputView
    from oms_hub.security.secret_store import SecretStore

SYNTHETIC_COURSE_ID = "task-2-8-synthetic-course"
SYNTHETIC_EXAM_ID = "task-2-8-synthetic-exam"
SYNTHETIC_LECTURE_ID = "task-2-8-synthetic-lecture"
SYNTHETIC_REVISION_ID = "sr_aaaaaaaaaaaaaaaaaaaaaaaaaa"
SYNTHETIC_MARKER = "cobalt-otter-28"
SYNTHETIC_FACT = f"The Task 2.8 synthetic marker is {SYNTHETIC_MARKER}."
WRONG_LECTURE_ID = "task-2-8-wrong-lecture"
PRIVATE_SHADOW_MODEL_CONTRACT = (
    "2.14.0",
    "gemini-3.7-flash",
    "models/gemini-embedding-2",
    "v1beta",
)
PRIVATE_SHADOW_WRONG_LECTURE_ID = "task-2-8-private-wrong-lecture"
_REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
_MAX_DIAGNOSTIC_BYTES = 16 * 1024 * 1024
_MAX_PRIVATE_DIAGNOSTIC_BYTES = 16 * 1024
_MAX_TRANSIENT_ATTEMPTS = 10_000
_RETRY_DELAYS = (1.0, 2.0)
_IS_WINDOWS = os.name == "nt"
_PRIVATE_DIAGNOSTIC_FILENAME = "provider-diagnostic.json"
_PRIVATE_DIAGNOSTIC_REASONS = frozenset(
    {
        "none",
        "invalid_argument",
        "provider_bad_request",
        "sdk_contract",
        "timeout",
        "transport_error",
        "unknown_provider",
        "unsupported_mime_type",
    }
)
_PRIVATE_DIAGNOSTIC_MESSAGES = {
    "invalid_argument": "Provider rejected the argument.",
    "provider_bad_request": "Provider rejected the request.",
    "sdk_contract": "Provider SDK contract was invalid.",
    "timeout": "Provider operation timed out.",
    "transport_error": "Provider transport failed.",
    "unknown_provider": "Provider failure classification was unavailable.",
    "unsupported_mime_type": "Unsupported MIME type.",
}
_CITATION_CHECKS = (
    "citation_presence",
    "citation_document_binding",
    "citation_file_binding",
    "citation_scope_binding",
    "citation_page_binding",
    "citation_excerpt_binding",
)
_PRIVATE_SLIDE_COORDINATE = re.compile(
    r"(?:slide )?([1-9][0-9]*)(?::[1-9][0-9]*)?\Z"
)
class SmokeContractError(RuntimeError):
    _SAFE_REASONS = frozenset(
        {
            "citation_document_identity_unavailable",
            "citation_document_uri_absent",
            "citation_document_uri_invalid",
            "citation_excerpt_invalid",
            "citation_excerpt_absent",
            "citation_excerpt_mismatch",
            "citation_excerpt_unavailable",
            "citation_file_absent",
            "citation_file_invalid",
            "citation_metadata_invalid",
            "citation_metadata_absent",
            "citation_annotations_absent",
            "citation_annotations_invalid",
            "citation_content_absent",
            "citation_content_invalid",
            "citation_page_invalid",
            "citation_page_absent",
            "citation_page_mismatch",
            "citation_scope_mismatch",
            "citation_steps_absent",
            "citation_steps_invalid",
            "citation_wrong_document",
            "citation_wrong_file",
            "diagnostic_overflow",
            "diagnostic_permissions_unavailable",
            "negative_answer_invalid",
            "positive_answer_invalid",
            "positive_answer_missing_marker",
            "positive_answer_unsupported",
            "positive_citation_missing",
            "positive_citation_unresolved",
            "private_cleanup_failed",
            "private_citation_unresolved",
            "private_reconciliation_failed",
            "private_usage_invalid",
            "private_wrong_scope_retrieved",
            "structured_output_invalid",
            "structured_output_absent",
            "structured_output_unavailable",
            "usage_input_absent",
            "usage_input_invalid",
            "usage_output_absent",
            "usage_output_invalid",
            "usage_count_invalid",
        }
    )

    def __init__(self, message: str, *, reason: str | None = None) -> None:
        super().__init__(message)
        self.reason = reason if reason in self._SAFE_REASONS else None


class SmokeTemporaryFailure(RuntimeError):
    pass


class LiveSmokeBlocked(RuntimeError):
    pass


def _canonical_slide_number(locator: str) -> int:
    match = _PRIVATE_SLIDE_COORDINATE.fullmatch(locator)
    if match is None:
        raise LiveSmokeBlocked("private shadow source has invalid slide evidence")
    try:
        return int(match.group(1))
    except ValueError:
        raise LiveSmokeBlocked(
            "private shadow source has invalid slide evidence"
        ) from None


def prepare_private_shadow_index_input(
    slide_revision_id: str,
    *,
    schema_version: int,
    artifacts: ArtifactService,
    materialization_root: Path,
    parser: Any | None = None,
) -> IndexInputView:
    return project_schema29_index_input(
        slide_revision_id,
        schema_version=schema_version,
        ingestion=artifacts.repository,
        catalog=artifacts.catalog,
        artifacts=artifacts,
        materialization_root=materialization_root,
        parser=parser,
    )


def run_private_shadow_preflight(
    slide_revision_id: str,
    *,
    schema_version: int,
    artifacts: ArtifactService,
    materialization_root: Path,
    parser: Any | None = None,
) -> dict[str, object]:
    view = prepare_private_shadow_index_input(
        slide_revision_id,
        schema_version=schema_version,
        artifacts=artifacts,
        materialization_root=materialization_root,
        parser=parser,
    )
    return _private_shadow_preflight_from_view(view)


def _private_shadow_manifest(view: IndexInputView) -> Any:
    from oms_hub.indexing.service import IndexManifest, IndexManifestInput, build_index_manifest

    manifest = build_index_manifest(view)
    all_evidence_ids = tuple(ref.evidence_id for ref in manifest.evidence)
    return IndexManifest(
        source_revision_id=manifest.source_revision_id,
        authority_class=manifest.authority_class,
        inputs=(
            IndexManifestInput(
                input_key="pptx",
                input_kind="pptx",
                path=view.pptx.path,
                media_type=view.pptx.media_type,
                sha256=view.pptx.sha256,
                evidence_ids=all_evidence_ids,
            ),
            *manifest.inputs,
        ),
        evidence=manifest.evidence,
    )


def _private_shadow_preflight_from_view(view: IndexInputView) -> dict[str, object]:
    return _expected_private_shadow_preflight(view, _private_shadow_manifest(view))


def _expected_private_shadow_preflight(
    view: IndexInputView,
    manifest: Any,
) -> dict[str, object]:
    from oms_hub.files.pdf import validate_pdf
    from oms_hub.knowledge.models import EvidenceLocatorKind

    if not view.evidence_units or any(
        unit.locator.kind is not EvidenceLocatorKind.SLIDE
        for unit in view.evidence_units
    ):
        raise LiveSmokeBlocked("private shadow source has invalid slide evidence")
    slide_numbers = {
        _canonical_slide_number(unit.locator.value) for unit in view.evidence_units
    }
    return {
        "status": "ready",
        "source_revision_hash": hashlib.sha256(
            view.source_revision_id.encode("utf-8")
        ).hexdigest(),
        "document_types": sorted(
            {item.input_kind for item in manifest.inputs}
        ),
        "page_count": validate_pdf(view.pdf.path).page_count,
        "slide_count": len(slide_numbers),
        "provider_operation_states": ["private_preflight_ready"],
        "byte_usage": {
            "index_inputs": sum(item.path.stat().st_size for item in manifest.inputs)
        },
        "warnings": [],
    }


class _OperationFailure(Exception):
    def __init__(self, status_code: int | None) -> None:
        self.status_code = status_code
        super().__init__("Gemini import operation failed")


class SmokeAnswer(BaseModel):
    model_config = ConfigDict(extra="forbid")

    answer: str
    supported: bool


@dataclass(frozen=True, slots=True)
class SmokeScope:
    course_id: str
    exam_id: str
    lecture_id: str


@dataclass(frozen=True, slots=True)
class SmokeCitation:
    document_name: str
    page_number: int | None
    excerpt: str


@dataclass(frozen=True, slots=True)
class SmokeQueryResult:
    answer: dict[str, object]
    citations: tuple[SmokeCitation, ...]
    input_tokens: int | None = None
    output_tokens: int | None = None
    citation_checks: tuple[tuple[str, str], ...] = ()
    usage_checks: tuple[tuple[str, str], ...] = ()


@dataclass(frozen=True, slots=True)
class CitationAudit:
    citations: tuple[SmokeCitation, ...]
    checks: dict[str, str]


@dataclass(frozen=True, slots=True)
class UsageAudit:
    input_tokens: int | None
    output_tokens: int | None
    checks: dict[str, str]


@dataclass(frozen=True, slots=True)
class PrivateShadowQueryAudit:
    citation_count: int
    resolved_citation_count: int
    input_tokens: int
    output_tokens: int
    supported: bool | None
    answer_empty: bool | None
    answer: str | None = None
    citation_page: int | None = None
    citation_excerpt: str | None = None


class DiagnosticSink(Protocol):
    def capture(self, label: str, value: object) -> None: ...

    def capture_exception(self, label: str, error: BaseException) -> None: ...


@dataclass(frozen=True, slots=True)
class _SyntheticDiagnosticRequest:
    output_path: Path
    course_id: str
    exam_id: str
    lecture_id: str
    source_revision_id: str
    fixture_sha256: str


@dataclass(slots=True)
class _SyntheticDiagnosticSink:
    request: _SyntheticDiagnosticRequest
    events: list[dict[str, object]] = field(default_factory=list)
    secrets: list[str] = field(default_factory=list)
    _closed: bool = False

    @classmethod
    def open(cls, request: _SyntheticDiagnosticRequest) -> _SyntheticDiagnosticSink:
        _validate_diagnostic_request(request)
        if request.output_path.exists():
            raise LiveSmokeBlocked("synthetic diagnostic output already exists")
        return cls(request=request)

    def add_secret(self, value: str) -> None:
        if value:
            self.secrets.append(value)

    def capture(self, label: str, value: object) -> None:
        self.events.append(
            {
                "label": label,
                "value": _diagnostic_value(value, tuple(self.secrets)),
            }
        )

    def capture_exception(self, label: str, error: BaseException) -> None:
        response = _field(error, "response")
        self.capture(
            label,
            {
                "exception_type": type(error).__name__,
                "message": str(error),
                "traceback": "".join(traceback.format_exception(error)),
                "status": _field(response, "status_code"),
                "headers": _field(response, "headers"),
                "body": _response_body(response),
            },
        )

    def close(self) -> None:
        if self._closed:
            return
        payload = json.dumps(
            {
                "schema_version": 1,
                "synthetic_only": True,
                "events": self.events,
            },
            ensure_ascii=False,
            separators=(",", ":"),
        ).encode("utf-8")
        if len(payload) > _MAX_DIAGNOSTIC_BYTES:
            raise SmokeContractError(
                "synthetic diagnostic overflow",
                reason="diagnostic_overflow",
            )
        output = self.request.output_path
        file_descriptor, temporary_name = tempfile.mkstemp(
            dir=output.parent,
            prefix=f".{output.name}.",
            suffix=".tmp",
        )
        temporary = Path(temporary_name)
        try:
            try:
                _restrict_diagnostic_file(file_descriptor, temporary)
            except BaseException:
                os.close(file_descriptor)
                raise
            with os.fdopen(file_descriptor, "wb") as handle:
                handle.write(payload)
                handle.flush()
                os.fsync(handle.fileno())
            os.link(temporary, output)
        finally:
            temporary.unlink(missing_ok=True)
        self._closed = True

    def delete(self) -> None:
        self.request.output_path.unlink(missing_ok=True)
        self.events.clear()
        self.secrets.clear()


def _restrict_diagnostic_file(file_descriptor: int, path: Path) -> None:
    if not _IS_WINDOWS:
        fchmod = getattr(os, "fchmod", None)
        if callable(fchmod):
            fchmod(file_descriptor, 0o600)
        else:
            os.chmod(path, 0o600)
        return
    try:
        identity = subprocess.run(
            ["whoami", "/user", "/fo", "csv", "/nh"],
            capture_output=True,
            text=True,
            timeout=10,
            check=False,
        )
        row = next(csv.reader([identity.stdout.strip()]))
        sid = row[1].strip() if len(row) == 2 else ""
        if identity.returncode != 0 or re.fullmatch(r"S-1(?:-\d+)+", sid) is None:
            raise ValueError
        secured = subprocess.run(
            [
                "icacls",
                str(path),
                "/inheritance:r",
                "/grant:r",
                f"*{sid}:(F)",
            ],
            capture_output=True,
            text=True,
            timeout=10,
            check=False,
        )
        verified = subprocess.run(
            ["icacls", str(path), "/verify"],
            capture_output=True,
            text=True,
            timeout=10,
            check=False,
        )
        if secured.returncode != 0 or verified.returncode != 0:
            raise ValueError
        inspection_environment = dict(os.environ)
        inspection_environment["OMS_TASK28_DIAGNOSTIC_PATH"] = str(path)
        inspected = subprocess.run(
            [
                "powershell.exe",
                "-NoProfile",
                "-NonInteractive",
                "-Command",
                (
                    "$acl=Get-Acl -LiteralPath $env:OMS_TASK28_DIAGNOSTIC_PATH;"
                    "$rules=@($acl.GetAccessRules($true,$false,"
                    "[System.Security.Principal.SecurityIdentifier]));"
                    "$items=@($rules|ForEach-Object{[pscustomobject]@{"
                    "Sid=$_.IdentityReference.Value;"
                    "Allow=($_.AccessControlType -eq "
                    "[System.Security.AccessControl.AccessControlType]::Allow);"
                    "FullControl=(($_.FileSystemRights -band "
                    "[System.Security.AccessControl.FileSystemRights]::FullControl) "
                    "-eq [System.Security.AccessControl.FileSystemRights]::FullControl);"
                    "Inherited=$_.IsInherited}});"
                    "[pscustomobject]@{Protected=$acl.AreAccessRulesProtected;"
                    "Rules=$items}|ConvertTo-Json -Compress -Depth 4"
                ),
            ],
            capture_output=True,
            text=True,
            timeout=10,
            check=False,
            env=inspection_environment,
        )
        inspection = json.loads(inspected.stdout)
        rules = inspection.get("Rules") if isinstance(inspection, Mapping) else None
        if isinstance(rules, Mapping):
            rule_list = [rules]
        elif isinstance(rules, list):
            rule_list = rules
        else:
            rule_list = []
        rule = rule_list[0] if len(rule_list) == 1 else None
        if (
            inspected.returncode != 0
            or not isinstance(inspection, Mapping)
            or inspection.get("Protected") is not True
            or not isinstance(rule, Mapping)
            or rule.get("Sid") != sid
            or rule.get("Allow") is not True
            or rule.get("FullControl") is not True
            or rule.get("Inherited") is not False
        ):
            raise ValueError
    except (OSError, StopIteration, subprocess.SubprocessError, ValueError):
        raise SmokeContractError(
            "synthetic diagnostic permissions were unavailable",
            reason="diagnostic_permissions_unavailable",
        ) from None


def _has_reparse_component(path: Path) -> bool:
    cursor = path
    while True:
        try:
            info = cursor.lstat()
        except OSError:
            return True
        attributes = getattr(info, "st_file_attributes", 0)
        if stat.S_ISLNK(info.st_mode) or (
            _IS_WINDOWS
            and bool(attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0))
        ):
            return True
        parent = cursor.parent
        if parent == cursor:
            return False
        cursor = parent


def _validate_private_diagnostic_capability(path: object) -> Path:
    if not isinstance(path, Path):
        raise LiveSmokeBlocked("private diagnostic capability is unavailable")
    supplied = os.getenv("OMS_TASK28_PRIVATE_DIAGNOSTIC_PATH")
    supplied_path = Path(supplied) if supplied else None
    if (
        supplied_path is None
        or not path.is_absolute()
        or not supplied_path.is_absolute()
        or ".." in path.parts
        or ".." in supplied_path.parts
        or path != supplied_path
    ):
        raise LiveSmokeBlocked("private diagnostic capability is unavailable")
    path_has_reparse = _has_reparse_component(path.parent)
    supplied_has_reparse = _has_reparse_component(supplied_path.parent)
    try:
        output = path.resolve(strict=False)
        launcher_output = supplied_path.resolve(strict=False)
    except (OSError, RuntimeError):
        raise LiveSmokeBlocked("private diagnostic capability is invalid") from None
    if output != launcher_output:
        raise LiveSmokeBlocked("private diagnostic capability did not match the launcher")
    if (
        output.name != _PRIVATE_DIAGNOSTIC_FILENAME
        or output.parent.name != "diagnostic"
        or re.fullmatch(r"[0-9a-f]{32}", output.parent.parent.name) is None
        or output.parent.parent.parent.name != "oms-task28-runs"
        or path_has_reparse
        or supplied_has_reparse
        or output.is_relative_to(_REPOSITORY_ROOT)
    ):
        raise LiveSmokeBlocked("private diagnostic capability is invalid")
    if output.exists() or output.is_symlink():
        raise LiveSmokeBlocked("private diagnostic output already exists")
    return output


def _private_terminal_diagnostic_payload(
    error: BaseException,
    *,
    failure_stage: str,
    input_identity: str,
) -> bytes:
    status_code = getattr(error, "provider_status_code", None)
    if type(status_code) is not int or not 100 <= status_code <= 599:
        status_code = None
    reason = getattr(error, "diagnostic_code", "none")
    if reason not in _PRIVATE_DIAGNOSTIC_REASONS:
        reason = "none"
    if reason == "unknown_provider" and status_code == 400:
        reason = "provider_bad_request"
    exception_kind = (
        "gemini_provider_error"
        if isinstance(error, GeminiProviderError)
        else "smoke_contract_error"
        if isinstance(error, SmokeContractError)
        else "unclassified_error"
    )
    payload = {
        "schema_version": 1,
        "exception_kind": exception_kind,
        "provider_status_code": status_code,
        "provider_reason": reason,
        "provider_message": _PRIVATE_DIAGNOSTIC_MESSAGES.get(
            reason, "Provider failure classification was unavailable."
        ),
        "failure_stage": failure_stage
        if failure_stage
        in {
            "prior_state_check",
            "create_store",
            "upload_input",
            "import_input",
            "wait_for_import",
            "positive_query",
            "positive_validation",
            "negative_query",
            "negative_validation",
            "cleanup",
            "unknown",
        }
        else "unknown",
        "failure_input_identity": input_identity
        if input_identity in {"none", "pptx", "pdf", "normalized_markdown", "visual_asset", "unknown"}
        else "unknown",
    }
    return (json.dumps(payload, ensure_ascii=True, sort_keys=True, separators=(",", ":")) + "\n").encode(
        "utf-8"
    )


def _write_private_terminal_diagnostic(
    output: Path,
    error: BaseException,
    *,
    failure_stage: str,
    input_identity: str,
) -> str:
    if output.exists() or output.is_symlink():
        raise LiveSmokeBlocked("private diagnostic output already exists")
    payload = _private_terminal_diagnostic_payload(
        error,
        failure_stage=failure_stage,
        input_identity=input_identity,
    )
    if len(payload) > _MAX_PRIVATE_DIAGNOSTIC_BYTES:
        raise SmokeContractError("private diagnostic overflow", reason="diagnostic_overflow")
    descriptor, temporary_name = tempfile.mkstemp(
        dir=output.parent,
        prefix=f".{output.name}.",
        suffix=".tmp",
    )
    temporary = Path(temporary_name)
    try:
        try:
            _restrict_diagnostic_file(descriptor, temporary)
        except BaseException:
            os.close(descriptor)
            raise
        with os.fdopen(descriptor, "wb") as handle:
            handle.write(payload)
            handle.flush()
            os.fsync(handle.fileno())
        os.link(temporary, output)
    finally:
        temporary.unlink(missing_ok=True)
    return hashlib.sha256(payload).hexdigest()


def _finalize_private_terminal_diagnostic(
    mode: Literal["public_matrix", "private_acceptance"],
    private_diagnostic_path: Path | None,
    error: BaseException,
    *,
    failure_stage: str,
    input_identity: str,
) -> str | None:
    if mode != "private_acceptance" or private_diagnostic_path is None:
        return None
    verified_path = _validate_private_diagnostic_capability(private_diagnostic_path)
    return _write_private_terminal_diagnostic(
        verified_path,
        error,
        failure_stage=failure_stage,
        input_identity=input_identity,
    )


class ShadowSession(Protocol):
    model_contract: tuple[str, str, str, str]

    async def create_store(self, display_name: str, embedding_model: str) -> str: ...

    async def find_stores(self, display_name: str) -> tuple[str, ...]: ...

    async def upload_input(
        self,
        display_name: str,
        path: Path,
        media_type: str,
    ) -> str: ...

    async def import_input(
        self,
        store_name: str,
        file_name: str,
        metadata: tuple[tuple[str, str], ...],
        chunking: object | None,
    ) -> str: ...

    async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]: ...

    async def wait_for_import(self, operation_name: str) -> str: ...

    async def query_private(
        self,
        store_name: str,
        prompt: str,
        scope: SmokeScope,
        *,
        source_revision_id: str,
        manifest: object,
        file_bindings: tuple[tuple[str, str], ...],
        require_structured_no_result: bool = False,
        require_structured_supported: bool = False,
    ) -> PrivateShadowQueryAudit: ...

    async def delete_document(self, document_name: str) -> None: ...

    async def delete_file(self, file_name: str) -> None: ...

    async def delete_store(self, store_name: str) -> None: ...


class GoogleGenaiSmokeSession:
    """Minimum live adapter for the pinned google-genai 2.14.0 contract."""

    def __init__(
        self,
        api_key: str,
        *,
        sdk_factory: SdkFactory | None = None,
        diagnostic_sink: DiagnosticSink | None = None,
        sleep: Callable[[float], Awaitable[None]] | None = None,
        clock: Callable[[], float] | None = None,
    ) -> None:
        self._config = GeminiConfig(api_key=SecretStr(api_key))
        self._clients = GeminiClientFactory(self._config, sdk_factory=sdk_factory)
        self._diagnostic_sink = diagnostic_sink
        self._sleep = asyncio.sleep if sleep is None else sleep
        self._clock = monotonic if clock is None else clock
        self._transient_attempts = 0
        self.model_contract = (
            self._config.sdk_version,
            self._config.file_search_model,
            self._config.embedding_model,
            self._config.api_version,
        )
        self._store_name: str | None = None
        self._document_name: str | None = None
        self._file_name: str | None = None

    @property
    def transient_attempts(self) -> int:
        return self._transient_attempts

    def _record_transient_retry(self) -> None:
        self._transient_attempts += 1

    async def create_store(self, display_name: str, embedding_model: str) -> str:
        async with self._clients.client() as client:
            created = await _provider_call(
                lambda: client.file_search_stores.create(
                    config={
                        "display_name": display_name,
                        "embedding_model": embedding_model,
                    }
                ),
                diagnostic_sink=self._diagnostic_sink,
                label="create_store",
            )
        self._store_name = _provider_identity(created, "store")
        return self._store_name

    async def find_stores(self, display_name: str) -> tuple[str, ...]:
        async with self._clients.client() as client:
            listed = await _provider_call(
                lambda: client.file_search_stores.list(config={"page_size": 20}),
                diagnostic_sink=self._diagnostic_sink,
                label="find_stores.request",
                capture_response=False,
                idempotent=True,
                sleep=self._sleep,
                clock=self._clock,
                on_transient_retry=self._record_transient_retry,
            )
            if not isinstance(listed, AsyncIterable):
                raise SmokeContractError(
                    "private shadow store reconciliation was unavailable",
                    reason="private_reconciliation_failed",
                )
            matched: list[str] = []
            inspected = 0
            try:
                async for item in listed:
                    inspected += 1
                    if inspected > 1000:
                        raise SmokeContractError(
                            "private shadow store reconciliation exceeded its bound",
                            reason="private_reconciliation_failed",
                        )
                    if _field(item, "display_name") == display_name:
                        matched.append(_provider_identity(item, "store"))
            except Exception as error:
                if self._diagnostic_sink is not None:
                    self._diagnostic_sink.capture_exception("find_stores.iteration", error)
                if isinstance(error, (GeminiProviderError, SmokeContractError)):
                    raise
                raise translate_gemini_error(error) from None
        return tuple(sorted(set(matched)))

    async def upload_input(
        self,
        display_name: str,
        path: Path,
        media_type: str,
    ) -> str:
        async with self._clients.client() as client:
            uploaded = await _provider_call(
                lambda: client.files.upload(
                    file=path,
                    config={"display_name": display_name, "mime_type": media_type},
                ), diagnostic_sink=self._diagnostic_sink, label="upload_input"
            )
        self._file_name = _provider_identity(uploaded, "file")
        return self._file_name

    async def import_input(
        self,
        store_name: str,
        file_name: str,
        metadata: tuple[tuple[str, str], ...],
        chunking: object | None,
    ) -> str:
        self._store_name = store_name
        async with self._clients.client() as client:
            operation = await _provider_call(
                lambda: client.file_search_stores.import_file(
                    file_search_store_name=store_name,
                    file_name=file_name,
                    config=build_import_file_config(
                        [{"key": key, "string_value": value} for key, value in metadata],
                        chunking,
                    ),
                ), diagnostic_sink=self._diagnostic_sink, label="import_input"
            )
        return _provider_identity(operation, "operation")

    async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]:
        expected = frozenset(display_names)
        if not expected or len(expected) != len(display_names):
            raise SmokeContractError(
                "private shadow reconciliation scope was invalid",
                reason="private_reconciliation_failed",
            )
        matched: list[str] = []
        inspected = 0
        async with self._clients.client() as client:
            listed = await _provider_call(
                lambda: client.files.list(config={"page_size": 100}),
                diagnostic_sink=self._diagnostic_sink,
                label="find_files.request",
                capture_response=False,
                idempotent=True,
                sleep=self._sleep,
                clock=self._clock,
                on_transient_retry=self._record_transient_retry,
            )
            if not isinstance(listed, AsyncIterable):
                raise SmokeContractError(
                    "private shadow file reconciliation was unavailable",
                    reason="private_reconciliation_failed",
                )
            try:
                async for item in listed:
                    inspected += 1
                    if inspected > 1000:
                        raise SmokeContractError(
                            "private shadow file reconciliation exceeded its bound",
                            reason="private_reconciliation_failed",
                        )
                    if _field(item, "display_name") in expected:
                        matched.append(_provider_identity(item, "file"))
            except Exception as error:
                if self._diagnostic_sink is not None:
                    self._diagnostic_sink.capture_exception("find_files.iteration", error)
                if isinstance(error, (GeminiProviderError, SmokeContractError)):
                    raise
                raise translate_gemini_error(error) from None
        return tuple(sorted(set(matched)))

    async def wait_for_import(self, operation_name: str) -> str:
        try:
            operation_type = import_module("google.genai.types").ImportFileOperation
            operation = operation_type(name=operation_name)
        except Exception as error:
            raise translate_gemini_error(error) from None
        deadline = self._clock() + self._config.operation_timeout_seconds
        async with self._clients.client() as client:
            while True:
                remaining = deadline - self._clock()
                if remaining <= 0:
                    raise SmokeTemporaryFailure("Gemini import operation timed out")
                try:
                    async with asyncio.timeout(remaining):
                        current_operation = operation
                        operation = await _provider_call(
                            partial(client.operations.get, current_operation),
                            diagnostic_sink=self._diagnostic_sink,
                            label="wait_for_import.response",
                            idempotent=True,
                            sleep=self._sleep,
                            clock=self._clock,
                            retry_deadline=deadline,
                            on_transient_retry=self._record_transient_retry,
                        )
                except TimeoutError:
                    if self._diagnostic_sink is not None:
                        self._diagnostic_sink.capture_exception(
                            "wait_for_import.timeout",
                            TimeoutError("Gemini import operation timed out"),
                        )
                    raise SmokeTemporaryFailure("Gemini import operation timed out") from None
                except GeminiProviderError:
                    raise
                except Exception as error:
                    if self._diagnostic_sink is not None:
                        self._diagnostic_sink.capture_exception(
                            "wait_for_import.request_failed",
                            error,
                        )
                    raise translate_gemini_error(error) from None
                if bool(_field(operation, "done")):
                    break
                await asyncio.sleep(self._config.operation_poll_seconds)
        operation_error = _field(operation, "error")
        if operation_error:
            if self._diagnostic_sink is not None:
                self._diagnostic_sink.capture("wait_for_import.operation_error", operation_error)
            status = _field(operation_error, "code")
            raise translate_gemini_error(
                _OperationFailure(status if isinstance(status, int) else None)
            ) from None
        response = _field(operation, "response")
        document_name = _provider_identity(response, "document", "document_name")
        if self._store_name is not None:
            parent = _field(response, "parent")
            if parent is not None and (
                not isinstance(parent, str)
                or not _resource_identity_matches(parent, self._store_name)
            ):
                raise SmokeContractError("Gemini document parent did not match the store")
            prefix = f"{self._store_name}/documents/"
            if "/" not in document_name:
                document_name = f"{prefix}{document_name}"
            elif not document_name.startswith(prefix) or "/" in document_name[len(prefix) :]:
                raise SmokeContractError("Gemini document identity did not match the store")
        self._document_name = document_name
        return self._document_name

    async def query(
        self,
        store_name: str,
        prompt: str,
        scope: SmokeScope,
        *,
        response_schema: type[SmokeAnswer] | None,
        omit_thinking: bool,
    ) -> SmokeQueryResult:
        if not omit_thinking:
            raise SmokeContractError("Task 2.8 smoke requires omitted thinking configuration")
        body: dict[str, object] = {
            "model": self._config.file_search_model,
            "input": prompt,
            "store": False,
            "tools": [
                {
                    "type": "file_search",
                    "file_search_store_names": [store_name],
                    "metadata_filter": _scope_filter(scope),
                }
            ],
        }
        if response_schema is not None:
            body["response_format"] = {
                "type": "text",
                "mime_type": "application/json",
                "schema": response_schema.model_json_schema(),
            }
        if self._diagnostic_sink is not None:
            self._diagnostic_sink.capture("query.request", body)
        async with self._clients.client() as client:
            response = await _provider_call(
                lambda: client.interactions.create(**body),
                diagnostic_sink=self._diagnostic_sink,
                label="query.response",
                idempotent=True,
                sleep=self._sleep,
                clock=self._clock,
                on_transient_retry=self._record_transient_retry,
            )
        output_text = _interaction_output(response)
        if response_schema is None:
            answer = {"answer": output_text, "supported": True}
        else:
            try:
                answer = response_schema.model_validate_json(output_text).model_dump(mode="json")
            except ValidationError:
                raise SmokeContractError(
                    "Gemini structured output did not match the required schema",
                    reason="structured_output_invalid",
                ) from None
        usage = _audit_usage(_field(response, "usage"))
        citations = _audit_citations(
            response,
            scope,
            self._document_name,
            self._file_name,
        )
        return SmokeQueryResult(
            answer=answer,
            citations=citations.citations,
            input_tokens=usage.input_tokens,
            output_tokens=usage.output_tokens,
            citation_checks=tuple(citations.checks.items()),
            usage_checks=tuple(usage.checks.items()),
        )

    async def query_private(
        self,
        store_name: str,
        prompt: str,
        scope: SmokeScope,
        *,
        source_revision_id: str,
        manifest: object,
        file_bindings: tuple[tuple[str, str], ...],
        require_structured_no_result: bool = False,
        require_structured_supported: bool = False,
    ) -> PrivateShadowQueryAudit:
        body: dict[str, object] = {
            "model": self._config.file_search_model,
            "input": prompt,
            "store": False,
            "tools": [
                {
                    "type": "file_search",
                    "file_search_store_names": [store_name],
                    "metadata_filter": _scope_filter(scope),
                }
            ],
        }
        if require_structured_no_result or require_structured_supported:
            body["response_format"] = {
                "type": "text",
                "mime_type": "application/json",
                "schema": SmokeAnswer.model_json_schema(),
            }
        async with self._clients.client() as client:
            response = await _provider_call(
                lambda: client.interactions.create(**body), diagnostic_sink=self._diagnostic_sink,
                label="query_private",
                idempotent=True,
                sleep=self._sleep,
                clock=self._clock,
                on_transient_retry=self._record_transient_retry,
            )
        output_text = _interaction_output(response)
        supported: bool | None = None
        answer_empty: bool | None = None
        if require_structured_no_result or require_structured_supported:
            try:
                answer = SmokeAnswer.model_validate_json(output_text)
            except ValidationError:
                raise SmokeContractError(
                    "private shadow structured answer was invalid",
                    reason=(
                        "negative_answer_invalid"
                        if require_structured_no_result
                        else "structured_output_invalid"
                    ),
                ) from None
            supported = answer.supported
            answer_empty = answer.answer == ""
        else:
            supported = True
            answer_empty = False
        usage = _audit_usage(_field(response, "usage"))
        if (
            usage.input_tokens is None
            or usage.output_tokens is None
            or any(value != "passed" for value in usage.checks.values())
        ):
            raise SmokeContractError(
                "private shadow usage was invalid",
                reason="private_usage_invalid",
            )
        citation_count, resolved_count, page, excerpt = _private_shadow_citation_counts(
            response,
            store_name=store_name,
            scope=scope,
            source_revision_id=source_revision_id,
            manifest=manifest,
            file_bindings=file_bindings,
        )
        return PrivateShadowQueryAudit(
            citation_count,
            resolved_count,
            usage.input_tokens,
            usage.output_tokens,
            supported,
            answer_empty,
            answer.answer if (require_structured_no_result or require_structured_supported) else output_text,
            page,
            excerpt,
        )

    async def list_documents(self, store_name: str) -> tuple[str, ...]:
        async with self._clients.client() as client:
            listed = await _provider_call(
                lambda: client.file_search_stores.documents.list(parent=store_name),
                diagnostic_sink=self._diagnostic_sink,
                label="list_documents",
                idempotent=True,
                sleep=self._sleep,
                clock=self._clock,
                on_transient_retry=self._record_transient_retry,
            )
            documents = await _collect(listed)
        return tuple(sorted(_provider_identity(item, "document") for item in documents))

    async def delete_document(self, document_name: str) -> None:
        await self._delete_resource(
            "document",
            lambda client: client.file_search_stores.documents.delete(
                    name=document_name,
                    config={"force": True},
            ),
        )

    async def delete_file(self, file_name: str) -> None:
        await self._delete_resource(
            "file",
            lambda client: client.files.delete(name=file_name),
        )

    async def delete_store(self, store_name: str) -> None:
        await self._delete_resource(
            "store",
            lambda client: client.file_search_stores.delete(
                    name=store_name,
                    config={"force": True},
            ),
        )

    async def _delete_resource(
        self,
        label: str,
        request: Callable[[Any], Awaitable[Any]],
    ) -> None:
        try:
            sdk_client = self._clients._build_sdk_client()  # noqa: SLF001
            client = sdk_client.aio
            close = client.aclose
        except Exception as error:
            if self._diagnostic_sink is not None:
                self._diagnostic_sink.capture_exception(
                    f"cleanup.{label}.client_setup_failed",
                    error,
                )
            raise translate_gemini_error(error) from None
        request_error: GeminiProviderError | None = None
        try:
            try:
                response = await request(client)
            except Exception as error:
                if self._diagnostic_sink is not None:
                    self._diagnostic_sink.capture_exception(
                        f"cleanup.{label}.delete_request_failed",
                        error,
                    )
                request_error = translate_gemini_error(error)
            else:
                if self._diagnostic_sink is not None:
                    self._diagnostic_sink.capture(
                        f"cleanup.{label}.delete_response",
                        response,
                    )
        finally:
            try:
                await close()
            except Exception as error:
                if self._diagnostic_sink is not None:
                    self._diagnostic_sink.capture_exception(
                        f"cleanup.{label}.context_close_failed",
                        error,
                    )
                if request_error is None:
                    request_error = translate_gemini_error(error)
        if request_error is not None:
            raise request_error from None


async def _provider_call(
    request: Callable[[], Awaitable[Any]],
    *,
    diagnostic_sink: DiagnosticSink | None = None,
    label: str = "provider_call",
    capture_response: bool = True,
    idempotent: bool = False,
    sleep: Callable[[float], Awaitable[None]] | None = None,
    clock: Callable[[], float] | None = None,
    retry_deadline: float | None = None,
    on_transient_retry: Callable[[], None] | None = None,
) -> Any:
    retry_sleep = asyncio.sleep if sleep is None else sleep
    retry_clock = monotonic if clock is None else clock
    for attempt in range(len(_RETRY_DELAYS) + 1):
        try:
            response = await request()
        except GeminiProviderError as error:
            translated = error
        except Exception as error:
            if diagnostic_sink is not None:
                diagnostic_sink.capture_exception(f"{label}.failed", error)
            translated = translate_gemini_error(error)
        else:
            if capture_response and diagnostic_sink is not None:
                diagnostic_sink.capture(label, response)
            return response
        if not idempotent or not isinstance(translated, GeminiTransientError) or attempt == len(_RETRY_DELAYS):
            raise translated from None
        delay = _RETRY_DELAYS[attempt]
        if retry_deadline is not None and retry_clock() + delay >= retry_deadline:
            raise translated from None
        if on_transient_retry is not None:
            on_transient_retry()
        await retry_sleep(delay)


def _session_transient_attempts(session: object) -> int:
    value = getattr(session, "transient_attempts", 0)
    if type(value) is int and 0 <= value <= _MAX_TRANSIENT_ATTEMPTS:
        return value
    return 0


def _provider_identity(value: object, label: str, field: str = "name") -> str:
    identity = _field(value, field)
    if not isinstance(identity, str) or not identity.strip():
        raise SmokeContractError(f"Gemini {label} identity was unavailable")
    normalized = identity.strip()
    if len(normalized) > 500 or not normalized.isprintable():
        raise SmokeContractError(f"Gemini {label} identity was invalid")
    return normalized


def _resource_identity_matches(actual: str, expected: str) -> bool:
    return actual == expected or ("/" not in actual and expected.endswith(f"/{actual}"))


def _field(value: object, name: str) -> object | None:
    if isinstance(value, Mapping):
        return value.get(name)
    try:
        return getattr(value, name, None)
    except Exception:
        return None


def _synthetic_diagnostic_request(output_path: Path) -> _SyntheticDiagnosticRequest:
    return _SyntheticDiagnosticRequest(
        output_path=output_path,
        course_id=SYNTHETIC_COURSE_ID,
        exam_id=SYNTHETIC_EXAM_ID,
        lecture_id=SYNTHETIC_LECTURE_ID,
        source_revision_id=SYNTHETIC_REVISION_ID,
        fixture_sha256=hashlib.sha256(synthetic_pdf_bytes()).hexdigest(),
    )


def _validate_diagnostic_request(request: _SyntheticDiagnosticRequest) -> None:
    expected = _synthetic_diagnostic_request(request.output_path)
    if request != expected:
        raise LiveSmokeBlocked("synthetic diagnostic scope mismatch")
    output = request.output_path
    if not output.is_absolute() or not output.parent.is_dir():
        raise LiveSmokeBlocked("synthetic diagnostic output path is unavailable")
    try:
        output.resolve().relative_to(_REPOSITORY_ROOT.resolve())
    except ValueError:
        pass
    else:
        raise LiveSmokeBlocked("synthetic diagnostic output must be outside the repository")


def _diagnostic_value(value: object, secrets: tuple[str, ...]) -> object:
    if isinstance(value, SecretStr):
        return "[REDACTED]"
    if isinstance(value, BaseModel):
        return _diagnostic_value(value.model_dump(mode="json"), secrets)
    if isinstance(value, Mapping):
        cleaned: dict[str, object] = {}
        for key, item in value.items():
            name = str(key)
            normalized = name.casefold().replace("_", "-")
            if normalized == "headers":
                cleaned[name] = "[REDACTED]"
            elif (
                normalized in {
                    "api-key",
                    "authorization",
                    "cookie",
                    "credentials",
                    "set-cookie",
                    "x-goog-api-key",
                }
                or normalized.endswith("-credential")
                or normalized.endswith("-secret")
            ):
                cleaned[name] = "[REDACTED]"
            else:
                cleaned[name] = _diagnostic_value(item, secrets)
        return cleaned
    if isinstance(value, (list, tuple, set, frozenset)):
        return [_diagnostic_value(item, secrets) for item in value]
    if isinstance(value, bytes):
        cleaned_bytes = value
        for secret in secrets:
            cleaned_bytes = cleaned_bytes.replace(
                secret.encode("utf-8"),
                b"[REDACTED]",
            )
        return {"base64": base64.b64encode(cleaned_bytes).decode("ascii")}
    if value is None or isinstance(value, (bool, int, float)):
        return value
    if isinstance(value, str):
        cleaned_text = value
        for secret in secrets:
            cleaned_text = cleaned_text.replace(secret, "[REDACTED]")
        return cleaned_text
    attributes = getattr(value, "__dict__", None)
    if isinstance(attributes, Mapping):
        return _diagnostic_value(attributes, secrets)
    return _diagnostic_value(repr(value), secrets)


def _response_body(response: object | None) -> object | None:
    if response is None:
        return None
    content = _field(response, "content")
    if content is not None:
        return content
    return _field(response, "text")


def _scope_filter(scope: SmokeScope) -> str:
    values = {
        "course_id": scope.course_id,
        "exam_id": scope.exam_id,
        "lecture_id": scope.lecture_id,
    }
    if any(
        not value
        or len(value) > 128
        or not all(character.isalnum() or character in ".:_-" for character in value)
        for value in values.values()
    ):
        raise SmokeContractError("Gemini metadata scope was invalid")
    return " AND ".join(f'{key}="{value}"' for key, value in values.items())


def _interaction_output(response: object) -> str:
    output = _field(response, "output_text")
    if output is None:
        raise SmokeContractError(
            "Gemini structured output was absent",
            reason="structured_output_absent",
        )
    if not isinstance(output, str) or not output:
        raise SmokeContractError(
            "Gemini structured output was invalid",
            reason="structured_output_invalid",
        )
    return output


def _audit_usage(value: object | None) -> UsageAudit:
    checks: dict[str, str] = {}
    counts: list[int | None] = []
    for field_name, check_name in (
        ("total_input_tokens", "usage_input"),
        ("total_output_tokens", "usage_output"),
    ):
        raw = _field(value, field_name) if value is not None else None
        if raw is None:
            counts.append(None)
            checks[check_name] = f"{check_name}_absent"
        elif isinstance(raw, bool) or not isinstance(raw, int) or raw < 0:
            counts.append(None)
            checks[check_name] = f"{check_name}_invalid"
        else:
            counts.append(raw)
            checks[check_name] = "passed"
    return UsageAudit(counts[0], counts[1], checks)


def _audit_citations(
    response: object,
    scope: SmokeScope,
    document_name: str | None,
    file_name: str | None,
) -> CitationAudit:
    steps = _field(response, "steps")
    blocked = {name: "blocked_by_citation_presence" for name in _CITATION_CHECKS}
    if steps is None:
        blocked["citation_presence"] = "citation_steps_absent"
        return CitationAudit((), blocked)
    if not isinstance(steps, Iterable) or isinstance(steps, (str, bytes, Mapping)):
        blocked["citation_presence"] = "citation_steps_invalid"
        return CitationAudit((), blocked)
    found: list[SmokeCitation] = []
    results: dict[str, list[str]] = {name: [] for name in _CITATION_CHECKS}
    expected = {
        "course_id": scope.course_id,
        "exam_id": scope.exam_id,
        "lecture_id": scope.lecture_id,
        "source_revision_id": SYNTHETIC_REVISION_ID,
    }
    store_name = None
    if document_name is not None:
        candidate, separator, suffix = document_name.partition("/documents/")
        if separator and candidate and suffix and "/" not in suffix:
            store_name = candidate
    saw_model_output = False
    saw_content = False
    saw_annotations = False
    invalid_content = False
    invalid_annotations = False
    for step in steps:
        if _field(step, "type") != "model_output":
            continue
        saw_model_output = True
        contents = _field(step, "content")
        if contents is None:
            continue
        if not isinstance(contents, Iterable) or isinstance(
            contents, (str, bytes, Mapping)
        ):
            invalid_content = True
            continue
        saw_content = True
        for content in contents:
            if _field(content, "type") != "text":
                continue
            annotations = _field(content, "annotations")
            if annotations is None:
                continue
            if not isinstance(annotations, Iterable) or isinstance(
                annotations, (str, bytes, Mapping)
            ):
                invalid_annotations = True
                continue
            saw_annotations = True
            for annotation in annotations:
                if _field(annotation, "type") != "file_citation":
                    continue
                results["citation_presence"].append("passed")
                current: dict[str, str] = {}
                raw_metadata = _field(annotation, "custom_metadata")
                if raw_metadata is None:
                    current["citation_scope_binding"] = "citation_metadata_absent"
                elif isinstance(raw_metadata, (str, bytes)) or not isinstance(
                    raw_metadata,
                    (Mapping, Iterable),
                ):
                    current["citation_scope_binding"] = "citation_metadata_invalid"
                else:
                    try:
                        actual = _string_metadata(raw_metadata)
                    except SmokeContractError:
                        current["citation_scope_binding"] = "citation_metadata_invalid"
                    else:
                        current["citation_scope_binding"] = (
                            "passed"
                            if all(actual.get(key) == value for key, value in expected.items())
                            else "citation_scope_mismatch"
                        )
                if document_name is None or store_name is None:
                    current["citation_document_binding"] = (
                        "citation_document_identity_unavailable"
                    )
                locator = _field(annotation, "document_uri")
                if document_name is not None and store_name is not None:
                    if locator is None:
                        current["citation_document_binding"] = (
                            "citation_document_uri_absent"
                        )
                    elif (
                        not isinstance(locator, str)
                        or not locator
                        or len(locator) > 500
                        or not locator.isprintable()
                    ):
                        current["citation_document_binding"] = (
                            "citation_document_uri_invalid"
                        )
                    else:
                        current["citation_document_binding"] = (
                            "passed"
                            if locator in {store_name, document_name}
                            else "citation_wrong_document"
                        )
                cited_file = _field(annotation, "file_name")
                if cited_file is None:
                    current["citation_file_binding"] = "citation_file_absent"
                elif (
                    not isinstance(cited_file, str)
                    or not cited_file
                    or len(cited_file) > 500
                    or not cited_file.isprintable()
                ):
                    current["citation_file_binding"] = "citation_file_invalid"
                else:
                    current["citation_file_binding"] = (
                        "passed"
                        if file_name is not None
                        and _resource_identity_matches(cited_file, file_name)
                        else "citation_wrong_file"
                    )
                page_value = _field(annotation, "page_number")
                if page_value is None:
                    page = None
                    current["citation_page_binding"] = "citation_page_absent"
                else:
                    try:
                        page = _optional_page(page_value)
                    except SmokeContractError:
                        page = None
                        current["citation_page_binding"] = "citation_page_invalid"
                    else:
                        current["citation_page_binding"] = (
                            "passed" if page == 1 else "citation_page_mismatch"
                        )
                try:
                    excerpt = _citation_excerpt(annotation, _field(content, "text"))
                except SmokeContractError as error:
                    excerpt = ""
                    current["citation_excerpt_binding"] = (
                        "citation_excerpt_absent"
                        if error.reason == "citation_excerpt_unavailable"
                        else error.reason or "citation_excerpt_invalid"
                    )
                else:
                    current["citation_excerpt_binding"] = (
                        "passed" if SYNTHETIC_FACT in excerpt else "citation_excerpt_mismatch"
                    )
                for name, outcome in current.items():
                    results[name].append(outcome)
                if document_name is not None and all(
                    current.get(name) == "passed"
                    or (
                        name == "citation_document_binding"
                        and current.get(name) == "citation_document_uri_absent"
                    )
                    for name in _CITATION_CHECKS
                    if name != "citation_presence"
                ):
                    found.append(
                        SmokeCitation(
                            document_name=document_name,
                            page_number=page,
                            excerpt=excerpt,
                        )
                    )
    if not results["citation_presence"]:
        if not saw_model_output:
            diagnosis = "citation_steps_absent"
        elif invalid_content:
            diagnosis = "citation_content_invalid"
        elif not saw_content:
            diagnosis = "citation_content_absent"
        elif invalid_annotations:
            diagnosis = "citation_annotations_invalid"
        elif not saw_annotations:
            diagnosis = "citation_annotations_absent"
        else:
            diagnosis = "positive_citation_missing"
        blocked["citation_presence"] = diagnosis
        return CitationAudit((), blocked)
    checks: dict[str, str] = {}
    for name in _CITATION_CHECKS:
        outcomes = results[name]
        checks[name] = "passed" if "passed" in outcomes else outcomes[0]
    return CitationAudit(tuple(found), checks)


def _private_shadow_citation_counts(
    response: object,
    *,
    store_name: str,
    scope: SmokeScope,
    source_revision_id: str,
    manifest: Any,
    file_bindings: tuple[tuple[str, str], ...],
) -> tuple[int, int, int | None, str | None]:
    from oms_hub.providers.gemini.citations import ProviderCitation, map_provider_citation

    steps = _field(response, "steps")
    if not isinstance(steps, Iterable) or isinstance(steps, (str, bytes, Mapping)):
        return 0, 0, None, None
    inputs = {item.input_key: item for item in manifest.inputs}
    expected_scope = {
        "course_id": scope.course_id,
        "exam_id": scope.exam_id,
        "lecture_id": scope.lecture_id,
        "source_revision_id": source_revision_id,
    }
    citation_count = 0
    resolved_count = 0
    first_page: int | None = None
    first_excerpt: str | None = None
    for step in steps:
        contents = _field(step, "content")
        if _field(step, "type") != "model_output" or not isinstance(contents, Iterable):
            continue
        for content in contents:
            annotations = _field(content, "annotations")
            if not isinstance(annotations, Iterable) or isinstance(
                annotations, (str, bytes, Mapping)
            ):
                continue
            for annotation in annotations:
                if _field(annotation, "type") != "file_citation":
                    continue
                citation_count += 1
                try:
                    metadata = _string_metadata(_field(annotation, "custom_metadata"))
                    input_key = metadata.get("input_key")
                    item = inputs.get(input_key)
                    if (
                        item is None
                        or any(metadata.get(key) != value for key, value in expected_scope.items())
                        or metadata.get("input_kind") != item.input_kind
                        or metadata.get("input_sha256") != item.sha256
                    ):
                        continue
                    cited_file = _field(annotation, "file_name")
                    if not isinstance(cited_file, str):
                        continue
                    matched_keys = {
                        key
                        for provider_file, key in file_bindings
                        if _resource_identity_matches(cited_file, provider_file)
                    }
                    if matched_keys != {input_key}:
                        continue
                    document_uri = _field(annotation, "document_uri")
                    if document_uri is not None and (
                        not isinstance(document_uri, str)
                        or (
                            document_uri != store_name
                            and not document_uri.startswith(f"{store_name}/documents/")
                        )
                    ):
                        continue
                    excerpt = " ".join(
                        _citation_excerpt(annotation, _field(content, "text")).split()
                    )
                    page = _optional_page(_field(annotation, "page_number"))
                    citation = ProviderCitation(item.path.name, excerpt, page)
                    if map_provider_citation(citation, manifest) is not None:
                        resolved_count += 1
                        if first_excerpt is None:
                            first_page, first_excerpt = page, excerpt
                except (SmokeContractError, TypeError, ValueError):
                    continue
    return citation_count, resolved_count, first_page, first_excerpt


def _citations(
    response: object,
    store_name: str,
    scope: SmokeScope,
    document_name: str | None,
    file_display_name: str | None,
) -> tuple[SmokeCitation, ...]:
    del store_name
    audit = _audit_citations(response, scope, document_name, file_display_name)
    for name in (
        "citation_scope_binding",
        "citation_file_binding",
        "citation_document_binding",
        "citation_page_binding",
        "citation_excerpt_binding",
    ):
        reason = audit.checks[name]
        if reason not in {
            "passed",
            "blocked_by_citation_presence",
            "citation_document_uri_absent",
        }:
            messages = {
                "citation_scope_binding": (
                    "Gemini citation metadata did not match the requested scope"
                ),
                "citation_file_binding": "Gemini citation referenced the wrong file",
                "citation_document_binding": "Gemini citation referenced the wrong document",
                "citation_page_binding": "Gemini citation page did not satisfy the contract",
                "citation_excerpt_binding": "Gemini citation excerpt did not satisfy the contract",
            }
            raise SmokeContractError(messages[name], reason=reason)
    return audit.citations


def _string_metadata(value: object) -> dict[str, str]:
    if isinstance(value, Mapping):
        if len(value) > 16 or not all(
            isinstance(key, str)
            and 0 < len(key) <= 64
            and key.isprintable()
            and isinstance(text, str)
            and 0 < len(text) <= 512
            and text.isprintable()
            for key, text in value.items()
        ):
            raise SmokeContractError(
                "Gemini citation metadata was invalid",
                reason="citation_metadata_invalid",
            )
        return dict(value)
    if not isinstance(value, Iterable) or isinstance(value, (str, bytes)):
        return {}
    metadata: dict[str, str] = {}
    for item in value:
        if len(metadata) >= 16:
            raise SmokeContractError(
                "Gemini citation metadata was invalid",
                reason="citation_metadata_invalid",
            )
        key = _field(item, "key")
        text = _field(item, "string_value")
        if (
            not isinstance(key, str)
            or not 0 < len(key) <= 64
            or not key.isprintable()
            or not isinstance(text, str)
            or not 0 < len(text) <= 512
            or not text.isprintable()
            or key in metadata
        ):
            raise SmokeContractError(
                "Gemini citation metadata was invalid",
                reason="citation_metadata_invalid",
            )
        metadata[key] = text
    return metadata


def _citation_excerpt(annotation: object, content_text: object) -> str:
    source = _field(annotation, "source")
    if isinstance(source, str) and source:
        return _bounded_excerpt(source)
    if not isinstance(content_text, str):
        raise SmokeContractError(
            "Gemini citation excerpt was unavailable",
            reason="citation_excerpt_unavailable",
        )
    start = _field(annotation, "start_index")
    end = _field(annotation, "end_index")
    encoded = content_text.encode("utf-8")
    if (
        isinstance(start, bool)
        or isinstance(end, bool)
        or not isinstance(start, int)
        or not isinstance(end, int)
        or not 0 <= start < end <= len(encoded)
    ):
        raise SmokeContractError(
            "Gemini citation excerpt was unavailable",
            reason="citation_excerpt_unavailable",
        )
    try:
        return _bounded_excerpt(encoded[start:end].decode("utf-8"))
    except UnicodeDecodeError:
        raise SmokeContractError(
            "Gemini citation excerpt was invalid",
            reason="citation_excerpt_invalid",
        ) from None


def _bounded_excerpt(value: str) -> str:
    normalized = value.strip()
    if (
        not normalized
        or len(value) > 4096
        or any(not character.isprintable() and character not in "\r\n\t" for character in value)
    ):
        raise SmokeContractError(
            "Gemini citation excerpt was invalid",
            reason="citation_excerpt_invalid",
        )
    return normalized


def _optional_page(value: object) -> int | None:
    if value is None:
        return None
    if (
        isinstance(value, bool)
        or not isinstance(value, int)
        or not 1 <= value <= 1_000_000
    ):
        raise SmokeContractError(
            "Gemini citation page number was invalid",
            reason="citation_page_invalid",
        )
    return value


def _optional_count(value: object) -> int | None:
    if value is None:
        return None
    if isinstance(value, bool) or not isinstance(value, int) or value < 0:
        raise SmokeContractError(
            "Gemini usage count was invalid",
            reason="usage_count_invalid",
        )
    return value


async def _collect(value: object) -> tuple[object, ...]:
    if isinstance(value, AsyncIterable):
        items = [item async for item in value]
        return tuple(items)
    if isinstance(value, Iterable) and not isinstance(value, (str, bytes, Mapping)):
        return tuple(value)
    raise SmokeContractError("Gemini document listing was unavailable")


class _TemporaryRetryService:
    def __init__(self) -> None:
        self.calls = 0

    async def index_revision(self, source_revision_id: str) -> IndexResult:
        from oms_hub.indexing.models import IndexState
        from oms_hub.indexing.service import IndexResult
        from oms_hub.providers.gemini.errors import GeminiTransientError

        self.calls += 1
        if self.calls == 1:
            raise GeminiTransientError("synthetic temporary failure")
        return IndexResult(source_revision_id, IndexState.READY)


def synthetic_pdf_bytes() -> bytes:
    output = BytesIO()
    page = Canvas(output, pagesize=letter, invariant=1, pageCompression=0)
    page.setTitle("Task 2.8 synthetic Gemini contract fixture")
    page.drawString(72, 720, SYNTHETIC_FACT)
    page.save()
    return output.getvalue()


def _normalize_pptx(path: Path) -> None:
    normalized = path.with_suffix(".normalized")
    with ZipFile(path) as source:
        members = {
            info.filename: source.read(info.filename)
            for info in source.infolist()
            if not info.is_dir()
        }
    core_properties = "docProps/core.xml"
    members[core_properties] = re.sub(
        rb"(<dcterms:(?:created|modified)[^>]*>).*?(</dcterms:(?:created|modified)>)",
        rb"\g<1>1980-01-01T00:00:00Z\g<2>",
        members[core_properties],
    )
    with ZipFile(normalized, "w", ZIP_DEFLATED, compresslevel=9) as target:
        for name in sorted(members):
            info = ZipInfo(name, (1980, 1, 1, 0, 0, 0))
            info.compress_type = ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = 0
            target.writestr(info, members[name], compress_type=ZIP_DEFLATED, compresslevel=9)
    normalized.replace(path)


def _synthetic_index_input(root: Path) -> IndexInputView:
    from PIL import Image, ImageDraw
    from pptx import Presentation

    from oms_hub.artifacts import ArtifactRole
    from oms_hub.document_processing.domain import DocumentLocator
    from oms_hub.files.atomic import sha256_file
    from oms_hub.knowledge.models import (
        EvidenceLocator,
        EvidenceLocatorKind,
        EvidenceUnit,
        SourceRevisionState,
    )
    from oms_hub.knowledge.service import (
        CanonicalInputArtifact,
        IndexAssetView,
        IndexInputView,
    )
    from oms_hub.providers.contracts import AuthorityClass

    pptx = root / "lecture.pptx"
    pdf = root / "lecture.pdf"
    markdown = root / "normalized.md"
    png = root / "visual.png"
    jpeg = root / "visual.jpg"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    assert slide.shapes.title is not None
    slide.shapes.title.text = SYNTHETIC_FACT
    presentation.save(str(pptx))
    _normalize_pptx(pptx)
    pdf.write_bytes(synthetic_pdf_bytes())
    markdown.write_text(f"# Synthetic\n\n{SYNTHETIC_FACT}\n", encoding="utf-8")
    for path, image_format, color in (
        (png, "PNG", (37, 99, 235)),
        (jpeg, "JPEG", (217, 119, 6)),
    ):
        image = Image.new("RGB", (160, 90), color)
        ImageDraw.Draw(image).text((8, 8), SYNTHETIC_MARKER, fill="white")
        image.save(path, format=image_format, quality=85)
    evidence = EvidenceUnit(
        evidence_id="synthetic-evidence", source_revision_id=SYNTHETIC_REVISION_ID,
        authority_class=AuthorityClass.COURSE_MATERIAL, course_id=SYNTHETIC_COURSE_ID,
        exam_id=SYNTHETIC_EXAM_ID, lecture_id=SYNTHETIC_LECTURE_ID,
        locator=EvidenceLocator(EvidenceLocatorKind.SLIDE, "1"), normalized_text=SYNTHETIC_FACT,
        content_sha256=hashlib.sha256(SYNTHETIC_FACT.encode()).hexdigest(),
    )
    return IndexInputView(
        source_document_id="task-2-8-synthetic-document", source_revision_id=SYNTHETIC_REVISION_ID,
        source_family="synthetic", revision_state=SourceRevisionState.READY,
        authority_class=AuthorityClass.COURSE_MATERIAL, course_id=SYNTHETIC_COURSE_ID,
        exam_id=SYNTHETIC_EXAM_ID, lecture_id=SYNTHETIC_LECTURE_ID,
        pptx=CanonicalInputArtifact(
            "synthetic:pptx",
            ArtifactRole.PPTX,
            pptx,
            sha256_file(pptx),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        pdf=CanonicalInputArtifact(
            "synthetic:pdf", ArtifactRole.PDF, pdf, sha256_file(pdf), "application/pdf"
        ),
        markdown=CanonicalInputArtifact(
            "synthetic:markdown",
            ArtifactRole.CLEANED,
            markdown,
            sha256_file(markdown),
            "text/markdown",
        ),
        evidence_units=(evidence,),
        assets=(
            IndexAssetView(
                asset_id="synthetic-png",
                path=png,
                media_type="image/png",
                sha256=sha256_file(png),
                locator=DocumentLocator("slide 1", slide_number=1),
                width=160,
                height=90,
                visual_semantic=True,
                evidence_ids=(evidence.evidence_id,),
            ),
            IndexAssetView(
                asset_id="synthetic-jpeg",
                path=jpeg,
                media_type="image/jpeg",
                sha256=sha256_file(jpeg),
                locator=DocumentLocator("slide 1", slide_number=1),
                width=160,
                height=90,
                visual_semantic=True,
                evidence_ids=(evidence.evidence_id,),
            ),
        ),
    )


def _redacted_identity(value: str) -> str:
    return f"sha256:{hashlib.sha256(value.encode('utf-8')).hexdigest()[:16]}"


def _failure_record(
    error: GeminiProviderError | SmokeContractError | SmokeTemporaryFailure,
    evidence: Mapping[str, object],
) -> dict[str, object]:
    if isinstance(error, GeminiProviderError):
        category = error.category
        retryable = error.retryable
        status_code = error.provider_status_code
    elif isinstance(error, SmokeTemporaryFailure):
        category = "transient"
        retryable = True
        status_code = None
    else:
        category = "contract"
        retryable = False
        status_code = None
    if (
        isinstance(status_code, bool)
        or not isinstance(status_code, int)
        or not 100 <= status_code <= 599
    ):
        status_code = None
    record: dict[str, object] = {
        "schema_version": 1,
        "status": "failed",
        "failure_stage": evidence.get("failure_stage", "unknown"),
        "error_category": category,
        "provider_status_code": status_code,
        "retryable": retryable,
        "resources_created": evidence.get("resources_created", {}),
        "cleanup": evidence.get("cleanup", {"attempted": 0, "status": "not_started"}),
        "reconciliation": evidence.get("reconciliation", "unknown"),
        "aggregate": evidence.get("aggregate", {"input_count": 0, "indexed_bytes": 0}),
        "input_results": evidence.get("input_results", []),
        "checks": evidence.get("checks", {}),
    }
    if isinstance(error, SmokeContractError) and error.reason is not None:
        record["contract_reason"] = error.reason
    if isinstance(error, GeminiProviderError) and error.diagnostic_code is not None:
        record["provider_reason"] = error.diagnostic_code
    return record


def run_temporary_failure_fixture() -> dict[str, object]:
    from oms_hub.db import Database
    from oms_hub.indexing.models import IndexJob, ProviderStore, StoreKey
    from oms_hub.indexing.repository import IndexRepository
    from oms_hub.indexing.worker import IndexWorker

    database = Database("sqlite://")
    database.create_schema()
    try:
        repository = IndexRepository(database)
        key = StoreKey.course(SYNTHETIC_COURSE_ID, SYNTHETIC_EXAM_ID)
        store = repository.create_store(
            ProviderStore(
                store_key=key,
                provider="gemini",
                provider_store_name="offlineFakeStores/task-2-8",
                embedding_model="models/gemini-embedding-2",
                authority_namespace=key.authority_namespace,
                course_id=key.course_id,
                exam_id=key.exam_id,
            )
        )
        job = repository.save_job(
            IndexJob(store_id=store.id, source_revision_id=SYNTHETIC_REVISION_ID)
        )
        clock = [datetime(2026, 8, 26, 12, 0, tzinfo=UTC)]
        service = _TemporaryRetryService()
        worker = IndexWorker(
            repository,
            service,
            worker_id="task-2-8-offline-retry",
            lease_seconds=60,
            now=lambda: clock[0],
        )

        worker.run_once()
        retry = repository.get_job(job.id)
        assert retry is not None and retry.next_attempt_at is not None
        next_attempt = datetime.fromisoformat(retry.next_attempt_at)
        backoff_seconds = round((next_attempt - clock[0]).total_seconds())
        clock[0] = next_attempt
        worker.run_once()
        resumed = repository.get_job(job.id)
        assert resumed is not None
        return {
            "first_state": retry.state.value,
            "retry_count": retry.retry_count,
            "error_category": retry.last_error_category,
            "backoff_seconds": backoff_seconds,
            "resumed_state": resumed.state.value,
            "service_calls": service.calls,
        }
    finally:
        database.close()


def _private_shadow_metadata(view: IndexInputView, item: Any) -> tuple[tuple[str, str], ...]:
    return (
        ("authority_class", view.authority_class.value),
        ("course_id", view.course_id),
        ("exam_id", view.exam_id),
        ("lecture_id", view.lecture_id),
        ("source_revision_id", view.source_revision_id),
        ("input_key", item.input_key),
        ("input_kind", item.input_kind),
        ("input_sha256", item.sha256),
    )


def _private_shadow_input_identity(input_key: object) -> str:
    if input_key in {"pptx", "pdf", "normalized_markdown"}:
        return str(input_key)
    if isinstance(input_key, str) and re.fullmatch(r"image\.[0-9a-f]{64}", input_key):
        return "visual_asset"
    return "unknown"


async def _run_shadow_sequence(
    session: ShadowSession,
    view: IndexInputView,
    preflight: dict[str, object],
    *,
    mode: Literal["public_matrix", "private_acceptance"],
    clock: Callable[[], float],
    failure_evidence: dict[str, object] | None = None,
    diagnostic_sink: DiagnosticSink | None = None,
    private_diagnostic_path: Path | None = None,
) -> dict[str, object]:
    if session.model_contract != PRIVATE_SHADOW_MODEL_CONTRACT:
        raise LiveSmokeBlocked("private shadow model contract mismatch")
    manifest = _private_shadow_manifest(view)
    input_results: list[dict[str, str]] = []
    aggregate = {
        "input_count": len(manifest.inputs),
        "indexed_bytes": sum(item.path.stat().st_size for item in manifest.inputs),
        "transient_attempts": 0,
    }

    def refresh_transient_attempts() -> None:
        aggregate["transient_attempts"] = _session_transient_attempts(session)
    store_name: str | None = None
    stores: list[str] = []
    files: list[tuple[str, str]] = []
    documents: list[str] = []
    operations: list[str] = []
    states: list[str] = []
    public_checks = {
        name: "not_run"
        for name in (
            "positive_answer",
            "citation_presence",
            "negative_structured_output",
            "create_store",
            "document_listing",
            "cleanup_store",
            "cleanup_document",
            "cleanup_file",
            "wrong_lecture_filtering",
        )
    }
    resource_states = {"document": "not_started", "file": "not_started", "store": "not_started"}
    cleanup_attempts = {"document": 0, "file": 0, "store": 0}
    cleanup_successes = {"document": 0, "file": 0, "store": 0}
    cleanup_failures = {"document": False, "file": False, "store": False}
    uncertain_resources: set[str] = set()
    failure: BaseException | None = None
    failure_stage = "unknown"
    input_failure: BaseException | None = None
    input_failure_stage = "unknown"
    input_failure_identity = "none"
    active_stage = "prior_state_check"
    active_input_identity = "none"
    cleanup_failed = False
    cleanup_unknown = False
    document_cleanup_failed = False
    reconciliation_unknown = False
    reconciliation_not_empty = False
    file_reconciliation_empty = False
    store_reconciliation_empty = False
    positive: PrivateShadowQueryAudit | None = None
    negative: PrivateShadowQueryAudit | None = None
    run_token = uuid4().hex
    display_names = tuple(
        f"task-2-8-{mode}-{run_token}-{ordinal:03d}"
        for ordinal, _ in enumerate(manifest.inputs, start=1)
    )
    store_display_name = f"task-2-8-{mode}-{run_token}"
    started = clock()
    if diagnostic_sink is not None and mode == "public_matrix":
        diagnostic_sink.capture(
            "contract.expected",
            {
                "course_id": view.course_id,
                "exam_id": view.exam_id,
                "lecture_id": view.lecture_id,
                "source_revision_id": view.source_revision_id,
                "fixture_sha256": hashlib.sha256(view.pdf.path.read_bytes()).hexdigest(),
            },
        )
    try:
        prior_state_present = bool(
            await session.find_stores(store_display_name)
            or await session.find_files(display_names)
        )
    except BaseException as prior_check_error:
        refresh_transient_attempts()
        diagnostic_sha256 = _finalize_private_terminal_diagnostic(
            mode,
            private_diagnostic_path,
            prior_check_error,
            failure_stage=active_stage,
            input_identity=active_input_identity,
        )
        if failure_evidence is not None:
            if mode == "public_matrix":
                failure_evidence.update(
                    {
                        "failure_stage": active_stage,
                        "resources_created": dict(resource_states),
                        "cleanup": {"attempted": 0, "status": "not_started"},
                        "reconciliation": "unknown",
                        "aggregate": aggregate,
                        "input_results": input_results,
                        "checks": dict(public_checks),
                    }
                )
            else:
                failure_evidence.update(
                    private_shadow_failure_record(
                        preflight,
                        prior_check_error,
                        failure_stage=active_stage,
                        states=states,
                        cleanup_outcome="unknown",
                        reconciliation_outcome="unknown",
                        transient_attempts=_session_transient_attempts(session),
                        diagnostic_sha256=diagnostic_sha256,
                    ).model_dump(mode="json")
                )
        if diagnostic_sink is not None and mode == "public_matrix":
            diagnostic_sink.capture("contract.check_matrix", public_checks)
        raise
    if prior_state_present:
        refresh_transient_attempts()
        prior_state_mismatch = LiveSmokeBlocked(
            "private shadow prior operator state mismatch"
        )
        diagnostic_sha256 = _finalize_private_terminal_diagnostic(
            mode,
            private_diagnostic_path,
            prior_state_mismatch,
            failure_stage=active_stage,
            input_identity=active_input_identity,
        )
        if failure_evidence is not None:
            if mode == "public_matrix":
                failure_evidence.update(
                    {
                        "failure_stage": active_stage,
                        "resources_created": dict(resource_states),
                        "cleanup": {"attempted": 0, "status": "not_started"},
                        "reconciliation": "not_empty",
                        "aggregate": aggregate,
                        "input_results": input_results,
                        "checks": dict(public_checks),
                    }
                )
            else:
                failure_evidence.update(
                    private_shadow_failure_record(
                        preflight,
                        prior_state_mismatch,
                        failure_stage=active_stage,
                        states=states,
                        cleanup_outcome="unknown",
                        reconciliation_outcome="not_empty",
                        transient_attempts=_session_transient_attempts(session),
                        diagnostic_sha256=diagnostic_sha256,
                    ).model_dump(mode="json")
                )
        if diagnostic_sink is not None and mode == "public_matrix":
            diagnostic_sink.capture("contract.check_matrix", public_checks)
        raise prior_state_mismatch
    states.append("prior_operator_state_empty")
    try:
        active_stage = "create_store"
        if mode == "public_matrix":
            resource_states["store"] = "unknown"
        store_name = await session.create_store(
            store_display_name,
            PRIVATE_SHADOW_MODEL_CONTRACT[2],
        )
        stores.append(store_name)
        if mode == "public_matrix":
            resource_states["store"] = "confirmed"
            public_checks["create_store"] = "passed"
        states.append("store_created")
        for item, display_name in zip(manifest.inputs, display_names, strict=True):
            active_input_identity = _private_shadow_input_identity(item.input_key)
            try:
                active_stage = "upload_input"
                if mode == "public_matrix":
                    resource_states["file"] = "unknown"
                file_name = await session.upload_input(
                    display_name,
                    item.path,
                    item.media_type,
                )
                files.append((file_name, item.input_key))
                if mode == "public_matrix":
                    resource_states["file"] = (
                        "unknown" if "file" in uncertain_resources else "confirmed"
                    )
                chunking = (
                    {
                        "white_space_config": {
                            "max_tokens_per_chunk": 700,
                            "max_overlap_tokens": 100,
                        }
                    }
                    if item.input_key == "normalized_markdown"
                    else None
                )
                active_stage = "import_input"
                if mode == "public_matrix":
                    resource_states["document"] = "unknown"
                operation = await session.import_input(
                    store_name,
                    file_name,
                    _private_shadow_metadata(view, item),
                    chunking,
                )
                operations.append(operation)
                active_stage = "wait_for_import"
                document_name = await session.wait_for_import(operation)
                documents.append(document_name)
                if mode == "public_matrix":
                    resource_states["document"] = (
                        "unknown" if "document" in uncertain_resources else "confirmed"
                    )
                    input_results.append(
                        {
                            "input_kind": item.input_kind,
                            "stage": "import_input",
                            "outcome": "passed",
                            "error_category": "none",
                        }
                    )
            except BaseException as error:
                if mode != "public_matrix":
                    raise
                uncertain_resource = {
                    "upload_input": "file",
                    "import_input": "document",
                    "wait_for_import": "document",
                }.get(active_stage)
                if uncertain_resource is not None:
                    uncertain_resources.add(uncertain_resource)
                    resource_states[uncertain_resource] = "unknown"
                category = (
                    error.category
                    if isinstance(error, GeminiProviderError)
                    else "transient"
                    if isinstance(error, SmokeTemporaryFailure)
                    else "contract"
                    if isinstance(error, SmokeContractError)
                    else "unknown"
                )
                input_results.append(
                    {
                        "input_kind": item.input_kind,
                        "stage": active_stage,
                        "outcome": "failed",
                        "error_category": category,
                    }
                )
                if input_failure is None:
                    input_failure = error
                    input_failure_stage = active_stage
                    input_failure_identity = active_input_identity
        if input_failure is not None:
            active_stage = input_failure_stage
            active_input_identity = input_failure_identity
            raise input_failure
        states.extend(
            (f"inputs_uploaded:{len(files)}", f"inputs_imported:{len(documents)}")
        )
        active_input_identity = "none"
        file_bindings = tuple(files)
        active_stage = "positive_query"
        positive = await session.query_private(
            store_name,
            (
                "Using only the indexed lecture, return one concise supported statement "
                "with a citation."
            ),
            SmokeScope(view.course_id, view.exam_id, view.lecture_id),
            source_revision_id=view.source_revision_id,
            manifest=manifest,
            file_bindings=file_bindings,
            require_structured_supported=mode == "public_matrix",
        )
        active_stage = "positive_validation"
        if mode == "public_matrix":
            public_checks["positive_answer"] = (
                "positive_answer_unsupported"
                if positive.supported is not True
                else "positive_answer_missing_marker"
                if positive.answer is None or SYNTHETIC_MARKER not in positive.answer
                else "passed"
            )
            public_checks["citation_presence"] = (
                "positive_citation_missing"
                if positive.citation_count == 0
                else "positive_citation_unresolved"
                if positive.resolved_citation_count != positive.citation_count
                else "citation_page_absent"
                if public_checks["positive_answer"] == "passed"
                and positive.citation_page is None
                else "citation_excerpt_absent"
                if public_checks["positive_answer"] == "passed"
                and not positive.citation_excerpt
                else "passed"
            )
            if positive.supported is not True:
                raise SmokeContractError(
                    "public smoke answer was unsupported",
                    reason="positive_answer_unsupported",
                )
            if positive.answer is None or SYNTHETIC_MARKER not in positive.answer:
                raise SmokeContractError(
                    "public smoke marker was invalid",
                    reason="positive_answer_missing_marker",
                )
            if positive.citation_count == 0:
                raise SmokeContractError(
                    "public smoke citation was missing",
                    reason="positive_citation_missing",
                )
            if positive.resolved_citation_count != positive.citation_count:
                raise SmokeContractError(
                    "public smoke citation was unresolved",
                    reason="positive_citation_unresolved",
                )
            if positive.citation_page is None:
                raise SmokeContractError(
                    "public smoke citation page was absent",
                    reason="citation_page_absent",
                )
            if not positive.citation_excerpt:
                raise SmokeContractError(
                    "public smoke citation excerpt was absent",
                    reason="citation_excerpt_absent",
                )
        elif (
            positive.citation_count < 1
            or positive.resolved_citation_count != positive.citation_count
        ):
            raise SmokeContractError(
                "private shadow citations were unresolved",
                reason="private_citation_unresolved",
            )
        if mode == "public_matrix" and (
            positive.supported is not True
            or positive.answer is None
            or SYNTHETIC_MARKER not in positive.answer
            or positive.citation_page is None
            or positive.citation_excerpt is None
        ):
            raise SmokeContractError("public smoke marker was invalid", reason="positive_answer_missing_marker")
        states.append("positive_query_complete")
        if view.lecture_id == PRIVATE_SHADOW_WRONG_LECTURE_ID:
            raise LiveSmokeBlocked("private shadow wrong-scope identity collided")
        active_stage = "negative_query"
        negative = await session.query_private(
            store_name,
            (
                "Use only files matching the requested lecture scope. If none match, "
                "return an empty answer and supported=false."
            ),
            SmokeScope(
                view.course_id,
                view.exam_id,
                WRONG_LECTURE_ID if mode == "public_matrix" else PRIVATE_SHADOW_WRONG_LECTURE_ID,
            ),
            source_revision_id=view.source_revision_id,
            manifest=manifest,
            file_bindings=file_bindings,
            require_structured_no_result=True,
        )
        active_stage = "negative_validation"
        if mode == "public_matrix":
            public_checks["negative_structured_output"] = (
                "passed"
                if negative.supported is False and negative.answer_empty is True
                else "negative_structured_output_invalid"
            )
            public_checks["wrong_lecture_filtering"] = (
                "passed"
                if negative.citation_count == 0
                else "wrong_lecture_retrieved"
            )
        if (
            negative.citation_count != 0
            or negative.supported is not False
            or negative.answer_empty is not True
        ):
            raise SmokeContractError(
                "private shadow wrong scope retrieved evidence",
                reason="private_wrong_scope_retrieved",
            )
        states.append("wrong_scope_query_complete")
    except BaseException as error:
        failure = error
        failure_stage = active_stage
    finally:
        for document_name in reversed(documents):
            cleanup_attempts["document"] += 1
            try:
                await session.delete_document(document_name)
                cleanup_successes["document"] += 1
            except BaseException:
                cleanup_failed = True
                document_cleanup_failed = True
                cleanup_failures["document"] = True
        states.append(f"documents_delete_attempted:{len(documents)}")
        try:
            discovered = await session.find_files(display_names)
        except BaseException:
            cleanup_failed = True
            cleanup_unknown = True
        else:
            known = {file_name for file_name, _ in files}
            files.extend(
                (file_name, "reconciled")
                for file_name in discovered
                if file_name not in known
            )
        for file_name, _ in reversed(files):
            cleanup_attempts["file"] += 1
            try:
                await session.delete_file(file_name)
                cleanup_successes["file"] += 1
            except BaseException:
                cleanup_failed = True
                cleanup_failures["file"] = True
        states.append(f"files_delete_attempted:{len(files)}")
        try:
            remaining_files = await session.find_files(display_names)
            if remaining_files:
                cleanup_failed = True
                reconciliation_not_empty = True
            else:
                file_reconciliation_empty = True
                states.append("file_reconciliation_empty")
        except BaseException:
            cleanup_failed = True
            cleanup_unknown = True
            reconciliation_unknown = True
        try:
            discovered_stores = await session.find_stores(store_display_name)
        except BaseException:
            cleanup_failed = True
            cleanup_unknown = True
        else:
            known_stores = set(stores)
            stores.extend(
                name for name in discovered_stores if name not in known_stores
            )
        for current_store in reversed(stores):
            cleanup_attempts["store"] += 1
            try:
                await session.delete_store(current_store)
                cleanup_successes["store"] += 1
            except BaseException:
                cleanup_failed = True
                cleanup_failures["store"] = True
        states.append(f"stores_delete_attempted:{len(stores)}")
        try:
            remaining_stores = await session.find_stores(store_display_name)
            if remaining_stores:
                cleanup_failed = True
                reconciliation_not_empty = True
            else:
                store_reconciliation_empty = True
                states.append("store_reconciliation_empty")
        except BaseException:
            cleanup_failed = True
            cleanup_unknown = True
            reconciliation_unknown = True
    if mode == "public_matrix" and "document" in uncertain_resources:
        cleanup_unknown = True
        reconciliation_unknown = True
    cleanup_outcome = (
        "unknown"
        if cleanup_unknown
        else "failed"
        if cleanup_failed
        else "complete"
    )
    reconciliation_outcome = (
        "unknown"
        if reconciliation_unknown or document_cleanup_failed
        else "not_empty"
        if reconciliation_not_empty
        else "empty"
        if file_reconciliation_empty and store_reconciliation_empty
        else "unknown"
    )
    if mode == "public_matrix":
        for resource, check in (
            ("document", "cleanup_document"),
            ("file", "cleanup_file"),
            ("store", "cleanup_store"),
        ):
            public_checks[check] = (
                "not_available"
                if cleanup_attempts[resource] == 0
                else "cleanup_delete_failed"
                if cleanup_failures[resource]
                or cleanup_successes[resource] != cleanup_attempts[resource]
                else "passed"
            )
    refresh_transient_attempts()
    if failure is not None:
        if mode == "public_matrix":
            if failure_stage == "create_store":
                public_checks["create_store"] = "create_store_failed"
            elif failure_stage == "positive_query":
                public_checks["positive_answer"] = (
                    failure.reason
                    if isinstance(failure, SmokeContractError) and failure.reason
                    else "positive_query_failed"
                )
                if public_checks["positive_answer"] == "positive_answer_missing_marker":
                    public_checks["citation_presence"] = "positive_citation_missing"
            elif failure_stage == "negative_query":
                public_checks["negative_structured_output"] = "negative_query_failed"
        diagnostic_sha256 = _finalize_private_terminal_diagnostic(
            mode,
            private_diagnostic_path,
            failure,
            failure_stage=failure_stage,
            input_identity=active_input_identity,
        )
        if failure_evidence is not None:
            if mode == "public_matrix":
                failure_evidence.update(
                    {
                        "failure_stage": failure_stage,
                        "resources_created": dict(resource_states),
                        "cleanup": {
                            "attempted": len(documents) + len(files) + len(stores),
                            "status": (
                                "unknown"
                                if failure_stage == "create_store"
                                else "completed"
                                if cleanup_outcome == "complete"
                                else cleanup_outcome
                            ),
                        },
                        "reconciliation": reconciliation_outcome,
                        "aggregate": aggregate,
                        "input_results": input_results,
                        "checks": dict(public_checks),
                    }
                )
            else:
                failure_evidence.update(
                    private_shadow_failure_record(
                        preflight,
                        failure,
                        failure_stage=failure_stage,
                        states=states,
                        cleanup_outcome=cleanup_outcome,
                        reconciliation_outcome=reconciliation_outcome,
                        input_identity=active_input_identity,
                        transient_attempts=_session_transient_attempts(session),
                        diagnostic_sha256=diagnostic_sha256,
                    ).model_dump(mode="json")
                )
        if diagnostic_sink is not None and mode == "public_matrix":
            diagnostic_sink.capture("contract.check_matrix", public_checks)
        raise failure
    if cleanup_failed:
        cleanup_error = SmokeContractError(
            "private shadow cleanup failed",
            reason="private_cleanup_failed",
        )
        diagnostic_sha256 = _finalize_private_terminal_diagnostic(
            mode,
            private_diagnostic_path,
            cleanup_error,
            failure_stage="cleanup",
            input_identity="none",
        )
        if failure_evidence is not None:
            if mode == "public_matrix":
                failure_evidence.update(
                    {
                        "failure_stage": "cleanup",
                        "resources_created": dict(resource_states),
                        "cleanup": {
                            "attempted": len(documents) + len(files) + len(stores),
                            "status": "completed" if cleanup_outcome == "complete" else cleanup_outcome,
                        },
                        "reconciliation": reconciliation_outcome,
                        "aggregate": aggregate,
                        "input_results": input_results,
                        "checks": dict(public_checks),
                    }
                )
            else:
                failure_evidence.update(
                private_shadow_failure_record(
                    preflight,
                    cleanup_error,
                    failure_stage="cleanup",
                    states=states,
                    cleanup_outcome=cleanup_outcome,
                    reconciliation_outcome=reconciliation_outcome,
                    input_identity="none",
                    transient_attempts=_session_transient_attempts(session),
                    diagnostic_sha256=diagnostic_sha256,
                ).model_dump(mode="json")
                )
        if diagnostic_sink is not None and mode == "public_matrix":
            diagnostic_sink.capture("contract.check_matrix", public_checks)
        raise cleanup_error from None
    assert positive is not None and negative is not None
    record = {
        **preflight,
        "status": "passed",
        "provider_operation_states": states,
        "transient_attempts": _session_transient_attempts(session),
        "failure_class": "none",
        "citation_resolution_rate": (
            positive.resolved_citation_count / positive.citation_count
        ),
        "duration_ms": round((clock() - started) * 1000),
        "token_usage": {
            "input": positive.input_tokens + negative.input_tokens,
            "output": positive.output_tokens + negative.output_tokens,
        },
    }
    if mode == "public_matrix":
        if diagnostic_sink is not None:
            diagnostic_sink.capture("contract.check_matrix", public_checks)
        assert store_name is not None
        assert files and operations and documents
        assert positive.citation_page is not None
        assert positive.citation_excerpt is not None
        return {
            "schema_version": 1,
            "status": "passed",
            "input_results": input_results,
            "aggregate": aggregate,
            "cleanup": {
                "attempted": len(documents) + len(files) + len(stores),
                "status": "completed",
            },
            "reconciliation": reconciliation_outcome,
            "error_category": "none",
        }
    return validate_private_shadow_record(record, process_exit_code=0)


class _CompositionProbeSession:
    """In-process fake used only by the tracked composition verifier."""

    model_contract = PRIVATE_SHADOW_MODEL_CONTRACT

    def __init__(self) -> None:
        self._files: dict[str, str] = {}
        self._stores: dict[str, str] = {}

    async def create_store(self, display_name: str, embedding_model: str) -> str:
        del embedding_model
        name = "composition-probe-store"
        self._stores[name] = display_name
        return name

    async def find_stores(self, display_name: str) -> tuple[str, ...]:
        return tuple(name for name, value in self._stores.items() if value == display_name)

    async def upload_input(self, display_name: str, path: Path, media_type: str) -> str:
        del path, media_type
        name = f"composition-probe-file-{len(self._files) + 1}"
        self._files[name] = display_name
        return name

    async def import_input(
        self,
        store_name: str,
        file_name: str,
        metadata: tuple[tuple[str, str], ...],
        chunking: object | None,
    ) -> str:
        del store_name, file_name, metadata, chunking
        return "composition-probe-operation"

    async def find_files(self, display_names: tuple[str, ...]) -> tuple[str, ...]:
        return tuple(name for name, value in self._files.items() if value in display_names)

    async def wait_for_import(self, operation_name: str) -> str:
        return f"composition-probe-document-{operation_name}"

    async def query_private(
        self,
        store_name: str,
        prompt: str,
        scope: SmokeScope,
        *,
        source_revision_id: str,
        manifest: object,
        file_bindings: tuple[tuple[str, str], ...],
        require_structured_no_result: bool = False,
        require_structured_supported: bool = False,
    ) -> PrivateShadowQueryAudit:
        del store_name, prompt, scope, source_revision_id, manifest, file_bindings
        if require_structured_no_result:
            return PrivateShadowQueryAudit(0, 0, 0, 0, False, True)
        if require_structured_supported:
            raise SmokeContractError("private composition probe has no public mode")
        return PrivateShadowQueryAudit(1, 1, 0, 0, None, None)

    async def delete_document(self, document_name: str) -> None:
        del document_name

    async def delete_file(self, file_name: str) -> None:
        self._files.pop(file_name, None)

    async def delete_store(self, store_name: str) -> None:
        self._stores.pop(store_name, None)


async def run_private_shadow_composition_probe() -> dict[str, object]:
    """Exercise the private lifecycle locally without credentials or provider access."""

    with tempfile.TemporaryDirectory(prefix="task-2-8-composition-probe-") as directory:
        view = _synthetic_index_input(Path(directory))
        return await _run_shadow_sequence(
            _CompositionProbeSession(),
            view,
            _private_shadow_preflight_from_view(view),
            mode="private_acceptance",
            clock=lambda: 0.0,
        )


async def run_authorized_private_shadow(
    slide_revision_id: str,
    *,
    schema_version: int,
    artifacts: ArtifactService,
    materialization_root: Path,
    approved_preflight: Mapping[str, object],
    diagnostic_path: Path,
    secret_store: SecretStore | None = None,
    session_factory: Callable[[str], ShadowSession] | None = None,
    parser: Any | None = None,
    clock: Callable[[], float] = monotonic,
    failure_evidence: dict[str, object] | None = None,
) -> dict[str, object]:
    if failure_evidence is not None:
        failure_evidence.clear()
    if os.getenv("RUN_PRIVATE_GEMINI_SHADOW") != "1":
        raise LiveSmokeBlocked(
            "RUN_PRIVATE_GEMINI_SHADOW=1 is required for a private shadow"
        )
    private_diagnostic_path = _validate_private_diagnostic_capability(diagnostic_path)
    try:
        sdk_version = importlib.metadata.version("google-genai")
    except importlib.metadata.PackageNotFoundError:
        raise LiveSmokeBlocked("private shadow model contract mismatch") from None
    if sdk_version != PRIVATE_SHADOW_MODEL_CONTRACT[0]:
        raise LiveSmokeBlocked("private shadow model contract mismatch")
    config_probe = GeminiConfig(api_key=SecretStr("private-shadow-contract-probe"))
    if (
        config_probe.sdk_version,
        config_probe.file_search_model,
        config_probe.embedding_model,
        config_probe.api_version,
    ) != PRIVATE_SHADOW_MODEL_CONTRACT:
        raise LiveSmokeBlocked("private shadow model contract mismatch")
    view = prepare_private_shadow_index_input(
        slide_revision_id,
        schema_version=schema_version,
        artifacts=artifacts,
        materialization_root=materialization_root,
        parser=parser,
    )
    preflight = _private_shadow_preflight_from_view(view)
    expected = _expected_private_shadow_preflight(view, _private_shadow_manifest(view))
    if preflight != expected or dict(approved_preflight) != expected:
        raise LiveSmokeBlocked("private shadow preflight mismatch")
    if secret_store is None:
        from oms_hub.security.secret_store import KeyringSecretStore

        secret_store = KeyringSecretStore()
    try:
        api_key = secret_store.get("gemini-api-key")
    except Exception:
        raise LiveSmokeBlocked("stored Gemini credential is unavailable") from None
    if not isinstance(api_key, str) or not api_key.strip():
        raise LiveSmokeBlocked("stored Gemini credential is unavailable")
    normalized_key = api_key.strip()
    session = (
        GoogleGenaiSmokeSession(normalized_key)
        if session_factory is None
        else session_factory(normalized_key)
    )
    return await _run_shadow_sequence(
        session,
        view,
        preflight,
        mode="private_acceptance",
        clock=clock,
        failure_evidence=failure_evidence,
        private_diagnostic_path=private_diagnostic_path,
    )


async def run_authorized_live_smoke(
    *,
    secret_store: SecretStore | None = None,
    session_factory: Callable[[str], ShadowSession] | None = None,
    failure_evidence: dict[str, object] | None = None,
    diagnostic_request: _SyntheticDiagnosticRequest | None = None,
) -> dict[str, object]:
    if failure_evidence is not None:
        failure_evidence.clear()
    diagnostic_sink: _SyntheticDiagnosticSink | None = None
    if diagnostic_request is not None:
        _validate_diagnostic_request(diagnostic_request)
        if session_factory is not None:
            raise LiveSmokeBlocked(
                "synthetic diagnostics require the default synthetic session"
            )
    if os.getenv("RUN_LIVE_GEMINI_TESTS") != "1":
        raise LiveSmokeBlocked("RUN_LIVE_GEMINI_TESTS=1 is required for a live smoke")
    if diagnostic_request is not None:
        diagnostic_sink = _SyntheticDiagnosticSink.open(diagnostic_request)
    if secret_store is None:
        from oms_hub.security.secret_store import KeyringSecretStore

        secret_store = KeyringSecretStore()
    try:
        api_key = secret_store.get("gemini-api-key")
    except Exception:
        raise LiveSmokeBlocked("stored Gemini credential is unavailable") from None
    if not isinstance(api_key, str) or not api_key.strip():
        raise LiveSmokeBlocked("stored Gemini credential is unavailable")
    normalized_key = api_key.strip()
    if diagnostic_sink is not None:
        diagnostic_sink.add_secret(normalized_key)
    if session_factory is None:
        session: ShadowSession = GoogleGenaiSmokeSession(
            normalized_key,
            diagnostic_sink=diagnostic_sink,
        )
    else:
        session = session_factory(normalized_key)
    try:
        with tempfile.TemporaryDirectory(prefix="task-2-8-synthetic-") as directory:
            view = _synthetic_index_input(Path(directory))
            return await _run_shadow_sequence(
                session,
                view,
                _private_shadow_preflight_from_view(view),
                mode="public_matrix",
                clock=monotonic,
                failure_evidence=failure_evidence,
                diagnostic_sink=diagnostic_sink,
            )
    except BaseException as error:
        if diagnostic_sink is not None:
            diagnostic_sink.capture_exception("contract.failure", error)
        raise
    finally:
        if diagnostic_sink is not None:
            diagnostic_sink.close()


def _plan() -> dict[str, object]:
    return {
        "status": "ready_after_independent_review",
        "calls_provider": False,
        "reads_secrets": False,
        "required_flag": "RUN_LIVE_GEMINI_TESTS=1",
        "required_command": "python scripts/run-gemini-contract-smoke.py --execute-live",
        "required_authorization": (
            "Connor must explicitly authorize one synthetic Gemini smoke, disposable provider "
            "create/query/delete operations, quota/cost, and approved secret-store access."
        ),
        "required_owner_action": (
            "Independent specification and quality/security reviews must approve the exact "
            "adapter commit before run_authorized_live_smoke crosses the provider boundary."
        ),
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--execute-live", action="store_true")
    parser.add_argument("--synthetic-diagnostic-output", type=Path)
    args = parser.parse_args(argv)
    if not args.execute_live:
        print(json.dumps(_plan(), indent=2, sort_keys=True))
        return 0
    failure_evidence: dict[str, object] = {}
    diagnostic_request = (
        _synthetic_diagnostic_request(args.synthetic_diagnostic_output)
        if args.synthetic_diagnostic_output is not None
        else None
    )
    try:
        if diagnostic_request is None:
            record = asyncio.run(
                run_authorized_live_smoke(failure_evidence=failure_evidence)
            )
        else:
            record = asyncio.run(
                run_authorized_live_smoke(
                    failure_evidence=failure_evidence,
                    diagnostic_request=diagnostic_request,
                )
            )
    except LiveSmokeBlocked as error:
        parser.error(str(error))
    except (GeminiProviderError, SmokeContractError, SmokeTemporaryFailure) as error:
        print(json.dumps(_failure_record(error, failure_evidence), indent=2, sort_keys=True))
        return 1
    print(json.dumps(record, indent=2, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
